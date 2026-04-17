"""AWS/Databricks Co-Exist Strategy v2 — cleaner, softer, win-win-win theme."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Refined palette — less orange, more professional
DARK = RGBColor(0x1B, 0x2A, 0x3D)
NAVY = RGBColor(0x23, 0x2F, 0x3E)
WHITE = RGBColor(0xF1, 0xF5, 0xF9)
LIGHT = RGBColor(0xCB, 0xD5, 0xE1)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
TEAL = RGBColor(0x06, 0xB6, 0xD4)
SOFT_BLUE = RGBColor(0x38, 0xBD, 0xF8)
AWS_ORANGE = RGBColor(0xFF, 0x99, 0x00)
ACCENT = RGBColor(0x22, 0xD3, 0xEE)  # cyan accent
GREEN = RGBColor(0x34, 0xD3, 0x99)
SOFT_GREEN = RGBColor(0x6E, 0xE7, 0xB7)
RED_SOFT = RGBColor(0xFB, 0x92, 0x3C)  # warm orange for warnings
PURPLE = RGBColor(0xA7, 0x8B, 0xFA)
PAL_GREEN = RGBColor(0x4A, 0xDE, 0x80)
CARD_BG = RGBColor(0x1E, 0x2D, 0x40)
INTERNAL_RED = RGBColor(0xF8, 0x71, 0x71)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def bg(slide, color=DARK):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def t(slide, l, tp, w, h, text, sz=18, c=WHITE, b=False, a=PP_ALIGN.LEFT):
    tx = slide.shapes.add_textbox(Inches(l), Inches(tp), Inches(w), Inches(h))
    tx.text_frame.word_wrap = True
    p = tx.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = c
    p.font.bold = b
    p.alignment = a
    return tx

def ap(tf, text, sz=14, c=WHITE, b=False, a=PP_ALIGN.LEFT, sb=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = c
    p.font.bold = b
    p.alignment = a
    if sb: p.space_before = Pt(sb)

def cd(slide, l, tp, w, h, bc=TEAL, fc=CARD_BG):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(tp), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fc
    s.line.color.rgb = bc
    s.line.width = Pt(1)
    return s

def line(slide, l, tp, w, h, c):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(tp), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = c
    s.line.fill.background()

def tag(slide, l, tp, text, c):
    """Small tag label"""
    t(slide, l, tp, 1.5, 0.25, text, 8, c, True)

# ==================== SLIDE 1: TITLE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
line(slide, 0, 0, 13.333, 0.04, TEAL)

t(slide, 1.2, 1.5, 10, 0.8, "AWS + Databricks", 42, WHITE, False)
t(slide, 1.2, 2.3, 10, 0.8, "Co-Existence Strategy", 42, ACCENT, True)
t(slide, 1.2, 3.5, 8, 0.5, "A win for customers. A win for Databricks. A win for AWS.", 16, MUTED)

line(slide, 1.2, 4.3, 3, Pt(1), RGBColor(0x33, 0x44, 0x55))

t(slide, 1.2, 4.6, 8, 0.8, "Open data lakes on S3 with Apache Iceberg let customers use the best service\nfor each workload — without forcing a single-vendor choice.", 14, LIGHT)

t(slide, 1.2, 6.0, 4, 0.3, "David Eyre", 13, ACCENT, True)
t(slide, 1.2, 6.3, 4, 0.3, "Emerging Tech Solutions — Amazon Web Services", 11, MUTED)

# Win-win-win cards on right
wins = [
    ("🏢 Customer Win", "Best tool for each job. No lock-in. Lower cost at scale.", GREEN),
    ("🤝 Databricks Win", "Keeps Spark workloads. Grows on S3/Iceberg foundation.", PURPLE),
    ("☁️ AWS Win", "More services consumed. Deeper account penetration.", TEAL),
]
for i, (title, body, color) in enumerate(wins):
    y = 1.8 + i * 1.5
    cd(slide, 8.5, y, 4.3, 1.2, color)
    t(slide, 8.7, y + 0.1, 4.0, 0.3, title, 13, color, True)
    t(slide, 8.7, y + 0.45, 4.0, 0.5, body, 11, LIGHT)

# ==================== SLIDE 2: THE OPPORTUNITY [CUSTOMER-FACING] ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
line(slide, 0, 0, 13.333, 0.04, TEAL)
tag(slide, 11.5, 0.2, "CUSTOMER-FACING", GREEN)
t(slide, 0.8, 0.3, 10, 0.6, "The Open Data Lake Opportunity", 28, WHITE, True)
t(slide, 0.8, 0.85, 10, 0.4, "Your data should work for you — not lock you into a single platform", 14, MUTED)

# Left: What customers get
cd(slide, 0.8, 1.5, 5.5, 5.0, TEAL)
t(slide, 1.0, 1.6, 5.0, 0.3, "What This Means for You", 16, ACCENT, True)
benefits = [
    "Store data once in S3 — query from any service",
    "Apache Iceberg: open format readable by every engine",
    "Keep Databricks for Spark workloads you've already built",
    "Add AWS-native services where they're the best fit",
    "One bill, one IAM, one security model",
    "Scale storage and compute independently",
    "Switch or add tools without migrating data",
]
txBox = t(slide, 1.0, 2.1, 5.0, 0.3, "→  " + benefits[0], 12, LIGHT)
for b in benefits[1:]:
    ap(txBox.text_frame, "→  " + b, 12, LIGHT, sb=6)

# Right: The architecture concept
cd(slide, 6.8, 1.5, 5.5, 5.0, SOFT_BLUE)
t(slide, 7.0, 1.6, 5.0, 0.3, "How It Works", 16, SOFT_BLUE, True)

# Simple stack
layers = [
    ("🪣 Amazon S3 + Apache Iceberg", "Open data lake — your data, open format", SOFT_BLUE),
    ("📊 Redshift  |  🔍 OpenSearch  |  📈 QuickSight", "SQL analytics, search & RAG, embedded BI", TEAL),
    ("🤖 Amazon Bedrock + SageMaker", "GenAI and ML across all your data", ACCENT),
    ("🔥 Databricks", "Spark ETL, data engineering, notebooks", PURPLE),
]
for i, (name, desc, color) in enumerate(layers):
    y = 2.2 + i * 1.0
    cd(slide, 7.0, y, 5.0, 0.8, color)
    t(slide, 7.2, y + 0.05, 4.6, 0.3, name, 11, color, True, PP_ALIGN.CENTER)
    t(slide, 7.2, y + 0.35, 4.6, 0.3, desc, 9, MUTED, False, PP_ALIGN.CENTER)

t(slide, 0.8, 6.7, 11.5, 0.4, "All services read the same Iceberg tables on S3. No data duplication. No vendor lock-in.", 13, ACCENT, False, PP_ALIGN.CENTER)

# ==================== SLIDE 3: THE PROBLEM [INTERNAL] ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
line(slide, 0, 0, 13.333, 0.04, INTERNAL_RED)
tag(slide, 11.8, 0.2, "INTERNAL", INTERNAL_RED)
t(slide, 0.8, 0.3, 10, 0.6, "The Databricks-Only Problem", 28, WHITE, True)
t(slide, 0.8, 0.85, 10, 0.4, "What happens when one vendor owns everything — and why customers are feeling the pain", 14, MUTED)

# Problem cards
problems = [
    ("💰 Premium Storage Costs", "All data imported into Databricks at premium rates. At petabyte scale, the gap vs S3 is millions per year.", INTERNAL_RED),
    ("🔒 Vendor Lock-In", "Migrating away means rewriting pipelines, notebooks, and integrations. Separate identity, networking, billing.", INTERNAL_RED),
    ("🤖 AI Limitations", "No native Bedrock connection. Models purchased through Databricks marketplace. Limited model selection.", INTERNAL_RED),
    ("📊 Weak BI", "Databricks dashboards are basic, not embeddable, per-user pricing. Every viewer needs a Databricks license.", INTERNAL_RED),
    ("🔍 No Search Engine", "No full-text or vector search. Must bolt on third-party (Pinecone, Elastic) for RAG — adding cost and complexity.", INTERNAL_RED),
    ("📈 Linear Cost Scaling", "Scaling data = scaling Databricks spend linearly. Compute + storage coupled. No serverless option for SQL.", INTERNAL_RED),
]

for i, (title, body, color) in enumerate(problems):
    col = i % 3
    row = i // 3
    x = 0.8 + col * 4.0
    y = 1.5 + row * 2.7
    cd(slide, x, y, 3.7, 2.4, color)
    t(slide, x + 0.2, y + 0.1, 3.3, 0.3, title, 13, color, True)
    t(slide, x + 0.2, y + 0.5, 3.3, 1.5, body, 11, LIGHT)

t(slide, 0.8, 6.8, 11.5, 0.4, "This is the competitive landscape. The customer-facing story is about adding value, not attacking Databricks.", 11, MUTED, False, PP_ALIGN.CENTER)

# ==================== SLIDE 4: S3 + ICEBERG FOUNDATION [CUSTOMER-FACING] ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
line(slide, 0, 0, 13.333, 0.04, TEAL)
tag(slide, 11.5, 0.2, "CUSTOMER-FACING", GREEN)
t(slide, 0.8, 0.3, 10, 0.6, "Amazon S3 + Apache Iceberg: The Open Foundation", 26, WHITE, True)
t(slide, 0.8, 0.85, 10, 0.4, "Store once. Query from anywhere. Scale without limits.", 14, MUTED)

# Center lake
cd(slide, 2.0, 1.5, 9.3, 1.5, SOFT_BLUE, RGBColor(0x15, 0x2A, 0x45))
t(slide, 2.2, 1.55, 9.0, 0.35, "🪣 Amazon S3 Data Lake + Apache Iceberg", 20, WHITE, True, PP_ALIGN.CENTER)
t(slide, 2.2, 1.95, 9.0, 0.5, "Open table format. ACID transactions. Time travel. Schema evolution.\nAny engine reads it — Redshift, Athena, Spark, Databricks, Trino, Snowflake, Flink.", 12, SOFT_BLUE, False, PP_ALIGN.CENTER)
t(slide, 2.2, 2.5, 9.0, 0.3, "$0.023/GB/month  •  Infinite scale  •  Zero vendor lock-in", 11, ACCENT, True, PP_ALIGN.CENTER)

# Stats
stats = [("90%", "Storage cost\nreduction vs import"), ("8+", "Engines read\nIceberg natively"), ("0", "Data migration\nto switch tools"), ("∞", "Scale without\ncluster limits")]
for i, (num, lbl) in enumerate(stats):
    x = 1.5 + i * 2.8
    t(slide, x, 3.3, 2.0, 0.4, num, 28, ACCENT, True, PP_ALIGN.CENTER)
    t(slide, x, 3.8, 2.0, 0.5, lbl, 10, MUTED, False, PP_ALIGN.CENTER)

# Iceberg vs Delta (softer framing)
t(slide, 0.8, 4.5, 5, 0.3, "Why Iceberg over Delta Lake?", 14, WHITE, True)
comparisons = [
    ("Governance", "Apache Foundation (community-driven)", "Databricks (vendor-managed)"),
    ("Engine support", "Spark, Trino, Flink, Redshift, Athena, Snowflake, Presto", "Optimized for Databricks runtime"),
    ("AWS integration", "First-class in Glue, Athena, Redshift, EMR", "Requires Databricks or UniForm bridge"),
    ("Portability", "Read/write from any engine freely", "Best experience within Databricks"),
]
t(slide, 0.8, 4.9, 2.0, 0.25, "", 9, MUTED, True)
t(slide, 2.8, 4.9, 4.5, 0.25, "Apache Iceberg", 9, GREEN, True)
t(slide, 7.5, 4.9, 4.5, 0.25, "Delta Lake", 9, PURPLE, True)
for i, (label, ice, delta) in enumerate(comparisons):
    y = 5.2 + i * 0.35
    t(slide, 0.8, y, 2.0, 0.25, label, 9, MUTED)
    t(slide, 2.8, y, 4.5, 0.25, ice, 9, SOFT_GREEN)
    t(slide, 7.5, y, 4.5, 0.25, delta, 9, PURPLE)

# ==================== SLIDE 5: AWS SERVICES OVERVIEW [CUSTOMER-FACING] ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
line(slide, 0, 0, 13.333, 0.04, TEAL)
tag(slide, 11.5, 0.2, "CUSTOMER-FACING", GREEN)
t(slide, 0.8, 0.3, 10, 0.6, "AWS Analytics Services That Complement Databricks", 26, WHITE, True)
t(slide, 0.8, 0.85, 10, 0.4, "Each service is purpose-built for a specific workload — all reading from the same S3 data lake", 14, MUTED)

services = [
    ("📊 Amazon Redshift", "SQL Analytics & AI Queries", SOFT_BLUE,
     ["Serverless — pay per query", "Native Bedrock Knowledge Base source", "Zero-ETL from Aurora, DynamoDB", "Queries S3/Iceberg via Spectrum"]),
    ("🔍 Amazon OpenSearch", "Search, Vectors & RAG", TEAL,
     ["Full-text + vector hybrid search", "Native Bedrock vector store", "Serverless — scales to zero", "Powers RAG applications"]),
    ("📈 Amazon QuickSight", "Embedded BI & Dashboards", GREEN,
     ["Natural language queries (Q)", "Embeddable in any web app", "Per-session: $0.30 vs $70/user", "SPICE in-memory engine"]),
    ("🧪 SageMaker Studio", "ML & AI Development", ACCENT,
     ["Unified IDE: SQL, Python, Spark, ML", "Native Bedrock integration", "Full MLOps pipeline", "Fine-tune foundation models"]),
    ("🤖 Amazon Bedrock", "Generative AI Layer", AWS_ORANGE,
     ["100+ foundation models", "Knowledge Bases across all data", "Queries Redshift + OpenSearch + S3", "One API, one bill"]),
    ("🔄 Athena + Glue + EMR", "Data Processing & Catalog", MUTED,
     ["Serverless SQL on Iceberg", "Data catalog & governance", "Spark at scale (EMR)", "Lake Formation permissions"]),
]

for i, (name, role, color, features) in enumerate(services):
    col = i % 3
    row = i // 3
    x = 0.8 + col * 4.0
    y = 1.5 + row * 2.8
    cd(slide, x, y, 3.7, 2.5, color)
    t(slide, x + 0.15, y + 0.1, 3.4, 0.3, name, 13, color, True)
    t(slide, x + 0.15, y + 0.4, 3.4, 0.25, role, 10, MUTED)
    txBox = t(slide, x + 0.15, y + 0.75, 3.4, 0.25, "✓ " + features[0], 10, LIGHT)
    for f in features[1:]:
        ap(txBox.text_frame, "✓ " + f, 10, LIGHT, sb=3)

# ==================== SLIDE 6: COMPETITIVE GAPS [INTERNAL] ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
line(slide, 0, 0, 13.333, 0.04, INTERNAL_RED)
tag(slide, 11.8, 0.2, "INTERNAL", INTERNAL_RED)
t(slide, 0.8, 0.3, 10, 0.6, "Where AWS Wins vs. Databricks — Service by Service", 26, WHITE, True)
t(slide, 0.8, 0.85, 10, 0.4, "Use these talking points internally. Lead with customer value externally.", 14, MUTED)

gaps = [
    ("Redshift vs Databricks SQL", "Serverless pay-per-query. Native Bedrock KB source. Zero-ETL from Aurora/DynamoDB. Databricks SQL requires always-on cluster compute.", SOFT_BLUE),
    ("OpenSearch vs ???", "Databricks has no search engine at all. Customers must bolt on Pinecone or Elastic for RAG — extra cost, extra integration, extra vendor.", TEAL),
    ("QuickSight vs DB Dashboards", "Embeddable, per-session pricing ($0.30 vs $70/user/mo), natural language Q. Databricks dashboards are basic and not embeddable.", GREEN),
    ("Bedrock vs DB Marketplace", "100+ foundation models. Native KB integration with Redshift + OpenSearch + S3. Databricks has no native Bedrock connection — must export stale copies.", AWS_ORANGE),
    ("SageMaker vs MLflow", "Full MLOps: train, tune, deploy, monitor. Native Bedrock fine-tuning. Databricks MLflow is strong but siloed from AWS AI services.", ACCENT),
    ("S3/Iceberg vs Delta Lake", "Apache Foundation governed. $0.023/GB vs premium Delta storage. 8+ engines read natively. Delta is optimized for Databricks runtime only.", SOFT_GREEN),
]

for i, (title, body, color) in enumerate(gaps):
    col = i % 2
    row = i // 2
    x = 0.8 + col * 6.2
    y = 1.5 + row * 1.8
    cd(slide, x, y, 5.8, 1.6, color)
    t(slide, x + 0.2, y + 0.1, 5.4, 0.3, title, 13, color, True)
    t(slide, x + 0.2, y + 0.45, 5.4, 0.9, body, 10, LIGHT)

# ==================== SLIDE 7: MULTI-VENDOR ARCHITECTURE [CUSTOMER-FACING] ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
line(slide, 0, 0, 13.333, 0.04, TEAL)
tag(slide, 11.5, 0.2, "CUSTOMER-FACING", GREEN)
t(slide, 0.8, 0.3, 10, 0.6, "Open Architecture: Best Service for Every Workload", 26, WHITE, True)
t(slide, 0.8, 0.85, 10, 0.4, "Iceberg on S3 makes your data accessible to every platform. No single vendor owns it.", 14, MUTED)

# Foundation
cd(slide, 1.5, 1.4, 10.3, 0.7, SOFT_BLUE, RGBColor(0x15, 0x2A, 0x45))
t(slide, 1.7, 1.45, 10.0, 0.3, "🧊 Apache Iceberg Tables on Amazon S3", 16, WHITE, True, PP_ALIGN.CENTER)
t(slide, 1.7, 1.75, 10.0, 0.25, "One copy of data. Open format. Readable by every platform simultaneously.", 11, SOFT_BLUE, False, PP_ALIGN.CENTER)

t(slide, 6.2, 2.2, 1.0, 0.3, "↓", 16, SOFT_BLUE, True, PP_ALIGN.CENTER)

# Three vendor cards — softer framing
vendors = [
    ("🔥 Databricks", "Data Engineering & Spark", PURPLE,
     ["Large-scale Spark ETL & pipelines", "Data science notebooks", "Complex transformations", "Existing Spark workloads"]),
    ("☁️ AWS Analytics", "AI, Search, BI & Real-Time", TEAL,
     ["Redshift: SQL + Bedrock AI", "OpenSearch: Search + RAG", "QuickSight: Embedded BI", "Bedrock: GenAI across all data"]),
    ("🛡️ Palantir Foundry", "Operational Intelligence", PAL_GREEN,
     ["Ontology-based applications", "Cross-domain data fusion", "Mission-critical workflows", "Reads Iceberg from S3"]),
]

for i, (name, role, color, items) in enumerate(vendors):
    x = 0.8 + i * 4.2
    cd(slide, x, 2.6, 3.8, 3.2, color)
    t(slide, x + 0.2, 2.7, 3.4, 0.3, name, 15, color, True, PP_ALIGN.CENTER)
    t(slide, x + 0.2, 3.0, 3.4, 0.25, role, 10, MUTED, False, PP_ALIGN.CENTER)
    txBox = t(slide, x + 0.2, 3.4, 3.4, 0.25, "→ " + items[0], 10, LIGHT)
    for item in items[1:]:
        ap(txBox.text_frame, "→ " + item, 10, LIGHT, sb=4)

# Key message
cd(slide, 0.8, 6.1, 11.7, 0.8, ACCENT)
t(slide, 1.0, 6.15, 11.3, 0.3, "The data stays in S3. The format stays open. The choice stays yours.", 14, ACCENT, True, PP_ALIGN.CENTER)
t(slide, 1.0, 6.5, 11.3, 0.25, "Add or remove platforms without data migration. Every vendor reads the same Iceberg tables.", 11, LIGHT, False, PP_ALIGN.CENTER)

# ==================== SLIDE 8: USE CASE MATRIX [CUSTOMER-FACING] ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
line(slide, 0, 0, 13.333, 0.04, TEAL)
tag(slide, 11.5, 0.2, "CUSTOMER-FACING", GREEN)
t(slide, 0.8, 0.3, 10, 0.6, "Best Service for Each Workload", 26, WHITE, True)
t(slide, 0.8, 0.85, 10, 0.4, "Align the right tool to each use case — data stays open on S3/Iceberg", 14, MUTED)

# Headers
hy = 1.3
t(slide, 0.5, hy, 2.8, 0.3, "USE CASE", 9, MUTED, True)
t(slide, 3.5, hy, 2.5, 0.3, "DATABRICKS", 9, PURPLE, True)
t(slide, 6.2, hy, 2.8, 0.3, "AWS ANALYTICS", 9, TEAL, True)
t(slide, 9.2, hy, 2.3, 0.3, "PALANTIR", 9, PAL_GREEN, True)
t(slide, 11.7, hy, 1.5, 0.3, "BEST FIT", 9, ACCENT, True)
line(slide, 0.5, 1.55, 12.3, Pt(1), RGBColor(0x33, 0x44, 0x55))

matrix = [
    ("Spark ETL & Data Eng.", "✓ Native strength", "EMR / Glue", "—", "Databricks", PURPLE),
    ("SQL Analytics / DW", "SQL (cluster-based)", "✓ Redshift Serverless", "—", "Redshift", TEAL),
    ("Generative AI / RAG", "Limited models", "✓ Bedrock + OpenSearch", "AIP (ontology)", "Bedrock", TEAL),
    ("Full-Text & Vector Search", "No native engine", "✓ OpenSearch hybrid", "—", "OpenSearch", TEAL),
    ("Business Intelligence", "Basic dashboards", "✓ QuickSight embedded", "Foundry apps", "QuickSight", TEAL),
    ("Custom ML Training", "✓ MLflow strong", "✓ SageMaker + Bedrock", "Foundry ML", "Both", ACCENT),
    ("Operational Decisions", "—", "Step Functions", "✓ Foundry / AIP", "Palantir", PAL_GREEN),
    ("Near Real-Time Ingest", "Custom streaming", "✓ Redshift zETL", "Connectors", "zETL", TEAL),
]

for i, (uc, db, aws, pal, best, best_color) in enumerate(matrix):
    y = 1.7 + i * 0.58
    t(slide, 0.5, y, 2.8, 0.25, uc, 10, WHITE, True)
    t(slide, 3.5, y, 2.5, 0.25, db, 9, LIGHT)
    t(slide, 6.2, y, 2.8, 0.25, aws, 9, LIGHT)
    t(slide, 9.2, y, 2.3, 0.25, pal, 9, LIGHT)
    t(slide, 11.7, y, 1.5, 0.25, best, 10, best_color, True)

t(slide, 0.8, 6.5, 11.5, 0.4, "All platforms read the same Iceberg tables on S3. No data duplication. Change your mind later without migrating data.", 12, ACCENT, False, PP_ALIGN.CENTER)

# ==================== SLIDE 9: BEDROCK AI LAYER [CUSTOMER-FACING] ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
line(slide, 0, 0, 13.333, 0.04, TEAL)
tag(slide, 11.5, 0.2, "CUSTOMER-FACING", GREEN)
t(slide, 0.8, 0.3, 10, 0.6, "Amazon Bedrock: Unified AI Across All Your Data", 26, WHITE, True)
t(slide, 0.8, 0.85, 10, 0.4, "One AI layer that queries structured, unstructured, and real-time data natively", 14, MUTED)

# Bedrock center
cd(slide, 2.0, 1.4, 9.3, 0.8, AWS_ORANGE, RGBColor(0x3D, 0x2A, 0x0A))
t(slide, 2.2, 1.45, 9.0, 0.3, "Amazon Bedrock Knowledge Bases", 18, AWS_ORANGE, True, PP_ALIGN.CENTER)
t(slide, 2.2, 1.8, 9.0, 0.25, "Natural language interface to all your data through one unified API", 12, LIGHT, False, PP_ALIGN.CENTER)

# Sources
sources = [
    ("📊 Redshift", "Structured Data", "Bedrock generates SQL, queries live tables.\nReal-time answers. Zero ETL.", "NATIVE", GREEN),
    ("🔍 OpenSearch", "Vector Store / RAG", "Semantic search across millions of docs.\nVector embeddings for retrieval.", "NATIVE", GREEN),
    ("🪣 S3 Data Lake", "Documents & Files", "PDFs, images, CSVs ingested directly.\nAuto-chunked and indexed.", "NATIVE", GREEN),
    ("🧪 SageMaker", "Custom Models", "Fine-tuned models as Bedrock endpoints.\nCustom ML alongside foundation models.", "DIRECT", SOFT_BLUE),
]
for i, (name, stype, desc, tag_text, tag_color) in enumerate(sources):
    x = 0.8 + i * 3.1
    cd(slide, x, 2.5, 2.9, 2.8, tag_color)
    t(slide, x + 0.1, 2.55, 2.7, 0.3, name, 13, WHITE, True, PP_ALIGN.CENTER)
    t(slide, x + 0.1, 2.85, 2.7, 0.2, stype, 10, MUTED, False, PP_ALIGN.CENTER)
    t(slide, x + 0.1, 3.2, 2.7, 0.8, desc, 9, LIGHT, False, PP_ALIGN.CENTER)
    t(slide, x + 0.1, 4.6, 2.7, 0.25, tag_text + " SOURCE", 8, tag_color, True, PP_ALIGN.CENTER)

# Bottom comparison — softer
cd(slide, 0.8, 5.6, 5.8, 1.2, INTERNAL_RED)
t(slide, 1.0, 5.65, 5.4, 0.25, "Without AWS-native AI", 12, INTERNAL_RED, True)
t(slide, 1.0, 5.95, 5.4, 0.6, "Must export data to S3, then Bedrock reads stale copies. AI models limited to one vendor's marketplace. No unified query across data types.", 10, LIGHT)

cd(slide, 7.0, 5.6, 5.8, 1.2, GREEN)
t(slide, 7.2, 5.65, 5.4, 0.25, "With Bedrock + AWS Analytics", 12, GREEN, True)
t(slide, 7.2, 5.95, 5.4, 0.6, "One Knowledge Base queries Redshift + OpenSearch + S3 simultaneously. 100+ models. One question gets answers from all data sources.", 10, LIGHT)

# ==================== SLIDE 10: CO-EXIST SUMMARY [CUSTOMER-FACING] ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
line(slide, 0, 0, 13.333, 0.04, TEAL)
tag(slide, 11.5, 0.2, "CUSTOMER-FACING", GREEN)
t(slide, 0.8, 0.3, 10, 0.6, "Why Co-Exist: Six Reasons to Complement Databricks with AWS", 24, WHITE, True)
t(slide, 0.8, 0.85, 10, 0.4, "Keep what works. Add what's missing. Let the data stay open.", 14, MUTED)

cards = [
    ("💰", "Cost at Scale", "S3 at $0.023/GB. Serverless Redshift, OpenSearch, QuickSight — pay only when you query.", "Up to 90% storage savings", TEAL),
    ("🤖", "Native AI", "Redshift, OpenSearch, S3 are native Bedrock sources. One AI layer, zero ETL.", "Zero custom integration", ACCENT),
    ("🔓", "No Lock-In", "Data in S3/Iceberg. Any engine reads it. Swap services without migration.", "Open formats, open choice", GREEN),
    ("📈", "Best Tool Per Job", "Redshift for SQL. OpenSearch for search. QuickSight for BI. Bedrock for GenAI.", "5 services, one data lake", SOFT_BLUE),
    ("🔒", "Unified Security", "One IAM, one VPC, one Lake Formation. Data never leaves your AWS account.", "Single governance model", TEAL),
    ("🤝", "Co-Existence", "Keep Databricks for Spark. Add AWS-native for analytics, search, BI, AI. Better together.", "Databricks + AWS", PURPLE),
]

for i, (icon, title, body, stat, color) in enumerate(cards):
    col = i % 3
    row = i // 3
    x = 0.8 + col * 4.0
    y = 1.4 + row * 2.7
    cd(slide, x, y, 3.7, 2.4, color)
    t(slide, x + 0.15, y + 0.1, 3.4, 0.3, icon, 20, WHITE)
    t(slide, x + 0.15, y + 0.4, 3.4, 0.3, title, 14, color, True)
    t(slide, x + 0.15, y + 0.75, 3.4, 0.8, body, 10, LIGHT)
    line(slide, x + 0.15, y + 1.7, 3.4, Pt(0.5), RGBColor(0x33, 0x44, 0x55))
    t(slide, x + 0.15, y + 1.8, 3.4, 0.3, stat, 10, color, True)

# CTA
cd(slide, 0.8, 6.6, 11.7, 0.7, ACCENT, RGBColor(0x0A, 0x1A, 0x2E))
t(slide, 1.0, 6.65, 8, 0.3, "Ready to see it in action?", 16, WHITE, True)
t(slide, 1.0, 6.95, 8, 0.25, "Start with a POC: S3 data lake → Redshift + OpenSearch + Bedrock. Results in days, not months.", 11, LIGHT)
t(slide, 10.0, 6.75, 2.5, 0.4, "Let's Build a POC →", 13, ACCENT, True, PP_ALIGN.CENTER)

# ==================== SAVE ====================
output_path = "AWS-Databricks-CoExist-v2.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
