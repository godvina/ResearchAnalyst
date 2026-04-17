#!/usr/bin/env python3
"""Convert AWS Analytics vs Databricks HTML presentation to PowerPoint (15 slides)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Widescreen dimensions ──
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Color palette ──
C_DARK_BG    = RGBColor(0x23, 0x2F, 0x3E)
C_DARKER_BG  = RGBColor(0x1A, 0x23, 0x32)
C_DARKEST_BG = RGBColor(0x0D, 0x1B, 0x2A)
C_ORANGE     = RGBColor(0xFF, 0x99, 0x00)
C_WHITE      = RGBColor(0xE2, 0xE8, 0xF0)
C_GRAY       = RGBColor(0xA0, 0xAE, 0xC0)
C_GRAY2      = RGBColor(0x71, 0x80, 0x96)
C_GRAY3      = RGBColor(0xCB, 0xD5, 0xE0)
C_RED        = RGBColor(0xDC, 0x26, 0x26)
C_RED_LIGHT  = RGBColor(0xFC, 0x81, 0x81)
C_RED_LIGHTER= RGBColor(0xFC, 0xA5, 0xA5)
C_RED_PALE   = RGBColor(0xFE, 0xCA, 0xCA)
C_GREEN      = RGBColor(0x16, 0xA3, 0x4A)
C_GREEN_LIGHT= RGBColor(0x86, 0xEF, 0xAC)
C_GREEN_PALE = RGBColor(0xD1, 0xFA, 0xE5)
C_BLUE       = RGBColor(0x3B, 0x82, 0xF6)
C_BLUE_LIGHT = RGBColor(0x60, 0xA5, 0xFA)
C_BLUE_PALE  = RGBColor(0x93, 0xC5, 0xFD)
C_BLUE_DARK  = RGBColor(0x1E, 0x40, 0xAF)
C_PURPLE     = RGBColor(0x8B, 0x5C, 0xF6)
C_PURPLE_LT  = RGBColor(0xA7, 0x8B, 0xFA)
C_PURPLE_PALE= RGBColor(0xC4, 0xB5, 0xFD)
C_PALANTIR   = RGBColor(0x22, 0xC5, 0x5E)
C_PALANTIR_LT= RGBColor(0x4A, 0xDE, 0x80)
C_YELLOW     = RGBColor(0xFC, 0xD3, 0x4D)
C_PURE_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK      = RGBColor(0x00, 0x00, 0x00)
C_SLATE      = RGBColor(0x33, 0x41, 0x55)
C_SLATE_LT   = RGBColor(0x47, 0x55, 0x69)
C_SLATE_DK   = RGBColor(0x1E, 0x29, 0x3B)
C_BG_CARD    = RGBColor(0x1A, 0x1A, 0x2E)

# ── Helper functions ──

def add_bg(slide, color):
    """Fill slide background with solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def tb(slide, left, top, width, height, text, size=12, color=C_WHITE, bold=False, align=PP_ALIGN.LEFT, font_name='Segoe UI'):
    """Add a text box and return it."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def ap(text_frame, text, size=12, color=C_WHITE, bold=False, align=PP_ALIGN.LEFT, space_before=Pt(0), space_after=Pt(0), font_name='Segoe UI'):
    """Add a paragraph to an existing text frame."""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    p.space_before = space_before
    p.space_after = space_after
    return p


def card(slide, left, top, width, height, border_color=C_ORANGE, fill_color=None):
    """Add a rounded rectangle card shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    # Smaller corner radius
    shape.adjustments[0] = 0.05
    return shape

def bar(slide, left, top, width, height, color):
    """Add a rectangular bar."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def footer(slide, left_text="Amazon Web Services", right_text=""):
    """Add footer bar."""
    bar(slide, 0, 7.1, 13.333, 0.4, RGBColor(0x1A, 0x1F, 0x2B))
    tb(slide, 0.6, 7.15, 3, 0.3, left_text, 10, C_ORANGE, bold=True)
    tb(slide, 4, 7.15, 8, 0.3, right_text, 9, C_GRAY2, align=PP_ALIGN.RIGHT)

def top_bar(slide, color=C_ORANGE):
    """Add top accent bar."""
    bar(slide, 0, 0, 13.333, 0.06, color)

def dark_header(slide, title, badge_text=None, y=0):
    """Add dark header bar with title and optional badge."""
    bar(slide, 0, y, 13.333, 0.55, C_DARK_BG)
    tb(slide, 0.6, y + 0.05, 9, 0.45, title, 20, C_PURE_WHITE, bold=True)
    if badge_text:
        c = card(slide, 10.5, y + 0.1, 2.2, 0.35, C_ORANGE, C_ORANGE)
        c.text_frame.word_wrap = True
        p = c.text_frame.paragraphs[0]
        p.text = badge_text
        p.font.size = Pt(9)
        p.font.color.rgb = C_DARK_BG
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

def bottom_callout(slide, icon, highlight, text, y=6.6):
    """Add bottom callout bar."""
    bar(slide, 0, y, 13.333, 0.5, C_DARK_BG)
    tb(slide, 0.4, y + 0.05, 0.4, 0.4, icon, 16, C_PURE_WHITE)
    bx = tb(slide, 0.9, y + 0.05, 11.5, 0.4, "", 11, C_PURE_WHITE, bold=True)
    tf = bx.text_frame
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = highlight + " "
    run1.font.color.rgb = C_ORANGE
    run1.font.bold = True
    run1.font.size = Pt(11)
    run2 = p.add_run()
    run2.text = text
    run2.font.color.rgb = C_PURE_WHITE
    run2.font.bold = False
    run2.font.size = Pt(11)


# ══════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════

def slide_01_title(prs):
    """Slide 1: Title / Problem Statement (dark bg)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, C_DARK_BG)
    top_bar(slide)

    # Left side - Title
    tb(slide, 0.6, 0.5, 6.5, 1.0,
       "Break Free from Vendor Lock-In:\nThe AWS Analytics Advantage",
       28, C_ORANGE, bold=True)
    tb(slide, 0.6, 1.8, 6.5, 1.2,
       "Build on Amazon S3 as your open data lake. Access best-of-breed analytics, search, BI, ML, and generative AI services — all natively integrated, all pay-per-use.",
       14, C_GRAY, align=PP_ALIGN.LEFT)

    # Tagline with left border (simulated with a thin bar)
    bar(slide, 0.6, 3.3, 0.04, 1.0, C_ORANGE)
    bx = tb(slide, 0.8, 3.3, 6, 1.0, "", 13, C_ORANGE, bold=True)
    tf = bx.text_frame
    tf.paragraphs[0].text = "Keep Databricks for what it does well."
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.color.rgb = C_ORANGE
    tf.paragraphs[0].font.bold = True
    ap(tf, "Stop forcing everything through one vendor.", 13, C_ORANGE, True)
    ap(tf, "Let AWS services work together natively.", 13, C_ORANGE, True)

    # Right side - Problem box
    c = card(slide, 7.5, 0.6, 5.2, 5.5, C_ORANGE, RGBColor(0x2A, 0x36, 0x48))
    tf = c.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "THE DATABRICKS-ONLY PROBLEM"
    p.font.size = Pt(10)
    p.font.color.rgb = C_ORANGE
    p.font.bold = True
    p.space_after = Pt(12)

    problems = [
        "All data imported into Databricks = expensive storage at premium rates",
        "AI/ML models purchased and maintained through Databricks marketplace",
        "No native connection to Amazon Bedrock, OpenSearch, QuickSight, or SageMaker",
        "Scaling data = scaling Databricks spend linearly (compute + storage)",
        "Separate identity, networking, billing, and support from your AWS environment",
        "Vendor lock-in: migrating away means rewriting pipelines, notebooks, and integrations",
    ]
    for prob in problems:
        ap(tf, "✗  " + prob, 11, C_GRAY3, space_before=Pt(6))

    footer(slide, "Amazon Web Services", "AWS Analytics Suite | Open Data Lake Strategy")


def slide_02_lockin(prs):
    """Slide 2: Lock-In vs Open Architecture (WHITE bg)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, C_PURE_WHITE)

    dark_header(slide, "Vendor Lock-In vs. Open Data Lake Architecture", "SIDE BY SIDE")

    # ── LEFT COLUMN: Locked ──
    c_left = card(slide, 0.4, 0.7, 6.1, 5.6, C_RED_LIGHTER, RGBColor(0xFE, 0xF2, 0xF2))
    tb(slide, 0.6, 0.8, 3.5, 0.3, "Databricks-Only Approach", 14, RGBColor(0x1E, 0x29, 0x3B), bold=True)
    c_badge = card(slide, 4.3, 0.82, 1.0, 0.25, C_RED_LIGHTER, C_RED_LIGHTER)
    c_badge.text_frame.paragraphs[0].text = "Locked In"
    c_badge.text_frame.paragraphs[0].font.size = Pt(8)
    c_badge.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x7F, 0x1D, 0x1D)
    c_badge.text_frame.paragraphs[0].font.bold = True
    c_badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    locked_items = [
        ("📥 Import ALL Data into Databricks", "Premium storage rates, data duplication"),
        ("⚙️ Databricks Compute for Everything", "ETL, analytics, ML — all one vendor"),
        ("🤖 Databricks AI/ML Marketplace", "Models purchased through Databricks"),
        ("📊 Databricks Dashboards", "Limited BI, no embedded analytics"),
        ("👤 End Users", "Single vendor dependency for everything"),
    ]
    y = 1.2
    for label, note in locked_items:
        c_item = card(slide, 0.7, y, 5.5, 0.35, C_RED_LIGHTER, RGBColor(0xFE, 0xE2, 0xE2))
        c_item.text_frame.paragraphs[0].text = label
        c_item.text_frame.paragraphs[0].font.size = Pt(10)
        c_item.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x99, 0x1B, 0x1B)
        c_item.text_frame.paragraphs[0].font.bold = True
        c_item.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        tb(slide, 0.7, y + 0.37, 5.5, 0.2, note, 8, C_RED, align=PP_ALIGN.CENTER)
        tb(slide, 2.8, y + 0.55, 0.6, 0.2, "↓", 14, C_RED, align=PP_ALIGN.CENTER)
        y += 0.78

    # ── RIGHT COLUMN: Open ──
    c_right = card(slide, 6.8, 0.7, 6.1, 5.6, C_GREEN_LIGHT, RGBColor(0xF0, 0xFD, 0xF4))
    tb(slide, 7.0, 0.8, 3.5, 0.3, "AWS Open Data Lake", 14, RGBColor(0x1E, 0x29, 0x3B), bold=True)
    c_badge2 = card(slide, 10.5, 0.82, 1.3, 0.25, C_GREEN_LIGHT, C_GREEN_LIGHT)
    c_badge2.text_frame.paragraphs[0].text = "Open & Flexible"
    c_badge2.text_frame.paragraphs[0].font.size = Pt(8)
    c_badge2.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x14, 0x53, 0x2D)
    c_badge2.text_frame.paragraphs[0].font.bold = True
    c_badge2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # S3 box
    y = 1.2
    c_s3 = card(slide, 7.1, y, 5.5, 0.35, C_BLUE, RGBColor(0xDB, 0xEA, 0xFE))
    c_s3.text_frame.paragraphs[0].text = "🪣 Amazon S3 Data Lake (Open Formats)"
    c_s3.text_frame.paragraphs[0].font.size = Pt(10)
    c_s3.text_frame.paragraphs[0].font.color.rgb = C_BLUE_DARK
    c_s3.text_frame.paragraphs[0].font.bold = True
    c_s3.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tb(slide, 7.1, y + 0.37, 5.5, 0.2, "Lowest cost storage, open formats, infinite scale", 8, C_GREEN, align=PP_ALIGN.CENTER)
    tb(slide, 9.3, y + 0.55, 0.6, 0.2, "↓", 14, C_GREEN, align=PP_ALIGN.CENTER)
    y += 0.78

    # 2x2 AWS services grid
    svc_pairs = [
        [("📊 Redshift", "SQL Analytics"), ("🔍 OpenSearch", "Search & RAG")],
        [("📈 QuickSight", "BI & Dashboards"), ("🧪 SageMaker", "ML Studio")],
    ]
    for row in svc_pairs:
        x = 7.1
        for label, sub in row:
            c_svc = card(slide, x, y, 2.65, 0.45, C_GREEN, RGBColor(0xD1, 0xFA, 0xE5))
            tf = c_svc.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = label
            tf.paragraphs[0].font.size = Pt(9)
            tf.paragraphs[0].font.color.rgb = RGBColor(0x06, 0x5F, 0x46)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            ap(tf, sub, 8, RGBColor(0x06, 0x5F, 0x46), align=PP_ALIGN.CENTER)
            x += 2.85
        y += 0.55

    tb(slide, 9.3, y, 0.6, 0.2, "↓", 14, C_GREEN, align=PP_ALIGN.CENTER)
    y += 0.25

    # Bedrock
    c_bed = card(slide, 7.1, y, 5.5, 0.35, C_DARK_BG, C_DARK_BG)
    c_bed.text_frame.paragraphs[0].text = "🤖 Amazon Bedrock — Unified AI Layer"
    c_bed.text_frame.paragraphs[0].font.size = Pt(10)
    c_bed.text_frame.paragraphs[0].font.color.rgb = C_ORANGE
    c_bed.text_frame.paragraphs[0].font.bold = True
    c_bed.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tb(slide, 7.1, y + 0.37, 5.5, 0.2, "All services feed into Bedrock natively", 8, C_GREEN, align=PP_ALIGN.CENTER)
    tb(slide, 9.3, y + 0.55, 0.6, 0.2, "↓", 14, C_GREEN, align=PP_ALIGN.CENTER)
    y += 0.78

    # End Users
    c_usr = card(slide, 7.1, y, 5.5, 0.35, C_PURPLE, RGBColor(0xF3, 0xE8, 0xFF))
    c_usr.text_frame.paragraphs[0].text = "👤 End Users"
    c_usr.text_frame.paragraphs[0].font.size = Pt(10)
    c_usr.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x6B, 0x21, 0xA8)
    c_usr.text_frame.paragraphs[0].font.bold = True
    c_usr.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tb(slide, 7.1, y + 0.37, 5.5, 0.2, "Best tool for each job, one AWS bill", 8, C_GREEN, align=PP_ALIGN.CENTER)

    bottom_callout(slide, "💡", "Co-existence strategy:",
        "Keep Databricks for Spark workloads and data engineering. Add AWS-native services for analytics, search, BI, and AI — all reading from the same S3 data lake.",
        y=6.45)



def slide_03_s3_datalake(prs):
    """Slide 3: S3 Data Lake Foundation (dark bg)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, C_DARK_BG)
    top_bar(slide)

    # Header
    tb(slide, 0.6, 0.2, 10, 0.4, "Amazon S3: The Foundation That Changes Everything", 24, C_ORANGE, bold=True)
    tb(slide, 0.6, 0.65, 10, 0.3, "Store once in open formats. Query from any service. Scale without limits. Pay only for what you store.", 12, C_GRAY)

    # Center lake box
    c_lake = card(slide, 0.6, 1.1, 7.5, 2.0, C_BLUE_LIGHT, C_BLUE)
    tf = c_lake.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🪣"
    p.font.size = Pt(36)
    p.alignment = PP_ALIGN.CENTER
    ap(tf, "Amazon S3 Data Lake + Apache Iceberg", 18, C_PURE_WHITE, True, PP_ALIGN.CENTER)
    ap(tf, "Open table format on S3. ACID transactions, time travel, schema evolution.", 11, RGBColor(0xBF, 0xDB, 0xFE), align=PP_ALIGN.CENTER)
    ap(tf, "Your data is never locked in. Any engine reads it. $0.023/GB/month.", 11, RGBColor(0xBF, 0xDB, 0xFE), align=PP_ALIGN.CENTER)

    # Stats row
    stats = [("90%", "Cost Reduction"), ("∞", "Scale"), ("0", "Vendor Lock-In"), ("5+", "AWS Services")]
    sx = 1.0
    for num, lbl in stats:
        c_st = card(slide, sx, 3.25, 1.5, 0.65, C_BLUE_LIGHT, RGBColor(0x1A, 0x30, 0x50))
        tf = c_st.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.color.rgb = C_ORANGE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        ap(tf, lbl, 8, C_BLUE_PALE, align=PP_ALIGN.CENTER)
        sx += 1.7

    # Service chips grid (3x3)
    services = [
        ("📊", "Redshift", "SQL Analytics"),
        ("🔍", "OpenSearch", "Search & RAG"),
        ("📈", "QuickSight", "BI & Dashboards"),
        ("🧪", "SageMaker", "ML & AI Studio"),
        ("🤖", "Bedrock", "Generative AI"),
        ("🧊", "Iceberg Tables", "Open Table Format"),
        ("🔧", "Glue", "Data Catalog"),
        ("🛡️", "Lake Formation", "Governance"),
        ("🔄", "Athena", "Ad-hoc Query"),
    ]
    sx, sy = 0.6, 4.1
    for i, (icon, name, role) in enumerate(services):
        c_svc = card(slide, sx, sy, 1.9, 0.7, C_ORANGE, RGBColor(0x2A, 0x30, 0x3E))
        tf = c_svc.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = icon
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        ap(tf, name, 9, C_ORANGE, True, PP_ALIGN.CENTER)
        ap(tf, role, 8, C_GRAY, align=PP_ALIGN.CENTER)
        sx += 2.05
        if (i + 1) % 3 == 0:
            sx = 0.6
            sy += 0.85

    # Right side benefits
    benefits = [
        ("Cost at Scale", "S3 storage: $0.023/GB/month. Databricks Delta Lake storage: $0.06+/GB/month plus compute. At petabyte scale, the gap is millions per year."),
        ("Apache Iceberg: True Openness", "Iceberg tables on S3 give you ACID transactions, time travel, and schema evolution — readable by Spark, Trino, Redshift, Athena, Flink, and Databricks itself."),
        ("Decouple Storage from Compute", "Scale storage and compute independently. Don't pay for idle compute just to keep data accessible. S3 is always on, always available."),
        ("Native AWS Integration", "Every AWS analytics service reads S3 natively. No connectors, no ETL, no data movement. One copy of data, many consumers."),
    ]
    by = 1.1
    for title, body in benefits:
        c_b = card(slide, 8.5, by, 4.4, 1.2, C_ORANGE, RGBColor(0x2A, 0x36, 0x48))
        tf = c_b.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.color.rgb = C_ORANGE
        tf.paragraphs[0].font.bold = True
        ap(tf, body, 9, C_GRAY3, space_before=Pt(4))
        by += 1.35

    footer(slide, "Amazon Web Services", "S3 Data Lake | Store Once, Query from Anywhere")



def slide_04_iceberg(prs):
    """Slide 4: Apache Iceberg (dark blue bg)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, RGBColor(0x0C, 0x19, 0x29))
    top_bar(slide)

    # Header
    tb(slide, 0.6, 0.2, 10, 0.4, "Apache Iceberg on S3: Your Data Goes Everywhere", 22, C_BLUE_LIGHT, bold=True)
    tb(slide, 0.6, 0.6, 10, 0.3, "Open table format means zero lock-in. Every major engine reads Iceberg natively — including Databricks itself.", 11, RGBColor(0x94, 0xA3, 0xB8))

    # Center Iceberg box
    c_ice = card(slide, 0.6, 0.95, 7.8, 1.3, C_BLUE_LIGHT, RGBColor(0x1E, 0x3A, 0x5F))
    tf = c_ice.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "🧊  Apache Iceberg Tables on Amazon S3"
    tf.paragraphs[0].font.size = Pt(17)
    tf.paragraphs[0].font.color.rgb = C_PURE_WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ap(tf, "ACID transactions · Time travel · Schema evolution · Partition evolution · Hidden partitioning", 10, RGBColor(0xBF, 0xDB, 0xFE), align=PP_ALIGN.CENTER, space_before=Pt(4))
    ap(tf, "One copy of data in S3. Every engine reads it. No proprietary format. No conversion. No export.", 10, RGBColor(0xBF, 0xDB, 0xFE), align=PP_ALIGN.CENTER)

    # Engine label
    tb(slide, 0.6, 2.35, 7.8, 0.25, "ENGINES THAT READ ICEBERG NATIVELY", 9, C_BLUE_LIGHT, bold=True, align=PP_ALIGN.CENTER)

    # Engine grid 4x2
    engines = [
        ("📊", "Redshift", "Query Iceberg via Spectrum, zero-copy", True),
        ("🔄", "Athena", "Serverless SQL on Iceberg tables", True),
        ("⚡", "EMR / Spark", "Native Iceberg read/write at scale", True),
        ("📋", "Glue", "Iceberg catalog + ETL jobs", True),
        ("🔥", "Databricks", "Reads Iceberg via UniForm", False),
        ("🔺", "Trino / Presto", "Full Iceberg connector", False),
        ("❄️", "Snowflake", "External Iceberg table support", False),
        ("🌊", "Flink", "Streaming read/write to Iceberg", False),
    ]
    ex, ey = 0.6, 2.65
    for i, (icon, name, how, aws) in enumerate(engines):
        border = C_ORANGE if aws else C_BLUE_LIGHT
        c_eng = card(slide, ex, ey, 1.85, 0.7, border, RGBColor(0x14, 0x20, 0x35))
        tf = c_eng.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = icon
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        name_color = C_ORANGE if aws else C_BLUE_LIGHT
        ap(tf, name, 9, name_color, True, PP_ALIGN.CENTER)
        ap(tf, how, 7, RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.CENTER)
        ex += 1.95
        if (i + 1) % 4 == 0:
            ex = 0.6
            ey += 0.8

    # Iceberg vs Delta comparison
    c_vs = card(slide, 0.6, 4.35, 7.8, 1.8, C_RED_LIGHTER, RGBColor(0x14, 0x20, 0x35))
    tf = c_vs.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "Iceberg vs. Databricks Delta Lake — Format Comparison"
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.color.rgb = C_RED_LIGHTER
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].space_after = Pt(6)

    comparisons = [
        ("Governance", "Apache Foundation (open)", "Databricks (vendor-controlled)"),
        ("Engine support", "Spark, Trino, Flink, Redshift, Athena, Snowflake, Presto", "Best in Databricks; limited elsewhere"),
        ("AWS integration", "First-class in Glue, Athena, Redshift, EMR", "Requires Databricks runtime or UniForm"),
        ("Portability", "Move data between engines freely", "Tied to Databricks ecosystem"),
    ]
    for label, ice, delta in comparisons:
        ap(tf, f"{label}:  {ice}  vs.  {delta}", 9, C_GRAY3, space_before=Pt(3))

    # Right side benefits
    ice_benefits = [
        ("🔓 True Data Portability", "Your data in Iceberg format on S3 can be read by any engine today and any engine tomorrow."),
        ("⏪ Time Travel Queries", "Query data as it existed at any point in time. Roll back bad writes. Audit historical state."),
        ("📐 Schema Evolution", "Add, rename, drop, or reorder columns without rewriting data."),
        ("📊 Partition Evolution", "Change partitioning strategy without rewriting existing data."),
        ("💰 S3 Economics", "Iceberg tables live on S3 at $0.023/GB. No proprietary storage layer."),
    ]
    by = 0.95
    for title, body in ice_benefits:
        c_b = card(slide, 8.8, by, 4.1, 0.95, C_BLUE_LIGHT, RGBColor(0x14, 0x20, 0x35))
        tf = c_b.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.color.rgb = C_BLUE_LIGHT
        tf.paragraphs[0].font.bold = True
        ap(tf, body, 8, C_GRAY3, space_before=Pt(3))
        by += 1.08

    footer(slide, "Amazon Web Services", "Apache Iceberg on S3 | Open Format, Universal Access")


def slide_05_redshift_opensearch(prs):
    """Slide 5: Redshift + OpenSearch (WHITE bg)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, C_PURE_WHITE)

    dark_header(slide, "AWS Analytics Services: Redshift + OpenSearch", "DATA QUERY & SEARCH")

    # Redshift card
    c_rs = card(slide, 0.4, 0.7, 6.1, 5.5, C_BLUE_PALE, RGBColor(0xEF, 0xF6, 0xFF))
    tf = c_rs.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "📊  Amazon Redshift"
    tf.paragraphs[0].font.size = Pt(17)
    tf.paragraphs[0].font.color.rgb = C_BLUE_DARK
    tf.paragraphs[0].font.bold = True
    ap(tf, "Serverless SQL Analytics & AI Query Layer", 10, RGBColor(0x64, 0x74, 0x8B), space_before=Pt(2))
    ap(tf, "", 6, C_BLACK)
    ap(tf, "The fastest cloud data warehouse with native Bedrock integration. Users ask questions in plain English — Bedrock generates SQL and queries live Redshift data in real time.", 10, C_SLATE, space_before=Pt(6))
    rs_features = [
        "Native Bedrock Knowledge Base structured data source",
        "Serverless — pay per query, no cluster management",
        "Zero-ETL from Aurora, DynamoDB, and S3",
        "Queries S3 data lake directly via Redshift Spectrum",
        "Federated queries across RDS, Aurora, S3 in one SQL",
        "ML inference directly in SQL with Redshift ML",
    ]
    for f in rs_features:
        ap(tf, "✓  " + f, 9, C_SLATE_LT, space_before=Pt(3))
    ap(tf, "", 4, C_BLACK)
    ap(tf, "Databricks gap: No native Bedrock connection. Must ETL to S3, then Bedrock reads stale copies. No serverless SQL warehouse option with AI integration.", 9, C_RED, space_before=Pt(6), bold=True)

    # OpenSearch card
    c_os = card(slide, 6.8, 0.7, 6.1, 5.5, RGBColor(0xD8, 0xB4, 0xFE), RGBColor(0xFD, 0xF4, 0xFF))
    tf = c_os.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "🔍  Amazon OpenSearch"
    tf.paragraphs[0].font.size = Pt(17)
    tf.paragraphs[0].font.color.rgb = RGBColor(0x7C, 0x3A, 0xED)
    tf.paragraphs[0].font.bold = True
    ap(tf, "Full-Text Search, Vector Search & RAG", 10, RGBColor(0x64, 0x74, 0x8B), space_before=Pt(2))
    ap(tf, "", 6, C_BLACK)
    ap(tf, "Enterprise search and vector database for RAG applications. Powers Bedrock Knowledge Bases for semantic search across millions of documents with sub-second response.", 10, C_SLATE, space_before=Pt(6))
    os_features = [
        "Native Bedrock Knowledge Base vector store",
        "Full-text search + vector (hybrid) search in one engine",
        "Serverless option — scales to zero, pay per use",
        "k-NN vector search for semantic similarity",
        "Ingests directly from S3, Kinesis, DynamoDB",
        "Dashboards built-in for log and search analytics",
    ]
    for f in os_features:
        ap(tf, "✓  " + f, 9, C_SLATE_LT, space_before=Pt(3))
    ap(tf, "", 4, C_BLACK)
    ap(tf, "Databricks gap: No equivalent full-text + vector search engine. Must use third-party (Pinecone, Elastic) for RAG, adding cost and integration complexity.", 9, C_RED, space_before=Pt(6), bold=True)

    bottom_callout(slide, "🔗", "Together:",
        "Redshift handles structured data queries (SQL, analytics). OpenSearch handles unstructured search (documents, RAG). Both feed into Bedrock as native data sources.",
        y=6.4)



def slide_06_quicksight_sagemaker(prs):
    """Slide 6: QuickSight + SageMaker (WHITE bg)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, C_PURE_WHITE)

    dark_header(slide, "AWS Analytics Services: QuickSight + SageMaker Unified Studio", "BI & MACHINE LEARNING")

    # QuickSight card
    c_qs = card(slide, 0.4, 0.7, 6.1, 5.5, RGBColor(0xFD, 0xE0, 0x47), RGBColor(0xFE, 0xFC, 0xE8))
    tf = c_qs.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "📈  Amazon QuickSight"
    tf.paragraphs[0].font.size = Pt(17)
    tf.paragraphs[0].font.color.rgb = RGBColor(0xA1, 0x62, 0x07)
    tf.paragraphs[0].font.bold = True
    ap(tf, "Serverless BI with Generative AI (Q)", 10, RGBColor(0x64, 0x74, 0x8B), space_before=Pt(2))
    ap(tf, "", 6, C_BLACK)
    ap(tf, "Cloud-native BI with natural language querying (QuickSight Q). Embed interactive dashboards in any application. Pay per session — not per named user.", 10, C_SLATE, space_before=Pt(6))
    qs_features = [
        "QuickSight Q: ask questions in plain English, get charts",
        "Connects natively to Redshift, S3, Athena, RDS",
        "Embeddable dashboards in web apps and portals",
        "Per-session pricing: $0.30/session vs. $70/user/month",
        "SPICE in-memory engine for sub-second dashboards",
        "Paginated reports for compliance and regulatory needs",
    ]
    for f in qs_features:
        ap(tf, "✓  " + f, 9, C_SLATE_LT, space_before=Pt(3))
    ap(tf, "", 4, C_BLACK)
    ap(tf, "Databricks gap: Databricks SQL dashboards are basic, not embeddable, and require Databricks compute. No per-session pricing — every viewer needs a Databricks license.", 9, C_RED, space_before=Pt(6), bold=True)

    # SageMaker card
    c_sm = card(slide, 6.8, 0.7, 6.1, 5.5, C_GREEN_LIGHT, RGBColor(0xF0, 0xFD, 0xF4))
    tf = c_sm.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "🧪  SageMaker Unified Studio"
    tf.paragraphs[0].font.size = Pt(17)
    tf.paragraphs[0].font.color.rgb = RGBColor(0x15, 0x80, 0x3D)
    tf.paragraphs[0].font.bold = True
    ap(tf, "End-to-End ML & AI Development Platform", 10, RGBColor(0x64, 0x74, 0x8B), space_before=Pt(2))
    ap(tf, "", 6, C_BLACK)
    ap(tf, "One IDE for data engineering, analytics, ML model development, and generative AI — all connected to your S3 data lake and Bedrock foundation models.", 10, C_SLATE, space_before=Pt(6))
    sm_features = [
        "Unified workspace: SQL, Python, Spark, and ML in one IDE",
        "Native Bedrock integration for generative AI apps",
        "Built-in MLOps: training, tuning, deployment, monitoring",
        "Access S3, Redshift, Glue catalog from one interface",
        "Fine-tune foundation models on your own data",
        "Governed collaboration with Lake Formation permissions",
    ]
    for f in sm_features:
        ap(tf, "✓  " + f, 9, C_SLATE_LT, space_before=Pt(3))
    ap(tf, "", 4, C_BLACK)
    ap(tf, "Databricks gap: Databricks notebooks are powerful but siloed from AWS AI services. Models trained in Databricks can't natively deploy to Bedrock or integrate with QuickSight Q.", 9, C_RED, space_before=Pt(6), bold=True)

    bottom_callout(slide, "🔗", "Together:",
        "QuickSight delivers BI to every user at per-session cost. SageMaker Unified Studio gives data teams one IDE for analytics, ML, and generative AI — all on the same data lake.",
        y=6.4)


def slide_07_multivendor(prs):
    """Slide 7: Multi-Vendor Open Architecture (dark bg)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, RGBColor(0x0F, 0x17, 0x2A))
    top_bar(slide)

    # Header
    tb(slide, 0.6, 0.15, 11, 0.4, "Open Data, Open Choice: Best Service for Every Use Case", 22, C_ORANGE, bold=True)
    tb(slide, 0.6, 0.55, 11, 0.3, "With Iceberg tables on S3, your data is accessible to every platform. No single vendor owns it. Pick the best tool for each job.", 11, RGBColor(0x94, 0xA3, 0xB8))

    # Iceberg foundation bar
    c_found = card(slide, 0.6, 0.95, 12.1, 0.65, C_BLUE_LIGHT, RGBColor(0x1E, 0x3A, 0x5F))
    tf = c_found.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "🧊 Apache Iceberg Tables on Amazon S3"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = C_PURE_WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ap(tf, "One copy of data. Open format. ACID transactions. Readable by every platform below — simultaneously, without export or conversion.", 9, RGBColor(0xBF, 0xDB, 0xFE), align=PP_ALIGN.CENTER)

    # Arrow row
    tb(slide, 0.6, 1.65, 12.1, 0.25, "↓    Same data, read by all    ↓", 13, C_BLUE_LIGHT, align=PP_ALIGN.CENTER)

    # Three vendor cards
    # Databricks
    c_db = card(slide, 0.6, 2.0, 3.8, 3.5, C_PURPLE, RGBColor(0x1E, 0x1B, 0x2E))
    tf = c_db.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "🔥"
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ap(tf, "Databricks", 15, C_PURPLE_LT, True, PP_ALIGN.CENTER)
    ap(tf, "Data Engineering & Spark Workloads", 9, C_PURPLE_PALE, align=PP_ALIGN.CENTER, space_before=Pt(2))
    db_uses = [
        "Large-scale Spark ETL and data pipelines",
        "Complex data transformations",
        "Data science notebooks (Python/Scala)",
        "Delta-to-Iceberg via UniForm",
        "Existing Spark workloads (no migration)",
    ]
    for u in db_uses:
        ap(tf, "→ " + u, 9, C_PURPLE_PALE, space_before=Pt(3))

    # AWS
    c_aws = card(slide, 4.75, 2.0, 3.8, 3.5, C_ORANGE, RGBColor(0x1A, 0x23, 0x32))
    tf = c_aws.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "☁️"
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ap(tf, "AWS Analytics Suite", 15, C_ORANGE, True, PP_ALIGN.CENTER)
    ap(tf, "AI, Search, BI & Real-Time Analytics", 9, C_YELLOW, align=PP_ALIGN.CENTER, space_before=Pt(2))
    aws_uses = [
        "Redshift: SQL analytics + Bedrock AI queries",
        "OpenSearch: Full-text search + RAG",
        "QuickSight: Embedded BI dashboards",
        "SageMaker: ML model training + deployment",
        "Bedrock: Generative AI across all data",
        "Athena: Ad-hoc serverless SQL on Iceberg",
    ]
    for u in aws_uses:
        ap(tf, "→ " + u, 9, C_YELLOW, space_before=Pt(3))

    # Palantir
    c_pal = card(slide, 8.9, 2.0, 3.8, 3.5, C_PALANTIR, RGBColor(0x1A, 0x23, 0x20))
    tf = c_pal.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "🛡️"
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ap(tf, "Palantir Foundry", 15, C_PALANTIR_LT, True, PP_ALIGN.CENTER)
    ap(tf, "Operational Intelligence & Decision Support", 9, C_GREEN_LIGHT, align=PP_ALIGN.CENTER, space_before=Pt(2))
    pal_uses = [
        "Ontology-based operational applications",
        "Cross-domain data fusion and linking",
        "Mission-critical decision workflows",
        "Reads Iceberg tables from S3 directly",
        "Deploys on AWS (GovCloud compatible)",
    ]
    for u in pal_uses:
        ap(tf, "→ " + u, 9, C_GREEN_LIGHT, space_before=Pt(3))

    # Key message
    c_msg = card(slide, 0.6, 5.7, 12.1, 0.7, C_ORANGE, RGBColor(0x1A, 0x23, 0x32))
    tf = c_msg.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "The data stays in S3. The format stays open. The choice stays yours."
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = C_ORANGE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ap(tf, "No vendor can hold your data hostage. Add or remove platforms without data migration. Iceberg on S3 is the universal foundation.", 9, C_GRAY3, align=PP_ALIGN.CENTER, space_before=Pt(3))

    footer(slide, "Amazon Web Services", "Open Data Architecture | Databricks + AWS + Palantir on Iceberg")



def _build_matrix_slide(slide, title, badge, rows, callout_highlight, callout_text):
    """Helper to build a use-case matrix table slide (WHITE bg)."""
    add_bg(slide, C_PURE_WHITE)
    dark_header(slide, title, badge)

    # Table headers
    col_x = [0.3, 2.8, 5.5, 8.5, 11.0]
    col_w = [2.4, 2.6, 2.9, 2.4, 2.0]
    headers = ["Use Case", "Databricks", "AWS Analytics", "Palantir", "Best Fit"]
    h_colors = [C_SLATE_LT, RGBColor(0x7C, 0x3A, 0xED), RGBColor(0xB4, 0x53, 0x09), RGBColor(0x15, 0x80, 0x3D), C_BLUE_DARK]
    h_bgs = [None, RGBColor(0xF5, 0xF3, 0xFF), RGBColor(0xFF, 0xFB, 0xEB), RGBColor(0xF0, 0xFD, 0xF4), RGBColor(0xEF, 0xF6, 0xFF)]

    for i, (hdr, hc, hbg) in enumerate(zip(headers, h_colors, h_bgs)):
        if hbg:
            bar(slide, col_x[i], 0.6, col_w[i], 0.3, hbg)
        tb(slide, col_x[i], 0.6, col_w[i], 0.3, hdr, 9, hc, bold=True, align=PP_ALIGN.LEFT)

    # Separator line
    bar(slide, 0.3, 0.92, 12.7, 0.03, RGBColor(0xE2, 0xE8, 0xF0))

    # Data rows
    ry = 1.0
    row_h = 0.62
    for ri, row in enumerate(rows):
        uc_name, uc_desc, db_text, aws_text, pal_text, best_text, best_type = row
        if ri % 2 == 1:
            bar(slide, 0.3, ry, 12.7, row_h, RGBColor(0xFA, 0xFA, 0xFA))

        # Use case
        bx = tb(slide, col_x[0], ry, col_w[0], row_h, uc_name, 10, RGBColor(0x1E, 0x29, 0x3B), bold=True)
        ap(bx.text_frame, uc_desc, 8, RGBColor(0x64, 0x74, 0x8B))

        # Databricks
        db_color = C_GREEN if "✓" in db_text else (RGBColor(0xCA, 0x8A, 0x04) if "good" not in db_text.lower() else RGBColor(0xCA, 0x8A, 0x04))
        if "✓" in db_text:
            db_color = C_GREEN
        elif "gap" in db_text.lower() or "No " in db_text or "Not " in db_text or "Limited" in db_text or "Basic" in db_text or "Requires" in db_text:
            db_color = C_RED
        else:
            db_color = RGBColor(0xCA, 0x8A, 0x04)
        tb(slide, col_x[1], ry, col_w[1], row_h, db_text, 9, db_color)

        # AWS
        aws_color = C_GREEN if "✓" in aws_text else RGBColor(0xCA, 0x8A, 0x04)
        tb(slide, col_x[2], ry, col_w[2], row_h, aws_text, 9, aws_color)

        # Palantir
        pal_color = C_GREEN if "✓" in pal_text else (C_RED if "Not " in pal_text or "No " in pal_text or "Requires" in pal_text else RGBColor(0xCA, 0x8A, 0x04))
        tb(slide, col_x[3], ry, col_w[3], row_h, pal_text, 9, pal_color)

        # Best fit
        best_colors = {
            "db": (RGBColor(0xED, 0xE9, 0xFE), RGBColor(0x6D, 0x28, 0xD9)),
            "aws": (RGBColor(0xFF, 0xF7, 0xED), RGBColor(0xC2, 0x41, 0x0C)),
            "pal": (RGBColor(0xDC, 0xFC, 0xE7), RGBColor(0x15, 0x80, 0x3D)),
            "both": (RGBColor(0xEF, 0xF6, 0xFF), C_BLUE_DARK),
        }
        bg_c, fg_c = best_colors.get(best_type, best_colors["aws"])
        c_best = card(slide, col_x[4], ry + 0.05, col_w[4] - 0.1, 0.3, bg_c, bg_c)
        c_best.text_frame.paragraphs[0].text = best_text
        c_best.text_frame.paragraphs[0].font.size = Pt(9)
        c_best.text_frame.paragraphs[0].font.color.rgb = fg_c
        c_best.text_frame.paragraphs[0].font.bold = True
        c_best.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        ry += row_h

    bottom_callout(slide, "🧊", callout_highlight, callout_text, y=6.4)


def slide_08_matrix1(prs):
    """Slide 8: Use Case Alignment Matrix (WHITE bg)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rows = [
        ("Large-Scale ETL / Spark", "Petabyte data transformations", "✓ Native Spark, Delta Live Tables", "EMR/Glue Spark, but less notebook UX", "Not designed for heavy ETL", "Databricks", "db"),
        ("Generative AI / RAG", "Natural language queries, chatbots", "No native Bedrock; must use DB models", "✓ Bedrock + OpenSearch + Redshift native", "AIP for operational AI, not general RAG", "AWS (Bedrock)", "aws"),
        ("SQL Analytics / DW", "Ad-hoc queries, reporting", "Databricks SQL, requires cluster compute", "✓ Redshift Serverless, pay-per-query", "Not a SQL analytics platform", "AWS (Redshift)", "aws"),
        ("Business Intelligence", "Dashboards, embedded analytics", "Basic dashboards, not embeddable", "✓ QuickSight: embeddable, per-session, Q", "Operational dashboards in Foundry", "AWS (QuickSight)", "aws"),
        ("Full-Text & Semantic Search", "Document search, vector search", "No search engine; requires third-party", "✓ OpenSearch: full-text + vector + serverless", "Object search within Foundry ontology", "AWS (OpenSearch)", "aws"),
        ("ML Model Training & MLOps", "Custom model development", "✓ MLflow, Feature Store, model serving", "✓ SageMaker Unified Studio, Bedrock fine-tuning", "Palantir AIP for operational ML", "Databricks / SageMaker", "both"),
        ("Operational Decision Support", "Mission workflows, ontology apps", "Not designed for operational workflows", "Step Functions + custom apps", "✓ Foundry ontology, AIP, decision workflows", "Palantir", "pal"),
        ("Data Science Notebooks", "Interactive Python/R/Scala", "✓ Collaborative notebooks, great UX", "✓ SageMaker Studio notebooks", "Code Workbook in Foundry", "Databricks / SageMaker", "both"),
    ]
    _build_matrix_slide(slide,
        "Use Case Alignment: Best Platform for Each Workload",
        "MAY THE BEST SERVICE WIN",
        rows,
        "The Iceberg advantage:",
        "All three platforms read the same Iceberg tables on S3. No data duplication. No export pipelines. Align the best service to each use case — and change your mind later without migration.")



def slide_09_bedrock(prs):
    """Slide 9: Bedrock AI Layer (dark bg)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, C_DARK_BG)
    top_bar(slide)

    # Header
    tb(slide, 0.6, 0.15, 11, 0.4, "Amazon Bedrock: The AI Layer That Ties It All Together", 22, C_ORANGE, bold=True)
    tb(slide, 0.6, 0.55, 11, 0.3, "Every AWS analytics service feeds into Bedrock natively — one AI layer for structured, unstructured, and real-time data", 11, C_GRAY)

    # Bedrock center bar
    c_bed = card(slide, 0.6, 0.95, 12.1, 0.65, C_ORANGE, C_ORANGE)
    tf = c_bed.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "Amazon Bedrock Knowledge Bases"
    tf.paragraphs[0].font.size = Pt(17)
    tf.paragraphs[0].font.color.rgb = C_DARK_BG
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ap(tf, "Natural language interface to ALL your data — structured and unstructured — through one unified API", 10, RGBColor(0x45, 0x1A, 0x03), align=PP_ALIGN.CENTER)

    # 4 source cards
    sources = [
        ("📊", "Redshift", "Structured Data", "Bedrock generates SQL, queries live tables. Real-time answers from production data. Zero ETL.", "NATIVE SOURCE"),
        ("🔍", "OpenSearch", "Vector Store / RAG", "Semantic search across millions of documents. Vector embeddings for retrieval-augmented generation.", "NATIVE SOURCE"),
        ("🪣", "S3 Data Lake", "Documents & Files", "PDFs, images, CSVs ingested directly. Auto-chunked, embedded, and indexed for retrieval.", "NATIVE SOURCE"),
        ("🧪", "SageMaker", "Custom Models", "Fine-tuned models deployed as Bedrock endpoints. Custom ML inference alongside foundation models.", "DIRECT INTEGRATION"),
    ]
    sx = 0.6
    for icon, name, stype, how, tag in sources:
        c_src = card(slide, sx, 1.8, 2.9, 2.4, C_ORANGE, RGBColor(0x2A, 0x36, 0x48))
        tf = c_src.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = icon
        tf.paragraphs[0].font.size = Pt(22)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        ap(tf, name, 11, C_ORANGE, True, PP_ALIGN.CENTER, space_before=Pt(4))
        ap(tf, stype, 8, C_GRAY, align=PP_ALIGN.CENTER)
        ap(tf, how, 8, C_GRAY3, space_before=Pt(6))
        tag_color = RGBColor(0x06, 0x5F, 0x46) if "NATIVE" in tag else C_BLUE_DARK
        tag_bg = C_GREEN_PALE if "NATIVE" in tag else RGBColor(0xDB, 0xEA, 0xFE)
        ap(tf, tag, 8, tag_color, True, PP_ALIGN.CENTER, space_before=Pt(6))
        sx += 3.1

    # Bottom comparison: What Databricks Can't Do vs What AWS Delivers
    c_cant = card(slide, 0.6, 4.4, 5.9, 1.5, C_ORANGE, RGBColor(0x2A, 0x36, 0x48))
    tf = c_cant.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "🚫 What Databricks Can't Do"
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.color.rgb = C_ORANGE
    tf.paragraphs[0].font.bold = True
    ap(tf, "No native Bedrock connection. Must export data to S3, then Bedrock reads stale copies. AI models purchased through Databricks marketplace — not Bedrock foundation models. No unified AI layer across structured + unstructured data.", 9, C_GRAY3, space_before=Pt(6))

    c_can = card(slide, 6.8, 4.4, 5.9, 1.5, C_ORANGE, RGBColor(0x2A, 0x36, 0x48))
    tf = c_can.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "✅ What AWS Delivers"
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.color.rgb = C_ORANGE
    tf.paragraphs[0].font.bold = True
    ap(tf, "One Bedrock Knowledge Base queries Redshift (structured), OpenSearch (vectors), and S3 (documents) simultaneously. Users ask one question, get answers from all data sources. One API, one bill, one security model.", 9, C_GRAY3, space_before=Pt(6))

    footer(slide, "Amazon Web Services", "Amazon Bedrock | Unified AI Across All Data Sources")


def slide_10_databricks_only(prs):
    """Slide 10: Databricks-Only Current State (purple-tinted dark bg)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, RGBColor(0x1A, 0x16, 0x25))
    top_bar(slide, C_PURPLE)

    # Header
    tb(slide, 0.6, 0.15, 10, 0.35, "Current State: Databricks-Only Architecture", 20, C_PURPLE_LT, bold=True)
    tb(slide, 0.6, 0.5, 10, 0.25, "All data flows through one vendor — creating cost, scale, governance, and AI limitations", 10, RGBColor(0x94, 0xA3, 0xB8))

    # Flow labels
    flow_labels = ["Data Sources", "Data Movement", "Data Lake /\nGovernance", "Data Warehouse", "AI", "BI"]
    lx = [0.3, 2.2, 4.2, 6.3, 8.4, 10.5]
    lw = [1.5, 1.8, 1.8, 1.8, 1.8, 1.2]
    for i, lbl in enumerate(flow_labels):
        tb(slide, lx[i], 0.8, lw[i], 0.3, lbl, 8, RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.CENTER)

    # Flow boxes
    # Sources
    src_items = [("Postgres", "AWS RDS"), ("Oracle", "On-Prem"), ("Other", "Sources")]
    sy = 1.15
    for name, sub in src_items:
        c_s = card(slide, 0.3, sy, 1.5, 0.55, RGBColor(0x64, 0x74, 0x8B), RGBColor(0x33, 0x41, 0x55))
        tf = c_s.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = name
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.color.rgb = C_WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        ap(tf, sub, 8, RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.CENTER)
        sy += 0.6

    # Arrows and boxes
    tb(slide, 1.85, 1.5, 0.3, 0.3, "→", 18, C_PURPLE)

    # Glue/Informatica
    c_glue = card(slide, 2.2, 1.2, 1.8, 0.7, C_PURPLE, RGBColor(0x1E, 0x1B, 0x2E))
    tf = c_glue.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = "Glue / Informatica"
    tf.paragraphs[0].font.size = Pt(11); tf.paragraphs[0].font.color.rgb = C_PURPLE_LT; tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ap(tf, "ETL to S3", 8, C_PURPLE_PALE, align=PP_ALIGN.CENTER)

    tb(slide, 4.05, 1.5, 0.3, 0.3, "→", 18, C_PURPLE)

    # Data Lake S3
    c_lake = card(slide, 4.2, 1.2, 1.8, 0.7, C_BLUE, RGBColor(0x1E, 0x3A, 0x5F))
    tf = c_lake.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = "Data Lake S3"
    tf.paragraphs[0].font.size = Pt(11); tf.paragraphs[0].font.color.rgb = C_BLUE_LIGHT; tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ap(tf, "Databricks Catalog", 8, C_BLUE_PALE, align=PP_ALIGN.CENTER)

    tb(slide, 6.05, 1.5, 0.3, 0.3, "→", 18, C_PURPLE)

    # Databricks
    c_dbr = card(slide, 6.3, 1.2, 1.8, 0.7, C_PURPLE, RGBColor(0x1E, 0x1B, 0x2E))
    tf = c_dbr.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = "Databricks"
    tf.paragraphs[0].font.size = Pt(11); tf.paragraphs[0].font.color.rgb = C_PURPLE_LT; tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ap(tf, "Delta / Proton Query", 8, C_PURPLE_PALE, align=PP_ALIGN.CENTER)

    tb(slide, 8.15, 1.5, 0.3, 0.3, "→", 18, C_PURPLE)

    # Databricks AI
    c_ai = card(slide, 8.4, 1.2, 1.8, 0.7, C_PURPLE, RGBColor(0x1E, 0x1B, 0x2E))
    tf = c_ai.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = "Databricks AI?"
    tf.paragraphs[0].font.size = Pt(11); tf.paragraphs[0].font.color.rgb = C_PURPLE_LT; tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ap(tf, "Limited models", 8, C_PURPLE_PALE, align=PP_ALIGN.CENTER)

    tb(slide, 10.25, 1.5, 0.3, 0.3, "→", 18, C_PURPLE)

    # BI Layer
    c_bi = card(slide, 10.5, 1.2, 1.5, 0.7, C_PURPLE, RGBColor(0x1E, 0x1B, 0x2E))
    tf = c_bi.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = "BI Layer"
    tf.paragraphs[0].font.size = Pt(10); tf.paragraphs[0].font.color.rgb = C_PURPLE_LT; tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ap(tf, "QuickSight? Tableau? PBI?", 7, C_PURPLE_PALE, align=PP_ALIGN.CENTER)

    # 4 problem callouts
    problems = [
        ("⚠ Near Real-Time Gap", "Dumping to S3 is not real time. Near real-time use cases are not supported. ETL latency = stale data for analytics."),
        ("⚠ Governance Lock-In", "Data governance done entirely through Databricks catalog. Vendor lock-in for access control and security."),
        ("⚠ AI Model Limitations", "Not all models available in Databricks. Limits use-case-driven AI solutions. Bedrock has more models at better prices."),
        ("⚠ Cost & Scale Challenges", "Loading all data into Databricks is costly and has had scaling challenges. BI tool integrations uncertain."),
    ]
    px = 0.3
    for title, body in problems:
        c_p = card(slide, px, 2.2, 2.95, 1.3, C_RED_LIGHTER, RGBColor(0x2A, 0x1A, 0x1A))
        tf = c_p.text_frame; tf.word_wrap = True
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = Pt(9); tf.paragraphs[0].font.color.rgb = C_RED_LIGHTER; tf.paragraphs[0].font.bold = True
        ap(tf, body, 8, C_RED_PALE, space_before=Pt(4))
        px += 3.1

    # Core question
    c_q = card(slide, 0.3, 3.7, 12.4, 0.55, C_RED_LIGHTER, RGBColor(0x2A, 0x1A, 0x1A))
    tf = c_q.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = "Core Question: "
    run1.font.size = Pt(11); run1.font.color.rgb = C_RED_LIGHTER; run1.font.bold = True
    run2 = p.add_run()
    run2.text = "Should one vendor own your data movement, governance, warehousing, AI, and BI? What happens when a better tool emerges for a specific use case?"
    run2.font.size = Pt(10); run2.font.color.rgb = C_RED_PALE; run2.font.bold = False

    footer(slide, "Current Architecture", "Databricks-Only | Single Vendor Dependency")



def slide_11_open_architecture(prs):
    """Slide 11: Open Architecture with Iceberg (dark bg)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, RGBColor(0x0F, 0x17, 0x2A))
    top_bar(slide)

    # Header
    tb(slide, 0.6, 0.1, 12, 0.35, "Open Architecture: Iceberg Tables + Best Service for Each Use Case", 19, C_ORANGE, bold=True)
    tb(slide, 0.6, 0.42, 12, 0.25, "zETL from Postgres/Oracle → Redshift Tables → Iceberg on S3. Every engine reads the same data. May the best service win.", 9, RGBColor(0x94, 0xA3, 0xB8))

    # Flow labels
    flow_labels = ["Data Sources", "zETL", "Redshift +\nIceberg on S3", "Warehouse /\nAnalytics", "AI Layer", "BI"]
    lx = [0.2, 2.0, 3.2, 5.3, 8.3, 10.8]
    lw = [1.3, 0.8, 1.8, 2.8, 2.2, 1.5]
    for i, lbl in enumerate(flow_labels):
        tb(slide, lx[i], 0.7, lw[i], 0.3, lbl, 7, RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.CENTER)

    # Sources
    src_items = [("Postgres", "AWS RDS"), ("Oracle", "On-Prem"), ("Other", "Sources")]
    sy = 1.05
    for name, sub in src_items:
        c_s = card(slide, 0.2, sy, 1.3, 0.45, RGBColor(0x64, 0x74, 0x8B), RGBColor(0x1A, 0x22, 0x30))
        tf = c_s.text_frame; tf.word_wrap = True
        tf.paragraphs[0].text = name
        tf.paragraphs[0].font.size = Pt(9); tf.paragraphs[0].font.color.rgb = C_WHITE; tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        ap(tf, sub, 7, RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.CENTER)
        sy += 0.5

    tb(slide, 1.55, 1.35, 0.3, 0.3, "→", 16, C_ORANGE)

    # zETL box
    c_zetl = card(slide, 1.9, 1.1, 1.0, 0.65, C_ORANGE, RGBColor(0x2A, 0x30, 0x3E))
    tf = c_zetl.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = "zETL"
    tf.paragraphs[0].font.size = Pt(11); tf.paragraphs[0].font.color.rgb = C_ORANGE; tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ap(tf, "Near real-time", 7, C_YELLOW, align=PP_ALIGN.CENTER)

    tb(slide, 2.95, 1.35, 0.3, 0.3, "→", 16, C_ORANGE)

    # Redshift + Iceberg
    c_rs = card(slide, 3.2, 1.0, 1.8, 0.4, C_BLUE, RGBColor(0x1A, 0x30, 0x50))
    tf = c_rs.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = "Redshift Tables"
    tf.paragraphs[0].font.size = Pt(10); tf.paragraphs[0].font.color.rgb = C_BLUE_LIGHT; tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    tb(slide, 3.8, 1.42, 0.6, 0.2, "↓", 12, C_BLUE_LIGHT, align=PP_ALIGN.CENTER)

    c_ice = card(slide, 3.2, 1.6, 1.8, 0.5, C_BLUE_LIGHT, RGBColor(0x1A, 0x30, 0x50))
    tf = c_ice.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = "🧊 Iceberg Tables on S3"
    tf.paragraphs[0].font.size = Pt(9); tf.paragraphs[0].font.color.rgb = C_BLUE_LIGHT; tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ap(tf, "Open format · Any engine reads", 7, C_BLUE_PALE, align=PP_ALIGN.CENTER)

    tb(slide, 5.05, 1.35, 0.3, 0.3, "→", 16, C_GREEN)

    # Warehouse/Analytics grid (2x2 + SageMaker)
    grid_items = [
        (5.3, 1.0, "Databricks", "Spark / ML", C_PURPLE, RGBColor(0x1E, 0x1B, 0x2E), C_PURPLE_LT, C_PURPLE_PALE),
        (6.65, 1.0, "Palantir", "Foundry / AIP", C_PALANTIR, RGBColor(0x0F, 0x23, 0x18), C_PALANTIR_LT, C_GREEN_LIGHT),
        (5.3, 1.5, "Redshift", "SQL Analytics", C_ORANGE, RGBColor(0x2A, 0x30, 0x3E), C_ORANGE, C_YELLOW),
        (6.65, 1.5, "Snowflake", "If needed", RGBColor(0x64, 0x74, 0x8B), RGBColor(0x1A, 0x22, 0x30), C_WHITE, RGBColor(0x94, 0xA3, 0xB8)),
    ]
    for gx, gy, gname, gsub, gborder, gfill, gname_c, gsub_c in grid_items:
        c_g = card(slide, gx, gy, 1.25, 0.42, gborder, gfill)
        tf = c_g.text_frame; tf.word_wrap = True
        tf.paragraphs[0].text = gname
        tf.paragraphs[0].font.size = Pt(9); tf.paragraphs[0].font.color.rgb = gname_c; tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        ap(tf, gsub, 7, gsub_c, align=PP_ALIGN.CENTER)

    # SageMaker
    c_sm = card(slide, 5.3, 2.0, 2.6, 0.3, C_ORANGE, RGBColor(0x2A, 0x30, 0x3E))
    tf = c_sm.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = "SageMaker Unified Studio"
    tf.paragraphs[0].font.size = Pt(8); tf.paragraphs[0].font.color.rgb = C_ORANGE; tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    tb(slide, 8.0, 1.35, 0.3, 0.3, "→", 16, C_GREEN)

    # AI Layer
    ai_items = [
        ("Bedrock", "GenAI / RAG", C_ORANGE, C_YELLOW),
        ("SageMaker", "Custom ML", C_ORANGE, C_YELLOW),
        ("Palantir AIP", "Ontology AI", C_PALANTIR_LT, C_GREEN_LIGHT),
        ("Databricks AI", "If best fit", C_PURPLE_LT, C_PURPLE_PALE),
    ]
    ay = 1.0
    for aname, asub, ac, asc in ai_items:
        c_a = card(slide, 8.3, ay, 1.8, 0.38, C_ORANGE if "Bedrock" in aname or "Sage" in aname else (C_PALANTIR if "Palantir" in aname else C_PURPLE), RGBColor(0x1A, 0x22, 0x30))
        tf = c_a.text_frame; tf.word_wrap = True
        tf.paragraphs[0].text = aname
        tf.paragraphs[0].font.size = Pt(9); tf.paragraphs[0].font.color.rgb = ac; tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        ap(tf, asub, 7, asc, align=PP_ALIGN.CENTER)
        ay += 0.42

    tb(slide, 10.15, 1.35, 0.3, 0.3, "→", 16, C_GREEN)

    # BI
    bi_items = ["QuickSight", "Tableau", "Power BI", "MicroStrategy"]
    by = 1.0
    for bname in bi_items:
        bc = C_ORANGE if bname == "QuickSight" else RGBColor(0x64, 0x74, 0x8B)
        bf = RGBColor(0x2A, 0x30, 0x3E) if bname == "QuickSight" else RGBColor(0x1A, 0x22, 0x30)
        nc = C_ORANGE if bname == "QuickSight" else C_WHITE
        c_bi = card(slide, 10.5, by, 1.5, 0.35, bc, bf)
        tf = c_bi.text_frame; tf.word_wrap = True
        tf.paragraphs[0].text = bname
        tf.paragraphs[0].font.size = Pt(9); tf.paragraphs[0].font.color.rgb = nc; tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        by += 0.4

    # 3 benefit cards
    benefits = [
        ("✓ Near Real-Time with zETL", "Postgres/Oracle connectors create Iceberg tables via zero-ETL through Redshift. Near real-time data availability — no batch dump to S3."),
        ("✓ Best AI Tool for Each Use Case", "Databricks ML for some use cases. Bedrock for GenAI. Palantir AIP for ontology-driven AI. SageMaker for custom models."),
        ("✓ Open Governance + BI Freedom", "Lake Formation for governance instead of Databricks catalog lock-in. Any BI tool connects to Iceberg."),
    ]
    bx = 0.2
    for btitle, bbody in benefits:
        c_b = card(slide, bx, 2.6, 4.0, 0.9, C_GREEN_LIGHT, RGBColor(0x14, 0x20, 0x25))
        tf = c_b.text_frame; tf.word_wrap = True
        tf.paragraphs[0].text = btitle
        tf.paragraphs[0].font.size = Pt(9); tf.paragraphs[0].font.color.rgb = C_GREEN_LIGHT; tf.paragraphs[0].font.bold = True
        ap(tf, bbody, 8, C_GREEN_PALE, space_before=Pt(3))
        bx += 4.2

    # Key principle
    c_kp = card(slide, 0.2, 3.7, 12.9, 0.5, C_ORANGE, RGBColor(0x1A, 0x22, 0x30))
    tf = c_kp.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = "Key Principle: "
    run1.font.size = Pt(10); run1.font.color.rgb = C_ORANGE; run1.font.bold = True
    run2 = p.add_run()
    run2.text = "Iceberg tables on S3 make your data open to every engine. Databricks, Palantir, Redshift, Snowflake, SageMaker — they all read the same tables."
    run2.font.size = Pt(9); run2.font.color.rgb = C_YELLOW; run2.font.bold = False

    footer(slide, "Amazon Web Services", "Open Architecture | Iceberg + zETL + Best Service Per Use Case")



def slide_12_aws_solves(prs):
    """Slide 12: AWS Solves Every Problem (dark bg)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, RGBColor(0x0A, 0x16, 0x28))
    top_bar(slide)

    # Header
    tb(slide, 0.6, 0.1, 11, 0.35, "AWS Solves Every Problem in the Pipeline", 19, C_ORANGE, bold=True)
    tb(slide, 0.6, 0.42, 11, 0.25, "Same architecture flow — but with an AWS benefit replacing each Databricks-only problem", 9, RGBColor(0x94, 0xA3, 0xB8))

    # Column headers
    col_labels = ["Data Sources", "Data Movement", "Data Lake /\nGovernance", "Data Warehouse", "AI", "BI"]
    cx = [0.2, 2.0, 4.0, 6.1, 8.2, 10.5]
    cw = [1.4, 1.8, 1.8, 1.8, 1.8, 1.5]
    for i, lbl in enumerate(col_labels):
        tb(slide, cx[i], 0.7, cw[i], 0.3, lbl, 7, RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.CENTER)

    # Sources
    src_items = [("Postgres", "AWS RDS"), ("Oracle", "On-Prem"), ("Other", "Sources")]
    sy = 1.05
    for name, sub in src_items:
        c_s = card(slide, 0.2, sy, 1.4, 0.42, RGBColor(0x64, 0x74, 0x8B), RGBColor(0x1A, 0x22, 0x30))
        tf = c_s.text_frame; tf.word_wrap = True
        tf.paragraphs[0].text = name; tf.paragraphs[0].font.size = Pt(9); tf.paragraphs[0].font.color.rgb = C_WHITE; tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        ap(tf, sub, 7, RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.CENTER)
        sy += 0.47

    # Arrows
    for ax in [1.65, 3.85, 5.95, 8.05, 10.3]:
        tb(slide, ax, 1.3, 0.3, 0.3, "→", 16, C_ORANGE)

    # Flow boxes
    # zETL
    c_z = card(slide, 2.0, 1.05, 1.8, 0.65, C_ORANGE, RGBColor(0x2A, 0x30, 0x3E))
    tf = c_z.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = "Redshift zETL"; tf.paragraphs[0].font.size = Pt(12); tf.paragraphs[0].font.color.rgb = C_ORANGE; tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ap(tf, "Zero-ETL connectors", 8, C_YELLOW, align=PP_ALIGN.CENTER)

    # Iceberg + Glue
    c_i = card(slide, 4.0, 1.0, 1.8, 0.5, C_BLUE_LIGHT, RGBColor(0x1A, 0x30, 0x50))
    tf = c_i.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = "🧊 Iceberg on S3"; tf.paragraphs[0].font.size = Pt(10); tf.paragraphs[0].font.color.rgb = C_BLUE_LIGHT; tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ap(tf, "Open format · Lake Formation", 7, C_BLUE_PALE, align=PP_ALIGN.CENTER)
    c_gc = card(slide, 4.0, 1.55, 1.8, 0.3, C_ORANGE, RGBColor(0x2A, 0x30, 0x3E))
    tf = c_gc.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = "Glue Data Catalog"; tf.paragraphs[0].font.size = Pt(8); tf.paragraphs[0].font.color.rgb = C_ORANGE; tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Redshift + OpenSearch
    c_r = card(slide, 6.1, 1.0, 1.8, 0.5, C_ORANGE, RGBColor(0x2A, 0x30, 0x3E))
    tf = c_r.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = "Redshift Serverless"; tf.paragraphs[0].font.size = Pt(10); tf.paragraphs[0].font.color.rgb = C_ORANGE; tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ap(tf, "+ OpenSearch", 8, C_YELLOW, align=PP_ALIGN.CENTER)
    c_ae = card(slide, 6.1, 1.55, 1.8, 0.3, C_ORANGE, RGBColor(0x2A, 0x30, 0x3E))
    tf = c_ae.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = "Athena · EMR"; tf.paragraphs[0].font.size = Pt(8); tf.paragraphs[0].font.color.rgb = C_ORANGE; tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Bedrock + SageMaker
    c_b = card(slide, 8.2, 1.0, 1.8, 0.5, C_ORANGE, RGBColor(0x2A, 0x30, 0x3E))
    tf = c_b.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = "Amazon Bedrock"; tf.paragraphs[0].font.size = Pt(10); tf.paragraphs[0].font.color.rgb = C_ORANGE; tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ap(tf, "+ SageMaker", 8, C_YELLOW, align=PP_ALIGN.CENTER)
    c_m = card(slide, 8.2, 1.55, 1.8, 0.3, C_ORANGE, RGBColor(0x2A, 0x30, 0x3E))
    tf = c_m.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = "Nova · Titan · 100+ models"; tf.paragraphs[0].font.size = Pt(7); tf.paragraphs[0].font.color.rgb = C_ORANGE; tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # QuickSight
    c_q = card(slide, 10.5, 1.05, 1.5, 0.55, C_ORANGE, RGBColor(0x2A, 0x30, 0x3E))
    tf = c_q.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = "QuickSight"; tf.paragraphs[0].font.size = Pt(10); tf.paragraphs[0].font.color.rgb = C_ORANGE; tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ap(tf, "+ Tableau / PBI", 7, C_YELLOW, align=PP_ALIGN.CENTER)

    # Green benefit boxes
    green_benefits = [
        (2.0, "✓ Eliminate Glue/Informatica", "No ETL tool to buy or maintain. Redshift zETL auto-replicates from Postgres/Aurora in near real-time."),
        (4.0, "✓ Open Governance, Not Vendor Catalog", "Lake Formation + Glue Catalog for governance — AWS-native, IAM-integrated. No Databricks catalog lock-in."),
        (6.1, "✓ Serverless, Scales, No Load Cost", "Redshift Serverless: pay per query, not per cluster. Queries Iceberg on S3 directly. OpenSearch for search/RAG."),
        (8.2, "✓ 100+ Models, Best Price, Native", "Bedrock: 100+ foundation models. Use-case-driven model selection. Native KB integration with Redshift + OpenSearch."),
    ]
    for gx, gtitle, gbody in green_benefits:
        c_g = card(slide, gx, 2.05, 1.8, 1.2, C_GREEN_LIGHT, RGBColor(0x14, 0x20, 0x25))
        tf = c_g.text_frame; tf.word_wrap = True
        tf.paragraphs[0].text = gtitle; tf.paragraphs[0].font.size = Pt(8); tf.paragraphs[0].font.color.rgb = C_GREEN_LIGHT; tf.paragraphs[0].font.bold = True
        ap(tf, gbody, 7, C_GREEN_PALE, space_before=Pt(3))

    # BI benefit
    c_gb = card(slide, 10.5, 2.05, 1.5, 1.2, C_GREEN_LIGHT, RGBColor(0x14, 0x20, 0x25))
    tf = c_gb.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = "✓ Native + Any BI Tool"; tf.paragraphs[0].font.size = Pt(8); tf.paragraphs[0].font.color.rgb = C_GREEN_LIGHT; tf.paragraphs[0].font.bold = True
    ap(tf, "QuickSight: embedded, per-session, Q natural language. Plus Tableau, PBI, MicroStrategy all connect natively.", 7, C_GREEN_PALE, space_before=Pt(3))

    # Red "Was:" boxes
    was_items = [
        (2.0, "Was: Glue/Informatica", "Expensive licenses. Dedicated team. Batch only — not real-time."),
        (4.0, "Was: Databricks Catalog", "Vendor lock-in for governance. All access control through Databricks."),
        (6.1, "Was: Load into Databricks", "Costly to load. Scaling challenges. Premium compute for every query."),
        (8.2, "Was: Databricks AI Only", "Limited models. Not all available. Higher prices than Bedrock."),
    ]
    for wx, wtitle, wbody in was_items:
        c_w = card(slide, wx, 3.4, 1.8, 0.85, C_RED_LIGHTER, RGBColor(0x2A, 0x1A, 0x1A))
        tf = c_w.text_frame; tf.word_wrap = True
        tf.paragraphs[0].text = wtitle; tf.paragraphs[0].font.size = Pt(7); tf.paragraphs[0].font.color.rgb = C_RED_LIGHTER; tf.paragraphs[0].font.bold = True
        ap(tf, wbody, 6, C_RED_PALE, space_before=Pt(2))

    c_wb = card(slide, 10.5, 3.4, 1.5, 0.85, C_RED_LIGHTER, RGBColor(0x2A, 0x1A, 0x1A))
    tf = c_wb.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = "Was: DB Dashboards"; tf.paragraphs[0].font.size = Pt(7); tf.paragraphs[0].font.color.rgb = C_RED_LIGHTER; tf.paragraphs[0].font.bold = True
    ap(tf, "Basic. Not embeddable. Connector risk with 3rd party BI.", 6, C_RED_PALE, space_before=Pt(2))

    # Bottom summary
    c_sum = card(slide, 0.2, 4.45, 12.9, 0.5, C_ORANGE, RGBColor(0x1A, 0x22, 0x30))
    tf = c_sum.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = "Every problem has an AWS-native answer. "
    run1.font.size = Pt(10); run1.font.color.rgb = C_ORANGE; run1.font.bold = True
    run2 = p.add_run()
    run2.text = "Eliminate ETL tools, remove vendor lock-in, get serverless scale, access 100+ AI models, and connect any BI tool — all on one bill, one IAM, one VPC."
    run2.font.size = Pt(9); run2.font.color.rgb = C_YELLOW; run2.font.bold = False

    footer(slide, "Amazon Web Services", "AWS Benefits | Problem → Solution at Every Step")



def slide_13_swim_lanes(prs):
    """Slide 13: Swim Lane Flow (dark bg)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, RGBColor(0x0F, 0x17, 0x2A))
    top_bar(slide)

    # Header
    tb(slide, 0.6, 0.1, 12, 0.35, "Open Architecture: Iceberg on S3 → Best Service for Each Use Case", 19, C_ORANGE, bold=True)
    tb(slide, 0.6, 0.42, 12, 0.25, "zETL from Postgres/Oracle → Redshift → Iceberg Tables on S3. Three paths fan out. Each vendor stays in its lane.", 9, RGBColor(0x94, 0xA3, 0xB8))

    # LEFT SIDE: Sources → zETL → Iceberg
    # Sources
    src_items = ["Postgres", "Oracle", "Other"]
    sy = 1.0
    for name in src_items:
        c_s = card(slide, 0.2, sy, 1.0, 0.3, RGBColor(0x64, 0x74, 0x8B), RGBColor(0x1A, 0x22, 0x30))
        tf = c_s.text_frame; tf.word_wrap = True
        tf.paragraphs[0].text = name; tf.paragraphs[0].font.size = Pt(8); tf.paragraphs[0].font.color.rgb = C_WHITE; tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        sy += 0.35

    tb(slide, 1.25, 1.2, 0.3, 0.3, "→", 14, C_ORANGE)

    # zETL
    c_z = card(slide, 1.55, 1.05, 0.8, 0.5, C_ORANGE, RGBColor(0x2A, 0x30, 0x3E))
    tf = c_z.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = "zETL"; tf.paragraphs[0].font.size = Pt(10); tf.paragraphs[0].font.color.rgb = C_ORANGE; tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ap(tf, "Real-time", 7, C_YELLOW, align=PP_ALIGN.CENTER)

    tb(slide, 2.4, 1.2, 0.3, 0.3, "→", 14, C_BLUE_LIGHT)

    # Redshift Tables
    c_rt = card(slide, 2.7, 0.95, 1.5, 0.35, C_BLUE, RGBColor(0x1A, 0x30, 0x50))
    tf = c_rt.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = "Redshift Tables"; tf.paragraphs[0].font.size = Pt(9); tf.paragraphs[0].font.color.rgb = C_BLUE_LIGHT; tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    tb(slide, 3.2, 1.32, 0.5, 0.2, "↓", 12, C_BLUE_LIGHT, align=PP_ALIGN.CENTER)

    # Iceberg
    c_ic = card(slide, 2.7, 1.5, 1.5, 0.45, C_BLUE_LIGHT, RGBColor(0x1A, 0x30, 0x50))
    tf = c_ic.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = "🧊 Iceberg on S3"; tf.paragraphs[0].font.size = Pt(9); tf.paragraphs[0].font.color.rgb = C_BLUE_LIGHT; tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ap(tf, "Open · Any engine reads", 7, C_BLUE_PALE, align=PP_ALIGN.CENTER)

    # RIGHT SIDE: 3 swim lanes
    # Lane column headers
    lane_headers = ["Analytics / Warehouse", "AI / ML", "BI"]
    lhx = [5.0, 7.8, 10.2]
    lhw = [2.5, 2.2, 1.8]
    for i, lh in enumerate(lane_headers):
        tb(slide, lhx[i], 0.75, lhw[i], 0.2, lh, 7, RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.CENTER)

    # DATABRICKS LANE (Purple)
    c_dl = card(slide, 4.4, 1.0, 8.5, 0.7, C_PURPLE, RGBColor(0x14, 0x10, 0x20))
    tb(slide, 4.5, 1.15, 0.3, 0.3, "→", 14, C_PURPLE)

    c_db1 = card(slide, 5.0, 1.1, 2.3, 0.5, C_PURPLE, RGBColor(0x1E, 0x1B, 0x2E))
    tf = c_db1.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = "Databricks"; tf.paragraphs[0].font.size = Pt(9); tf.paragraphs[0].font.color.rgb = C_PURPLE_LT; tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ap(tf, "Spark ETL · Delta Proton", 7, C_PURPLE_PALE, align=PP_ALIGN.CENTER)

    tb(slide, 7.4, 1.25, 0.3, 0.3, "→", 12, C_PURPLE)

    c_db2 = card(slide, 7.8, 1.1, 2.0, 0.5, C_PURPLE, RGBColor(0x1E, 0x1B, 0x2E))
    tf = c_db2.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = "Databricks AI"; tf.paragraphs[0].font.size = Pt(9); tf.paragraphs[0].font.color.rgb = C_PURPLE_LT; tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ap(tf, "MLflow · Mosaic", 7, C_PURPLE_PALE, align=PP_ALIGN.CENTER)

    tb(slide, 9.9, 1.25, 0.3, 0.3, "→", 12, C_PURPLE)

    c_db3 = card(slide, 10.3, 1.15, 1.5, 0.35, C_PURPLE, RGBColor(0x1E, 0x1B, 0x2E))
    tf = c_db3.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = "DB SQL Dash"; tf.paragraphs[0].font.size = Pt(8); tf.paragraphs[0].font.color.rgb = C_PURPLE_LT; tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # PALANTIR LANE (Green)
    c_pl = card(slide, 4.4, 1.8, 8.5, 0.7, C_PALANTIR, RGBColor(0x0A, 0x18, 0x12))
    tb(slide, 4.5, 1.95, 0.3, 0.3, "→", 14, C_PALANTIR)

    c_p1 = card(slide, 5.0, 1.9, 2.3, 0.5, C_PALANTIR, RGBColor(0x0F, 0x23, 0x18))
    tf = c_p1.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = "Palantir Foundry"; tf.paragraphs[0].font.size = Pt(9); tf.paragraphs[0].font.color.rgb = C_PALANTIR_LT; tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ap(tf, "Ontology · Pipelines", 7, C_GREEN_LIGHT, align=PP_ALIGN.CENTER)

    tb(slide, 7.4, 2.05, 0.3, 0.3, "→", 12, C_PALANTIR)

    c_p2 = card(slide, 7.8, 1.9, 2.0, 0.5, C_PALANTIR, RGBColor(0x0F, 0x23, 0x18))
    tf = c_p2.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = "Palantir AIP"; tf.paragraphs[0].font.size = Pt(9); tf.paragraphs[0].font.color.rgb = C_PALANTIR_LT; tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ap(tf, "Ontology-driven AI · Ops", 7, C_GREEN_LIGHT, align=PP_ALIGN.CENTER)

    tb(slide, 9.9, 2.05, 0.3, 0.3, "→", 12, C_PALANTIR)

    c_p3 = card(slide, 10.3, 1.95, 1.5, 0.35, C_PALANTIR, RGBColor(0x0F, 0x23, 0x18))
    tf = c_p3.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = "Foundry Apps"; tf.paragraphs[0].font.size = Pt(8); tf.paragraphs[0].font.color.rgb = C_PALANTIR_LT; tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # AWS LANE (Orange)
    c_al = card(slide, 4.4, 2.6, 8.5, 0.8, C_ORANGE, RGBColor(0x14, 0x18, 0x10))
    tb(slide, 4.5, 2.8, 0.3, 0.3, "→", 14, C_ORANGE)

    # Redshift + OpenSearch side by side
    c_a1 = card(slide, 5.0, 2.7, 1.1, 0.5, C_ORANGE, RGBColor(0x2A, 0x30, 0x3E))
    tf = c_a1.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = "Redshift"; tf.paragraphs[0].font.size = Pt(8); tf.paragraphs[0].font.color.rgb = C_ORANGE; tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ap(tf, "SQL · Serverless", 6, C_YELLOW, align=PP_ALIGN.CENTER)

    c_a2 = card(slide, 6.2, 2.7, 1.1, 0.5, C_ORANGE, RGBColor(0x2A, 0x30, 0x3E))
    tf = c_a2.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = "OpenSearch"; tf.paragraphs[0].font.size = Pt(8); tf.paragraphs[0].font.color.rgb = C_ORANGE; tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ap(tf, "Search · RAG", 6, C_YELLOW, align=PP_ALIGN.CENTER)

    tb(slide, 7.4, 2.85, 0.3, 0.3, "→", 12, C_ORANGE)

    # Bedrock + SageMaker
    c_a3 = card(slide, 7.8, 2.7, 1.0, 0.5, C_ORANGE, RGBColor(0x2A, 0x30, 0x3E))
    tf = c_a3.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = "Bedrock"; tf.paragraphs[0].font.size = Pt(8); tf.paragraphs[0].font.color.rgb = C_ORANGE; tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ap(tf, "GenAI · RAG", 6, C_YELLOW, align=PP_ALIGN.CENTER)

    c_a4 = card(slide, 8.9, 2.7, 1.0, 0.5, C_ORANGE, RGBColor(0x2A, 0x30, 0x3E))
    tf = c_a4.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = "SageMaker"; tf.paragraphs[0].font.size = Pt(8); tf.paragraphs[0].font.color.rgb = C_ORANGE; tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ap(tf, "Custom ML", 6, C_YELLOW, align=PP_ALIGN.CENTER)

    tb(slide, 9.95, 2.85, 0.3, 0.3, "→", 12, C_ORANGE)

    # QuickSight + Tableau/PBI
    c_a5 = card(slide, 10.3, 2.65, 1.5, 0.3, C_ORANGE, RGBColor(0x2A, 0x30, 0x3E))
    tf = c_a5.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = "QuickSight"; tf.paragraphs[0].font.size = Pt(8); tf.paragraphs[0].font.color.rgb = C_ORANGE; tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    c_a6 = card(slide, 10.3, 3.0, 1.5, 0.3, RGBColor(0x64, 0x74, 0x8B), RGBColor(0x1A, 0x22, 0x30))
    tf = c_a6.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = "Tableau / PBI"; tf.paragraphs[0].font.size = Pt(8); tf.paragraphs[0].font.color.rgb = C_GRAY3; tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Bottom 3 path descriptions
    paths = [
        ("Databricks Path", "Best for large-scale Spark ETL, data engineering, and ML workloads where Databricks runtime excels. Reads Iceberg natively.", C_PURPLE_LT, C_PURPLE_PALE, C_PURPLE),
        ("Palantir Path", "Best for ontology-driven operations, mission-critical decision-making, and AIP-powered workflows. Connects to S3/Iceberg directly.", C_PALANTIR_LT, C_GREEN_LIGHT, C_PALANTIR),
        ("AWS Analytics Path", "Best for SQL analytics, GenAI/RAG, semantic search, BI dashboards, and custom ML. Native Bedrock integration. Serverless pricing.", C_ORANGE, C_YELLOW, C_ORANGE),
    ]
    px = 0.2
    for ptitle, pbody, ptc, pbc, pbc2 in paths:
        c_p = card(slide, px, 3.65, 4.0, 0.75, pbc2, RGBColor(0x14, 0x18, 0x20))
        tf = c_p.text_frame; tf.word_wrap = True
        tf.paragraphs[0].text = ptitle; tf.paragraphs[0].font.size = Pt(9); tf.paragraphs[0].font.color.rgb = ptc; tf.paragraphs[0].font.bold = True
        ap(tf, pbody, 8, pbc, space_before=Pt(3))
        px += 4.2

    # Iceberg key message
    c_key = card(slide, 0.2, 4.55, 12.9, 0.45, C_BLUE_LIGHT, RGBColor(0x14, 0x1E, 0x30))
    tf = c_key.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = "🧊 Iceberg is the key: "
    run1.font.size = Pt(10); run1.font.color.rgb = C_BLUE_LIGHT; run1.font.bold = True
    run2 = p.add_run()
    run2.text = "One copy of data on S3. Three vendors read it natively. May the best service be aligned to the appropriate use case. Change your mind later — the data stays open."
    run2.font.size = Pt(9); run2.font.color.rgb = C_BLUE_PALE; run2.font.bold = False

    footer(slide, "Amazon Web Services", "Open Architecture | Iceberg + zETL + Best Service Per Use Case")



def slide_14_matrix2(prs):
    """Slide 14: Second Use Case Matrix (WHITE bg)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rows = [
        ("Large-Scale Spark ETL", "Batch data engineering", "✓ Strong — native Spark", "EMR / Glue Spark jobs", "Not primary use case", "Databricks", "db"),
        ("SQL Analytics & Warehousing", "Ad-hoc SQL, dashboards, reporting", "Databricks SQL (costly)", "✓ Redshift Serverless — native Bedrock, zETL", "Not primary use case", "Redshift", "aws"),
        ("Generative AI / RAG", "Natural language Q&A over docs and data", "Limited models, no native RAG", "✓ Bedrock + OpenSearch + Redshift — native KB", "Palantir AIP (ontology-based)", "Bedrock", "aws"),
        ("Full-Text & Semantic Search", "Search millions of docs, hybrid vector", "No native search engine", "✓ OpenSearch — full-text + vector + serverless", "Not primary use case", "OpenSearch", "aws"),
        ("Business Intelligence", "Dashboards, embedded analytics", "Basic SQL dashboards only", "✓ QuickSight — embedded, per-session, Q", "Not primary use case", "QuickSight", "aws"),
        ("Custom ML Model Training", "Train, tune, deploy custom models", "MLflow notebooks", "✓ SageMaker — full MLOps, Bedrock fine-tuning", "Palantir Foundry ML", "SageMaker", "aws"),
        ("Operational Decision-Making", "Ontology-driven workflows, ops", "Not designed for this", "Step Functions + Bedrock Agents", "✓ Palantir Foundry / AIP — purpose-built", "Palantir", "pal"),
        ("Near Real-Time Ingestion", "Stream data from Postgres/Oracle", "Requires custom Spark streaming", "✓ Redshift zETL — auto-replication, near real-time", "Requires connectors", "Redshift zETL", "aws"),
    ]
    _build_matrix_slide(slide,
        "Best Service for Each Use Case — Let the Data Decide",
        "USE CASE ALIGNMENT",
        rows,
        "Iceberg makes this possible:",
        "All services read the same Iceberg tables on S3. No data duplication, no vendor lock-in. Pick the best tool for each use case — and change your mind later without migrating data.")


def slide_15_summary_cta(prs):
    """Slide 15: Summary / CTA (dark bg)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, C_DARK_BG)
    top_bar(slide)

    # Header
    tb(slide, 0.6, 0.2, 10, 0.4, "The AWS Analytics Advantage: Why Add Native Services", 24, C_ORANGE, bold=True)
    tb(slide, 0.6, 0.65, 10, 0.3, "Six reasons to complement Databricks with AWS-native analytics on your S3 data lake", 12, C_GRAY)

    # 3x2 summary grid
    cards_data = [
        ("💰", "Cost at Scale", "S3 storage at $0.023/GB vs. importing into Databricks. Serverless Redshift, OpenSearch, and QuickSight — pay only when you query. No idle cluster costs.", "Up to 90% storage cost reduction"),
        ("🤖", "Native AI Integration", "Redshift, OpenSearch, and S3 are all native Bedrock data sources. One AI layer queries all your data — structured and unstructured — with zero ETL.", "Zero custom integration code"),
        ("🔓", "No Vendor Lock-In", "Data stays in S3 in open formats (Iceberg, Parquet). Any engine can read it. Swap or add services without data migration. You own your data.", "Open formats, open architecture"),
        ("📈", "Best Tool for Each Job", "Redshift for SQL analytics. OpenSearch for search and RAG. QuickSight for BI. SageMaker for ML. Bedrock for GenAI. Each purpose-built, all integrated.", "5 services, one data lake"),
        ("🔒", "Unified Security", "One IAM identity model, one VPC, one Lake Formation governance layer. Data never leaves your AWS account. No cross-platform auth to manage.", "Single identity, single governance"),
        ("🤝", "Co-Existence, Not Replace", "Keep Databricks for Spark workloads and data engineering. Add AWS-native services for analytics, search, BI, and AI. Data flows from Databricks → S3 → everywhere.", "Databricks + AWS = better together"),
    ]
    sx, sy = 0.6, 1.1
    for i, (icon, title, body, stat) in enumerate(cards_data):
        c_sc = card(slide, sx, sy, 3.9, 2.2, C_ORANGE, RGBColor(0x2A, 0x36, 0x48))
        tf = c_sc.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = icon
        tf.paragraphs[0].font.size = Pt(24)
        ap(tf, title, 13, C_ORANGE, True, space_before=Pt(4))
        ap(tf, body, 9, C_GRAY3, space_before=Pt(6))
        ap(tf, stat, 9, C_ORANGE, True, space_before=Pt(8))
        sx += 4.1
        if (i + 1) % 3 == 0:
            sx = 0.6
            sy += 2.4

    # CTA bar
    c_cta = bar(slide, 0, 5.95, 13.333, 0.85, C_ORANGE)
    tb(slide, 0.6, 6.0, 7, 0.35, "Ready to Unlock the Full AWS Analytics Suite?", 17, C_DARK_BG, bold=True)
    tb(slide, 0.6, 6.35, 7, 0.3, "Start with a POC: S3 data lake → Redshift + OpenSearch + Bedrock. See results in days, not months.", 10, C_DARK_BG)

    c_btn = card(slide, 10.0, 6.1, 2.8, 0.55, C_DARK_BG, C_DARK_BG)
    tf = c_btn.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "Let's Build a POC →"
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.color.rgb = C_ORANGE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER



# ══════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Building 15 slides...")

    slide_01_title(prs)
    print("  ✓ Slide 1: Title / Problem Statement")

    slide_02_lockin(prs)
    print("  ✓ Slide 2: Lock-In vs Open Architecture")

    slide_03_s3_datalake(prs)
    print("  ✓ Slide 3: S3 Data Lake Foundation")

    slide_04_iceberg(prs)
    print("  ✓ Slide 4: Apache Iceberg")

    slide_05_redshift_opensearch(prs)
    print("  ✓ Slide 5: Redshift + OpenSearch")

    slide_06_quicksight_sagemaker(prs)
    print("  ✓ Slide 6: QuickSight + SageMaker")

    slide_07_multivendor(prs)
    print("  ✓ Slide 7: Multi-Vendor Open Architecture")

    slide_08_matrix1(prs)
    print("  ✓ Slide 8: Use Case Alignment Matrix")

    slide_09_bedrock(prs)
    print("  ✓ Slide 9: Bedrock AI Layer")

    slide_10_databricks_only(prs)
    print("  ✓ Slide 10: Databricks-Only Current State")

    slide_11_open_architecture(prs)
    print("  ✓ Slide 11: Open Architecture with Iceberg")

    slide_12_aws_solves(prs)
    print("  ✓ Slide 12: AWS Solves Every Problem")

    slide_13_swim_lanes(prs)
    print("  ✓ Slide 13: Swim Lane Flow")

    slide_14_matrix2(prs)
    print("  ✓ Slide 14: Second Use Case Matrix")

    slide_15_summary_cta(prs)
    print("  ✓ Slide 15: Summary / CTA")

    out_dir = os.path.dirname(os.path.abspath(__file__))
    out_path = os.path.join(out_dir, "AWS-Analytics-vs-Databricks-all-slides.pptx")
    prs.save(out_path)
    print(f"\n✅ Saved: {out_path}")
    print(f"   {len(prs.slides)} slides created")


if __name__ == "__main__":
    main()
