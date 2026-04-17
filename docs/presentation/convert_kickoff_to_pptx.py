"""Convert doj-antitrust-kickoff.html to PowerPoint."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Colors
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ORANGE = RGBColor(0xFF, 0x99, 0x00)
LIGHT_ORANGE = RGBColor(0xFF, 0xB8, 0x4D)
WHITE = RGBColor(0xE2, 0xE8, 0xF0)
GRAY = RGBColor(0xA0, 0xAE, 0xC0)
DARK_GRAY = RGBColor(0x71, 0x80, 0x96)
GREEN = RGBColor(0x48, 0xBB, 0x78)
RED = RGBColor(0xFC, 0x81, 0x81)
YELLOW = RGBColor(0xF6, 0xE0, 0x5E)
BLUE = RGBColor(0x63, 0xB3, 0xED)
PURPLE = RGBColor(0x9F, 0x7A, 0xEA)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox

def add_paragraph(tf, text, font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, space_before=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    if space_before:
        p.space_before = Pt(space_before)
    return p

def add_card(slide, left, top, width, height, border_color=ORANGE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x3A)
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    return shape

# ==================== SLIDE 1: TITLE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 1.2, 1.2, 10, 1.0, "Investigative Intelligence:", 40, WHITE, False)
add_text_box(slide, 1.2, 2.0, 10, 1.0, "GenAI for Antitrust Enforcement", 40, WHITE, True)
add_text_box(slide, 1.2, 3.2, 8, 0.5, "DOJ Antitrust Pilot Kickoff 2026", 16, GRAY, True)
add_text_box(slide, 1.2, 4.0, 10, 0.5, "Department of Justice — Antitrust Division", 18, WHITE, True)
add_text_box(slide, 1.2, 4.5, 10, 0.8, "How generative AI, knowledge graphs, and intelligent search transform\n500TB of case evidence into actionable intelligence", 14, DARK_GRAY)
add_text_box(slide, 1.2, 5.6, 4, 0.4, "David Eyre", 14, ORANGE, True)
add_text_box(slide, 1.2, 6.0, 4, 0.4, "Emerging Tech Solutions — Amazon Web Services", 12, GRAY)

# ==================== SLIDE 2: AGENDA ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.3, 11, 0.7, "Agenda", 28, ORANGE, True)

agenda_items = [
    ("1", "Introductions", "Team members, roles, and responsibilities", "5 min"),
    ("2", "Pilot Objectives & Expected Outcomes", "What we're building, why it matters, how we measure success", "10 min"),
    ("3", "Live Demo — Key Platform Components", "Walk through capabilities to help define pilot scope", "20 min"),
    ("4", "PoC Questions — Business & Technical", "Scope, data, environment, timeline, and priorities", "25 min"),
]

for i, (num, title, desc, time) in enumerate(agenda_items):
    y = 1.2 + i * 1.3
    # Number circle
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.5), Inches(y), Inches(0.6), Inches(0.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ORANGE
    shape.line.fill.background()
    shape.text_frame.paragraphs[0].text = num
    shape.text_frame.paragraphs[0].font.size = Pt(18)
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].font.color.rgb = DARK_BG
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text_box(slide, 4.3, y - 0.05, 5.0, 0.4, title, 16, WHITE, True)
    add_text_box(slide, 4.3, y + 0.35, 5.0, 0.4, desc, 12, DARK_GRAY)
    add_text_box(slide, 9.5, y + 0.05, 1.0, 0.4, time, 12, DARK_GRAY, False, PP_ALIGN.RIGHT)

    # Divider line (except last)
    if i < 3:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(y + 0.85), Inches(7.0), Pt(0.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x50)
        shape.line.fill.background()

# ==================== SLIDE 3: PILOT OBJECTIVES ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.3, 11, 0.7, "Pilot Objectives", 28, ORANGE, True)

# What We're Building card
add_card(slide, 0.8, 1.2, 5.6, 2.0, ORANGE)
add_text_box(slide, 1.1, 1.3, 5.0, 0.4, "What We're Building", 16, LIGHT_ORANGE, True)
add_text_box(slide, 1.1, 1.8, 5.0, 1.2, "An AI-powered investigative intelligence platform that ingests your evidence into a centralized data lake, extracts entities and relationships, and surfaces patterns, connections, and insights that accelerate investigation resolution.", 12, WHITE)

# Pilot Scope card
add_card(slide, 0.8, 3.5, 5.6, 2.8, GREEN)
add_text_box(slide, 1.1, 3.6, 5.0, 0.4, "Pilot Scope", 16, GREEN, True)
scope = [
    ("Data Volume", "500 TB of case evidence"),
    ("Division", "DOJ Antitrust Division"),
    ("Timeline", "12-week pilot"),
    ("Budget", "~$200K (ProServe + AWS services)"),
    ("Environment", "AWS GovCloud"),
]
for i, (label, value) in enumerate(scope):
    y = 4.1 + i * 0.45
    add_text_box(slide, 1.1, y, 1.8, 0.4, label, 12, DARK_GRAY)
    add_text_box(slide, 3.0, y, 3.0, 0.4, value, 12, WHITE, True)

# Key Capabilities
add_text_box(slide, 6.8, 1.2, 5.5, 0.4, "Key Capabilities", 16, LIGHT_ORANGE, True)

capabilities = [
    ("📥", "Document ingestion at scale — S3 data lake with AI-powered parsing", ORANGE),
    ("🧠", "AI case briefings — automated intelligence summaries per case", LIGHT_ORANGE),
    ("🕸️", "Entity network discovery — knowledge graph across all evidence", GREEN),
    ("🔍", "Semantic search — natural language queries across 500TB", BLUE),
    ("📊", "Cross-case pattern analysis — connections humans miss", YELLOW),
    ("⚖️", "Prosecution readiness scoring — case strength assessment", RED),
]

for i, (icon, text, color) in enumerate(capabilities):
    y = 1.8 + i * 0.75
    add_card(slide, 6.8, y, 5.5, 0.65, color)
    add_text_box(slide, 7.0, y + 0.05, 5.0, 0.5, f"{icon} {text}", 11, WHITE)

# ==================== SLIDE 4: PILOT USE CASES ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.3, 11, 0.7, "Pilot Use Cases", 28, ORANGE, True)
add_text_box(slide, 0.8, 0.9, 11, 0.4, "Three core capabilities we'll validate during the 12-week pilot", 14, GRAY)

use_cases = [
    ("1", ORANGE, "AI Intelligence Briefing",
     "Analyst selects an investigation → system generates a comprehensive intelligence assessment with 3-level progressive disclosure:",
     ["① Executive summary with key findings & entity statistics",
      "② Finding detail with confidence, graph viz, AI justification",
      "③ Supporting documents with semantic search & excerpts"]),
    ("2", GREEN, "Cross-Case Pattern Analysis",
     "Neptune knowledge graph identifies entities appearing across multiple investigations. Surfaces connections analysts didn't know to look for.",
     ["Entity network visualization",
      "Relationship types & degree centrality",
      "AI-generated hypotheses on connections"]),
    ("3", BLUE, "Semantic Document Search",
     "Natural language queries across all ingested documents. Results ranked by semantic relevance using vector embeddings (Bedrock Titan → OpenSearch kNN).",
     ["Natural language queries",
      "Highlighted excerpts & source drill-down",
      "Cross-investigation relevance ranking"]),
]

for i, (num, color, title, desc, tags) in enumerate(use_cases):
    y = 1.5 + i * 1.9
    add_card(slide, 0.8, y, 11.5, 1.7, color)

    # Number circle
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.1), Inches(y + 0.15), Inches(0.45), Inches(0.45))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.text_frame.paragraphs[0].text = num
    shape.text_frame.paragraphs[0].font.size = Pt(14)
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].font.color.rgb = DARK_BG
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text_box(slide, 1.8, y + 0.1, 10.0, 0.35, title, 15, color, True)
    add_text_box(slide, 1.8, y + 0.5, 10.0, 0.5, desc, 11, WHITE)

    # Tags
    tag_text = "    |    ".join(tags)
    add_text_box(slide, 1.8, y + 1.1, 10.0, 0.3, tag_text, 10, color)

# ==================== SLIDE 5: EXPECTED OUTCOMES ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.3, 11, 0.7, "Expected Outcomes", 28, ORANGE, True)

# Success Criteria table
add_text_box(slide, 0.8, 1.1, 5.5, 0.4, "Success Criteria", 16, LIGHT_ORANGE, True)

# Table headers
add_text_box(slide, 0.8, 1.6, 3.5, 0.3, "Criteria", 11, ORANGE, True)
add_text_box(slide, 4.5, 1.6, 1.8, 0.3, "Target", 11, ORANGE, True, PP_ALIGN.CENTER)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.85), Inches(5.5), Pt(1.5))
shape.fill.solid()
shape.fill.fore_color.rgb = ORANGE
shape.line.fill.background()

criteria = [
    ("Ingestion throughput", "TBD docs/hr"),
    ("AI summary relevance (analyst rating)", "≥ 4.0 / 5"),
    ("Cross-case pattern accuracy", "SME validated"),
    ("Search relevance (precision@10)", "≥ 80%"),
    ("Entity extraction accuracy", "≥ 85%"),
    ("End-to-end query latency", "< 5 seconds"),
    ("GovCloud deployment", "All services ✓"),
]

for i, (crit, target) in enumerate(criteria):
    y = 2.0 + i * 0.38
    add_text_box(slide, 0.8, y, 3.5, 0.3, crit, 11, WHITE)
    add_text_box(slide, 4.5, y, 1.8, 0.3, target, 11, LIGHT_ORANGE, True, PP_ALIGN.CENTER)

# Right side - What Changes
add_text_box(slide, 6.8, 1.1, 5.5, 0.4, "What Changes for Your Team", 16, LIGHT_ORANGE, True)

# 70% stat
add_card(slide, 6.8, 1.6, 5.5, 0.9, GREEN)
add_text_box(slide, 7.0, 1.65, 1.5, 0.7, "70%", 36, GREEN, True)
add_text_box(slide, 8.5, 1.85, 3.5, 0.4, "reduction in manual document review time", 13, WHITE)

change_cards = [
    ("AI discovers patterns across corporate communications, financial records, and market data that keyword search can't find", LIGHT_ORANGE),
    ("Cross-case intelligence surfaces entity connections across historical matters automatically", BLUE),
    ("Prosecution readiness scoring gives leadership visibility into case viability earlier in the process", YELLOW),
]

for i, (text, color) in enumerate(change_cards):
    y = 2.7 + i * 0.85
    add_card(slide, 6.8, y, 5.5, 0.75, color)
    add_text_box(slide, 7.0, y + 0.1, 5.0, 0.5, text, 11, WHITE)

# Beta customer benefit
add_card(slide, 6.8, 5.3, 5.5, 0.8, ORANGE)
add_text_box(slide, 7.0, 5.35, 5.0, 0.25, "BETA CUSTOMER BENEFIT", 10, ORANGE, True, PP_ALIGN.CENTER)
add_text_box(slide, 7.0, 5.65, 5.0, 0.3, "Your requirements shape the AWS OpenSearch Batch Loader roadmap", 11, GRAY, False, PP_ALIGN.CENTER)

# ==================== SLIDE 6: DEMO OVERVIEW ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.3, 11, 0.7, "Demo — Key Platform Components", 28, ORANGE, True)
add_text_box(slide, 0.8, 0.9, 11, 0.4, "Live walkthrough of the investigative intelligence workflow", 14, GRAY)

# Three columns
columns = [
    ("🔍 Investigator Workflow", GREEN, [
        "🤖 AI Briefing & Case Summary",
        "🔎 Semantic & Keyword Search",
        "🕸️ Entity Network Graph",
        "📋 Evidence Triage & Library",
        "📖 Investigative Playbooks",
        "⏳ Timeline Intelligence",
        "🗺️ Geospatial Evidence Map",
        "📌 Entity Tracking & Dossiers",
    ]),
    ("🧠 Intelligence Engine", LIGHT_ORANGE, [
        "📊 Cross-Case Pattern Analysis",
        "🎯 Lead Generation & Scoring",
        "🔗 Entity Resolution",
        "📸 Photo & Face Intelligence",
        "📝 Research Notebook",
        "💡 AI Hypothesis Testing",
    ]),
    ("⚖️ Prosecutor Workflow", YELLOW, [
        "📊 Case Strength Scorecard",
        "✅ Element Assessment",
        "📚 Precedent Matching",
        "📄 Court Document Assembly",
        "🎯 Prosecution Funnel",
        "🏛️ Subject Assessment",
    ]),
]

for col_i, (title, color, items) in enumerate(columns):
    x = 0.8 + col_i * 4.2
    add_card(slide, x, 1.5, 3.8, 4.5, color)
    add_text_box(slide, x + 0.2, 1.6, 3.4, 0.4, title, 14, color, True)

    txBox = add_text_box(slide, x + 0.2, 2.1, 3.4, 0.3, items[0], 12, WHITE)
    for item in items[1:]:
        add_paragraph(txBox.text_frame, item, 12, WHITE, space_before=6)

    # Arrow between columns
    if col_i < 2:
        add_text_box(slide, x + 3.7, 3.2, 0.5, 0.5, "→", 24, ORANGE, True, PP_ALIGN.CENTER)

# Bottom banner
add_text_box(slide, 0.8, 6.3, 11.5, 0.5, "Goal: identify which capabilities matter most for your antitrust workflow to define pilot scope", 14, GREEN, False, PP_ALIGN.CENTER)

# ==================== SLIDE 7: PROJECT TEAM ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.3, 11, 0.7, "Project Team", 28, ORANGE, True)

# Row 1: Executive Sponsors
add_card(slide, 0.8, 1.1, 5.5, 1.2, BLUE)
add_text_box(slide, 1.0, 1.15, 5.0, 0.25, "AWS — DOJ ACCOUNT TEAM", 9, BLUE, True, PP_ALIGN.CENTER)
txBox = add_text_box(slide, 1.0, 1.45, 5.0, 0.3, "Account Mgr: Jack Lan / Brittany Carr", 10, WHITE, False, PP_ALIGN.CENTER)
add_paragraph(txBox.text_frame, "Account SA: Mike Reeves    |    CSM: MJ Yenser", 10, WHITE, False, PP_ALIGN.CENTER)

# Connector
add_text_box(slide, 6.3, 1.4, 0.7, 0.4, "—", 16, DARK_GRAY, False, PP_ALIGN.CENTER)

add_card(slide, 7.0, 1.1, 5.5, 1.2, YELLOW)
add_text_box(slide, 7.2, 1.15, 5.0, 0.25, "DOJ ANTITRUST (ATR) — EXECUTIVE SPONSOR", 9, YELLOW, True, PP_ALIGN.CENTER)
add_text_box(slide, 7.2, 1.5, 5.0, 0.25, "CIO", 10, DARK_GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, 7.2, 1.75, 5.0, 0.3, "Carlos Azuero", 13, WHITE, True, PP_ALIGN.CENTER)

# Row 2: Project Leads
add_card(slide, 0.8, 2.6, 5.5, 1.0, ORANGE)
add_text_box(slide, 1.0, 2.65, 5.0, 0.25, "AWS PROJECT LEAD", 9, ORANGE, True, PP_ALIGN.CENTER)
add_text_box(slide, 1.0, 2.9, 5.0, 0.3, "Emerging Tech Solutions", 13, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, 1.0, 3.2, 5.0, 0.25, "David Eyre • Gabriel Nguyen", 10, GRAY, False, PP_ALIGN.CENTER)

add_text_box(slide, 6.3, 2.85, 0.7, 0.4, "—", 16, DARK_GRAY, False, PP_ALIGN.CENTER)

add_card(slide, 7.0, 2.6, 5.5, 1.0, YELLOW)
add_text_box(slide, 7.2, 2.65, 5.0, 0.25, "ATR PROJECT LEAD", 9, YELLOW, True, PP_ALIGN.CENTER)
add_text_box(slide, 7.2, 2.95, 5.0, 0.3, "TBD", 13, DARK_GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, 7.2, 3.25, 5.0, 0.25, "Technical Decision Maker & Day-to-Day POC", 9, GRAY, False, PP_ALIGN.CENTER)

# Row 3: Implementation Teams - AWS side
aws_teams = [
    ("OpenSearch Service", "Bulk Loader", "Karl Meadows\nSarat Vemulapalli", ORANGE),
    ("Neptune Service", "Graph & Analytics", "Mridula Grandhi", LIGHT_ORANGE),
    ("SSA", "OS: Kevin Fallis / Aruna G.\nNeptune: Renuka Uttarala", "", GREEN),
    ("ProServe", "OpenSearch: TBD\nNeptune: TBD", "", PURPLE),
]

for i, (team, role, names, color) in enumerate(aws_teams):
    x = 0.8 + (i % 2) * 2.85
    y = 4.0 + (i // 2) * 1.6
    add_card(slide, x, y, 2.6, 1.4, color)
    add_text_box(slide, x + 0.1, y + 0.05, 2.4, 0.2, team.upper(), 8, color, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.1, y + 0.3, 2.4, 0.3, role, 9, LIGHT_ORANGE, True, PP_ALIGN.CENTER)
    if names:
        add_text_box(slide, x + 0.1, y + 0.7, 2.4, 0.6, names, 9, WHITE, False, PP_ALIGN.CENTER)

# Connector
add_text_box(slide, 6.3, 4.8, 0.7, 0.4, "—", 16, DARK_GRAY, False, PP_ALIGN.CENTER)

# Row 3: Customer teams
cust_teams = [
    ("ATR Engineering", "TBD", "Infrastructure & DevOps"),
    ("ATR Analysts / SMEs", "TBD", "Domain expertise & UAT"),
    ("ATR Data Team", "TBD", "Data access & pilot dataset"),
    ("ATR Security", "TBD", "GovCloud & compliance"),
]

for i, (team, name, role) in enumerate(cust_teams):
    x = 7.0 + (i % 2) * 2.85
    y = 4.0 + (i // 2) * 1.6
    add_card(slide, x, y, 2.6, 1.4, YELLOW)
    add_text_box(slide, x + 0.1, y + 0.05, 2.4, 0.2, team.upper(), 8, YELLOW, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.1, y + 0.35, 2.4, 0.3, name, 11, DARK_GRAY, False, PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.1, y + 0.7, 2.4, 0.4, role, 9, GRAY, False, PP_ALIGN.CENTER)

# Labels
add_text_box(slide, 0.8, 7.1, 5.5, 0.3, "AMAZON WEB SERVICES", 10, ORANGE, True, PP_ALIGN.CENTER)
add_text_box(slide, 7.0, 7.1, 5.5, 0.3, "DOJ ANTITRUST DIVISION", 10, YELLOW, True, PP_ALIGN.CENTER)

# ==================== SLIDE 8: PROJECT TEAM (SIMPLIFIED) ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.3, 11, 0.7, "Project Team — Roles & Engagement", 28, ORANGE, True)

# AWS Side header
add_text_box(slide, 0.8, 1.1, 5.5, 0.3, "AMAZON WEB SERVICES", 10, ORANGE, True, PP_ALIGN.CENTER)

aws_roles = [
    ("Account Team", "Jack Lan / Brittany Carr • Mike Reeves • MJ Yenser", "Exec Sponsor", BLUE),
    ("Emerging Tech Solutions", "David Eyre • Gabriel Nguyen", "Project Lead", ORANGE),
    ("Service Teams", "OS: Karl Meadows, Sarat Vemulapalli • Neptune: Mridula Grandhi", "Build & Advise", LIGHT_ORANGE),
    ("SSA", "OS: Kevin Fallis / Aruna G. • Neptune: Renuka Uttarala", "Technical Depth", GREEN),
    ("ProServe", "OpenSearch & Neptune — TBD", "Implementation", PURPLE),
]

for i, (role, names, badge, color) in enumerate(aws_roles):
    y = 1.5 + i * 0.85
    add_card(slide, 0.8, y, 5.5, 0.75, color)
    add_text_box(slide, 1.0, y + 0.05, 3.8, 0.3, role, 13, WHITE, True)
    add_text_box(slide, 1.0, y + 0.35, 3.8, 0.3, names, 9, GRAY)
    add_text_box(slide, 4.8, y + 0.15, 1.3, 0.3, badge, 8, color, False, PP_ALIGN.RIGHT)

# Center divider
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(1.5), Pt(2), Inches(4.5))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x50)
shape.line.fill.background()

# Customer Side header
add_text_box(slide, 7.0, 1.1, 5.5, 0.3, "DOJ ANTITRUST DIVISION", 10, YELLOW, True, PP_ALIGN.CENTER)

cust_roles = [
    ("Executive Sponsor — Carlos Azuero (CIO)", None, "Exec Sponsor", None),
    ("Project Lead — TBD", "Technical decision maker & day-to-day POC", "Day-to-Day", None),
    ("Data Engineer / Pipeline — TBD", "Data access, pilot dataset, ingestion pipeline support", "Key role", "⬆ Key role — ongoing engagement throughout pilot"),
    ("Functional SMEs / Analysts — TBD", "Domain expertise, use case validation, UAT feedback", "Key role", "⬆ Key role — guides what we build and validates output"),
    ("Security / Compliance — TBD", "GovCloud access, ATO guidance, encryption requirements", "Light touch", "Light touch — a few hours for initial setup & approvals"),
]

for i, (role, desc, badge, note) in enumerate(cust_roles):
    y = 1.5 + i * 0.85
    add_card(slide, 7.0, y, 5.5, 0.75, YELLOW)
    add_text_box(slide, 7.2, y + 0.05, 4.0, 0.3, role, 12, WHITE, True)
    if desc:
        add_text_box(slide, 7.2, y + 0.32, 4.0, 0.25, desc, 9, GRAY)
    if note:
        note_color = GREEN if "Key role" in note else DARK_GRAY
        add_text_box(slide, 7.2, y + 0.52, 4.0, 0.2, note, 8, note_color)

# Bottom: Implementation model question
add_card(slide, 0.8, 6.0, 11.5, 1.1, ORANGE)
add_text_box(slide, 1.0, 6.05, 11.0, 0.25, "📋 SCOPING QUESTION — IMPLEMENTATION MODEL", 10, ORANGE, True)

options = [
    ("Option A:", "ProServe builds, ATR team maintains and extends"),
    ("Option B:", "ATR team implements with AWS guidance (depends on AWS experience)"),
    ("Option C:", "Hybrid — ProServe leads complex components, ATR handles the rest"),
]
for i, (label, desc) in enumerate(options):
    x = 1.0 + i * 3.8
    add_text_box(slide, x, 6.4, 3.5, 0.25, label, 10, LIGHT_ORANGE, True)
    add_text_box(slide, x, 6.65, 3.5, 0.35, desc, 9, WHITE)

# ==================== SLIDE 9: POC QUESTIONS - BUSINESS ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.3, 11, 0.7, "PoC Questions — Business", 28, ORANGE, True)
add_text_box(slide, 0.8, 0.9, 11, 0.4, "Help us understand your priorities so we can scope the pilot correctly", 14, GRAY)

biz_questions = [
    ("🎯 Mission Outcome", "What is the primary outcome you want from this pilot? Faster case resolution? Cross-case pattern detection? Analyst productivity?", ORANGE),
    ("👥 End Users", "Who will use this? Investigators, analysts, prosecutors, leadership? How many users?", LIGHT_ORANGE),
    ("✅ Success Definition", "What would make you say \"let's go to production\"? What does success look like at the end of 12 weeks?", GREEN),
    ("📋 Case Type", "Is there a specific case or investigation type you want to pilot with first?", BLUE),
    ("⭐ Top Capabilities", "From the demo — which 3 capabilities matter most? (cross-case analysis, AI summaries, pattern discovery, entity graph, search, other)", YELLOW),
    ("⏰ Timeline Driver", "Is there a budget cycle, leadership review, or active investigation driving the timeline?", RED),
]

for i, (title, question, color) in enumerate(biz_questions):
    col = i % 2
    row = i // 2
    x = 0.8 + col * 6.2
    y = 1.5 + row * 1.8
    add_card(slide, x, y, 5.8, 1.6, color)
    add_text_box(slide, x + 0.2, y + 0.1, 5.4, 0.35, title, 14, color, True)
    add_text_box(slide, x + 0.2, y + 0.5, 5.4, 0.9, question, 12, GRAY)

# ==================== SLIDE 10: POC QUESTIONS - TECHNICAL ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.3, 11, 0.7, "PoC Questions — Technical", 28, ORANGE, True)
add_text_box(slide, 0.8, 0.9, 11, 0.4, "Data landscape, environment, and implementation details", 14, GRAY)

tech_cols = [
    ("📦 DATA", ORANGE, [
        "File types? (PDF, Word, Excel, images, email)",
        "Machine-readable or scanned (OCR)?",
        "Already in S3 or needs migration?",
        "Existing hierarchy? (Matter → Case → Doc)",
        "Existing metadata fields?",
        "Sensitivity level? (CUI, LES, PII)",
        "Pilot subset size? (1TB, 10TB, 50TB)",
    ]),
    ("🏗️ ENVIRONMENT", GREEN, [
        "GovCloud? Which region?",
        "Existing AWS services in use?",
        "VPC / networking setup?",
        "Auth requirements? (PIV/CAC, SAML)",
        "FedRAMP / FISMA / ATO for pilot?",
        "Encryption requirements?",
    ]),
    ("👷 IMPLEMENTATION", BLUE, [
        "Internal team or ProServe?",
        "Team size & AWS familiarity?",
        "Preferred model? (AWS builds, hybrid)",
        "Technical decision maker?",
        "Existing tools to integrate?",
        "Budget flexibility beyond $200K?",
    ]),
]

for col_i, (title, color, items) in enumerate(tech_cols):
    x = 0.8 + col_i * 4.0
    add_card(slide, x, 1.5, 3.7, 4.0, color)
    add_text_box(slide, x + 0.2, 1.6, 3.3, 0.3, title, 11, color, True)

    txBox = add_text_box(slide, x + 0.2, 2.0, 3.3, 0.3, items[0], 11, GRAY)
    for item in items[1:]:
        add_paragraph(txBox.text_frame, item, 11, GRAY, space_before=6)

# Scope Priorities
add_card(slide, 0.8, 5.8, 11.5, 1.3, ORANGE)
add_text_box(slide, 1.0, 5.9, 11.0, 0.3, "SCOPE PRIORITIES — RANK 1-5 FOR THE PILOT", 11, ORANGE, True)

priorities = [
    "Document Ingestion & Search", "Entity Extraction & Graph", "AI Case Summaries",
    "Cross-Case Patterns", "Geospatial Mapping", "Access Control", "Custom Batch Loader",
]
priority_text = "    |    ".join(priorities)
add_text_box(slide, 1.0, 6.3, 11.0, 0.5, priority_text, 11, WHITE, False, PP_ALIGN.CENTER)

# ==================== SLIDE 11: NEXT STEPS ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.3, 11, 0.7, "Next Steps", 28, ORANGE, True)

# Timeline
timeline = [
    ("1", "This Week", "Compile answers from today's session. Finalize PoC Success Plan."),
    ("2", "Week 2", "Data assessment. Select pilot dataset. Confirm ProServe engagement."),
    ("3", "Week 3", "Technical deep-dive with engineering team. Environment setup begins."),
    ("4", "Week 4", "Ingestion pipeline build begins. First data flowing through the platform."),
]

for i, (num, time_label, desc) in enumerate(timeline):
    y = 1.2 + i * 1.3

    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), Inches(y), Inches(0.45), Inches(0.45))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x40)
    shape.line.color.rgb = ORANGE
    shape.line.width = Pt(2)
    shape.text_frame.paragraphs[0].text = num
    shape.text_frame.paragraphs[0].font.size = Pt(12)
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].font.color.rgb = ORANGE
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    if i < len(timeline) - 1:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(y + 0.45), Pt(2), Inches(0.85))
        line.fill.solid()
        line.fill.fore_color.rgb = ORANGE
        line.line.fill.background()

    add_text_box(slide, 1.7, y - 0.05, 3.0, 0.3, time_label, 14, LIGHT_ORANGE, True)
    add_text_box(slide, 1.7, y + 0.3, 4.5, 0.5, desc, 12, WHITE)

# Action Items card
add_card(slide, 7.0, 1.2, 5.5, 3.5, GREEN)
add_text_box(slide, 7.2, 1.3, 5.0, 0.4, "📋", 24, GREEN, False, PP_ALIGN.CENTER)
add_text_box(slide, 7.2, 1.7, 5.0, 0.35, "Action Items from Today", 16, GREEN, True, PP_ALIGN.CENTER)

actions = [
    "→ Compile questionnaire answers into Project Statement",
    "→ Identify pilot dataset and begin data assessment",
    "→ Confirm customer technical POC and day-to-day contact",
    "→ Schedule technical deep-dive (Week 3)",
    "→ Finalize ProServe resource allocation",
]
txBox = add_text_box(slide, 7.4, 2.2, 4.8, 0.3, actions[0], 12, WHITE)
for action in actions[1:]:
    add_paragraph(txBox.text_frame, action, 12, WHITE, space_before=6)

# Bottom card
add_card(slide, 7.0, 5.0, 5.5, 1.0, ORANGE)
add_text_box(slide, 7.2, 5.1, 5.0, 0.3, "12 weeks to prove the platform.", 14, ORANGE, True, PP_ALIGN.CENTER)
add_text_box(slide, 7.2, 5.5, 5.0, 0.3, "Solution is built. Code is free. You pay only for AWS services.", 11, GRAY, False, PP_ALIGN.CENTER)

# ==================== SAVE ====================
output_path = "doj-antitrust-kickoff.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
