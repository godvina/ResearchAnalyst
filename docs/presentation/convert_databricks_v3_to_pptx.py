"""AWS/Databricks Co-Exist v3 — title slide + teal theme + all 14 original slides preserved."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Teal/cyan palette
DARK = RGBColor(0x1B, 0x2A, 0x3D)
WHITE = RGBColor(0xF1, 0xF5, 0xF9)
LIGHT = RGBColor(0xCB, 0xD5, 0xE1)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
TEAL = RGBColor(0x06, 0xB6, 0xD4)
ACCENT = RGBColor(0x22, 0xD3, 0xEE)
SOFT_BLUE = RGBColor(0x38, 0xBD, 0xF8)
LIGHT_BLUE = RGBColor(0x60, 0xA5, 0xFA)
GREEN = RGBColor(0x34, 0xD3, 0x99)
SOFT_GREEN = RGBColor(0x6E, 0xE7, 0xB7)
WARN = RGBColor(0xFB, 0x92, 0x3C)
SOFT_RED = RGBColor(0xFC, 0xA5, 0xA5)
RED_TEXT = RGBColor(0xFE, 0xCA, 0xCA)
PURPLE = RGBColor(0xA7, 0x8B, 0xFA)
PAL_GREEN = RGBColor(0x4A, 0xDE, 0x80)
CARD = RGBColor(0x1E, 0x2D, 0x40)
NAVY = RGBColor(0x23, 0x2F, 0x3E)
DEEP = RGBColor(0x0F, 0x17, 0x2A)
LAKE = RGBColor(0x15, 0x2A, 0x45)
ICE_BG = RGBColor(0x0C, 0x19, 0x29)
DB_BG = RGBColor(0x1A, 0x16, 0x25)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def bg(s, c=DARK):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = c

def t(s, l, tp, w, h, txt, sz=18, c=WHITE, b=False, a=PP_ALIGN.LEFT):
    tx = s.shapes.add_textbox(Inches(l), Inches(tp), Inches(w), Inches(h))
    tx.text_frame.word_wrap = True
    p = tx.text_frame.paragraphs[0]
    p.text = txt; p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b; p.alignment = a
    return tx

def ap(tf, txt, sz=14, c=WHITE, b=False, a=PP_ALIGN.LEFT, sb=0):
    p = tf.add_paragraph()
    p.text = txt; p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b; p.alignment = a
    if sb: p.space_before = Pt(sb)

def cd(s, l, tp, w, h, bc=TEAL, fc=CARD):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(tp), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fc; sh.line.color.rgb = bc; sh.line.width = Pt(1)
    return sh

def ln(s, l, tp, w, h, c):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(tp), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()

# ==================== SLIDE 1: CO-EXIST TITLE (NEW) ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
ln(slide, 0, 0, 13.333, 0.04, TEAL)
t(slide, 1.2, 1.5, 10, 0.8, "AWS + Databricks", 42, WHITE)
t(slide, 1.2, 2.3, 10, 0.8, "Co-Existence Strategy", 42, ACCENT, True)
t(slide, 1.2, 3.5, 8, 0.5, "A win for customers. A win for Databricks. A win for AWS.", 16, MUTED)
ln(slide, 1.2, 4.3, 3, Pt(1), RGBColor(0x33, 0x44, 0x55))
t(slide, 1.2, 4.6, 8, 0.8, "Open data lakes on S3 with Apache Iceberg let customers use the best service\nfor each workload — without forcing a single-vendor choice.", 14, LIGHT)
t(slide, 1.2, 6.0, 4, 0.3, "David Eyre", 13, ACCENT, True)
t(slide, 1.2, 6.3, 4, 0.3, "Emerging Tech Solutions — Amazon Web Services", 11, MUTED)
wins = [("🏢 Customer Win", "Best tool for each job. No lock-in. Lower cost at scale.", GREEN),
        ("🤝 Databricks Win", "Keeps Spark workloads. Grows on S3/Iceberg foundation.", PURPLE),
        ("☁️ AWS Win", "More services consumed. Deeper account penetration.", TEAL)]
for i, (title, body, color) in enumerate(wins):
    y = 1.8 + i * 1.5
    cd(slide, 8.5, y, 4.3, 1.2, color)
    t(slide, 8.7, y + 0.1, 4.0, 0.3, title, 13, color, True)
    t(slide, 8.7, y + 0.45, 4.0, 0.5, body, 11, LIGHT)

# ==================== SLIDE 2: ORIGINAL TITLE (softened) ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
ln(slide, 0, 0, 13.333, 0.04, TEAL)
t(slide, 1.0, 1.0, 7, 1.2, "Complement Databricks with\nAWS-Native Analytics", 32, ACCENT, True)
t(slide, 1.0, 2.5, 6.5, 1.2, "Build on Amazon S3 as your open data lake. Access best-of-breed analytics, search, BI, ML, and generative AI services — all natively integrated, all pay-per-use.", 15, MUTED)
t(slide, 1.0, 4.0, 6, 1.2, "Keep Databricks for what it does well.\nExtend with AWS-native services where they add value.\nLet the data stay open on S3.", 14, ACCENT, True)
cd(slide, 7.5, 1.0, 5.3, 5.5, WARN)
t(slide, 7.7, 1.1, 5.0, 0.3, "CHALLENGES WITH A SINGLE-VENDOR APPROACH", 10, WARN, True)
problems = [
    "• All data imported at premium storage rates",
    "• AI/ML models limited to one marketplace",
    "• No native connection to Bedrock, OpenSearch, QuickSight, SageMaker",
    "• Scaling data = scaling spend linearly",
    "• Separate identity, networking, billing from AWS",
    "• Switching costs increase over time",
]
txBox = t(slide, 7.7, 1.5, 5.0, 0.3, problems[0], 12, SOFT_RED)
for p in problems[1:]:
    ap(txBox.text_frame, p, 12, SOFT_RED, sb=6)
t(slide, 0.5, 7.0, 3, 0.3, "Amazon Web Services", 11, TEAL, True)

# ==================== SLIDE 3: LOCK-IN vs OPEN (softened) ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide, RGBColor(0xFF, 0xFF, 0xFF))
ln(slide, 0, 0, 13.333, 0.7, DARK)
t(slide, 0.8, 0.1, 9, 0.5, "Single-Vendor vs. Open Data Lake Architecture", 24, WHITE, True)
t(slide, 10.5, 0.15, 2.5, 0.4, "SIDE BY SIDE", 11, WHITE, True, PP_ALIGN.CENTER)

cd(slide, 0.5, 1.0, 5.8, 5.5, SOFT_RED, RGBColor(0xFE, 0xF2, 0xF2))
t(slide, 0.7, 1.1, 4.5, 0.4, "Single-Vendor Approach", 16, RGBColor(0x1E, 0x29, 0x3B), True)
t(slide, 4.5, 1.15, 1.5, 0.3, "Constrained", 10, RGBColor(0x7F, 0x1D, 0x1D), True, PP_ALIGN.CENTER)
locked = [("📥 Import ALL Data into one platform", "Premium storage, data duplication"),
          ("⚙️ Single compute for everything", "ETL, analytics, ML — one vendor"),
          ("🤖 Vendor AI/ML marketplace", "Limited model selection"),
          ("📊 Basic dashboards", "Not embeddable, per-user pricing"),
          ("👤 End Users", "Single vendor dependency")]
for i, (box, note) in enumerate(locked):
    y = 1.6 + i * 0.85
    t(slide, 0.8, y, 5.2, 0.3, box, 12, RGBColor(0x99, 0x1B, 0x1B), True, PP_ALIGN.CENTER)
    t(slide, 0.8, y + 0.3, 5.2, 0.25, note, 9, RGBColor(0xDC, 0x26, 0x26), False, PP_ALIGN.CENTER)
    if i < 4: t(slide, 3.2, y + 0.55, 0.5, 0.2, "↓", 14, RGBColor(0xDC, 0x26, 0x26), True, PP_ALIGN.CENTER)

cd(slide, 7.0, 1.0, 5.8, 5.5, RGBColor(0x86, 0xEF, 0xAC), RGBColor(0xF0, 0xFD, 0xF4))
t(slide, 7.2, 1.1, 4.5, 0.4, "AWS Open Data Lake", 16, RGBColor(0x1E, 0x29, 0x3B), True)
t(slide, 11.0, 1.15, 1.5, 0.3, "Open & Flexible", 10, RGBColor(0x14, 0x53, 0x2D), True, PP_ALIGN.CENTER)
opn = [("🪣 Amazon S3 Data Lake (Open Formats)", "Lowest cost, open formats, infinite scale"),
       ("📊 Redshift | 🔍 OpenSearch | 📈 QuickSight | 🧪 SageMaker", "Best tool for each job"),
       ("🤖 Amazon Bedrock — Unified AI Layer", "All services feed into Bedrock natively"),
       ("👤 End Users", "Best tool for each job, one AWS bill")]
for i, (box, note) in enumerate(opn):
    y = 1.6 + i * 1.05
    t(slide, 7.3, y, 5.4, 0.35, box, 12, RGBColor(0x06, 0x5F, 0x46), True, PP_ALIGN.CENTER)
    t(slide, 7.3, y + 0.35, 5.4, 0.25, note, 9, RGBColor(0x16, 0xA3, 0x4A), False, PP_ALIGN.CENTER)
    if i < 3: t(slide, 9.7, y + 0.6, 0.5, 0.2, "↓", 14, RGBColor(0x16, 0xA3, 0x4A), True, PP_ALIGN.CENTER)

ln(slide, 0, 6.7, 13.333, 0.8, DARK)
t(slide, 0.8, 6.8, 12, 0.5, "💡 Co-existence: Keep Databricks for Spark workloads. Add AWS-native services for analytics, search, BI, and AI — all reading from the same S3 data lake.", 13, WHITE)

# ==================== SLIDE 4: S3 DATA LAKE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
ln(slide, 0, 0, 13.333, 0.04, TEAL)
t(slide, 0.8, 0.3, 11, 0.5, "Amazon S3: The Open Foundation", 28, ACCENT, True)
t(slide, 0.8, 0.8, 11, 0.4, "Store once in open formats. Query from any service. Scale without limits.", 13, MUTED)

cd(slide, 1.5, 1.4, 7.5, 2.0, SOFT_BLUE, LAKE)
t(slide, 1.7, 1.5, 7.0, 0.4, "🪣 Amazon S3 Data Lake + Apache Iceberg", 20, WHITE, True, PP_ALIGN.CENTER)
t(slide, 1.7, 2.0, 7.0, 0.6, "Open table format on S3. ACID transactions, time travel, schema evolution.\nYour data is never locked in. Any engine reads it. $0.023/GB/month.", 12, SOFT_BLUE, False, PP_ALIGN.CENTER)

stats = [("90%", "Cost Reduction"), ("∞", "Scale"), ("0", "Vendor Lock-In"), ("5+", "AWS Services")]
for i, (num, lbl) in enumerate(stats):
    x = 2.0 + i * 1.8
    t(slide, x, 2.7, 1.5, 0.3, num, 18, ACCENT, True, PP_ALIGN.CENTER)
    t(slide, x, 3.0, 1.5, 0.2, lbl, 9, SOFT_BLUE, False, PP_ALIGN.CENTER)

svcs = [("📊 Redshift", "SQL Analytics"), ("🔍 OpenSearch", "Search & RAG"), ("📈 QuickSight", "BI & Dashboards"),
        ("🧪 SageMaker", "ML & AI Studio"), ("🤖 Bedrock", "Generative AI"), ("🧊 Iceberg", "Open Table Format"),
        ("📋 Glue", "Data Catalog"), ("🔒 Lake Formation", "Governance"), ("🔄 Athena", "Ad-hoc Query")]
for i, (name, role) in enumerate(svcs):
    x = 0.8 + (i % 5) * 2.4; y = 3.5 + (i // 5) * 0.8
    cd(slide, x, y, 2.2, 0.7, TEAL)
    t(slide, x + 0.1, y + 0.05, 2.0, 0.25, name, 11, TEAL, True, PP_ALIGN.CENTER)
    t(slide, x + 0.1, y + 0.35, 2.0, 0.25, role, 9, MUTED, False, PP_ALIGN.CENTER)

bens = [("Cost at Scale", "S3: $0.023/GB vs Databricks Delta: $0.06+/GB. At petabyte scale, millions/yr savings."),
        ("Apache Iceberg: True Openness", "Readable by Spark, Trino, Redshift, Athena, Flink, and Databricks itself."),
        ("Decouple Storage from Compute", "Scale independently. Don't pay for idle compute to keep data accessible."),
        ("Native AWS Integration", "Every AWS analytics service reads S3 natively. No connectors, no ETL.")]
for i, (title, body) in enumerate(bens):
    y = 5.2 + i * 0.55
    t(slide, 0.8, y, 2.5, 0.25, title, 11, ACCENT, True)
    t(slide, 3.5, y, 9.0, 0.25, body, 10, WHITE)

# ==================== SLIDE 5: APACHE ICEBERG ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide, ICE_BG)
ln(slide, 0, 0, 13.333, 0.04, TEAL)
t(slide, 0.8, 0.3, 11, 0.5, "Apache Iceberg on S3: Your Data Goes Everywhere", 26, LIGHT_BLUE, True)
t(slide, 0.8, 0.8, 11, 0.4, "Open table format means zero lock-in. Every major engine reads Iceberg natively — including Databricks.", 13, MUTED)

cd(slide, 0.8, 1.3, 8.0, 1.3, LIGHT_BLUE, RGBColor(0x1E, 0x3A, 0x5F))
t(slide, 1.0, 1.35, 7.5, 0.35, "🧊 Apache Iceberg Tables on Amazon S3", 18, WHITE, True, PP_ALIGN.CENTER)
t(slide, 1.0, 1.75, 7.5, 0.6, "ACID transactions · Time travel · Schema evolution · Partition evolution\nOne copy of data in S3. Every engine reads it. No proprietary format.", 11, SOFT_BLUE, False, PP_ALIGN.CENTER)

t(slide, 0.8, 2.8, 8, 0.3, "ENGINES THAT READ ICEBERG NATIVELY", 10, LIGHT_BLUE, True, PP_ALIGN.CENTER)
engines = [("📊 Redshift", "Zero-copy via Spectrum", TEAL), ("🔄 Athena", "Serverless SQL", TEAL),
           ("⚡ EMR/Spark", "Native read/write", TEAL), ("📋 Glue", "Catalog + ETL", TEAL),
           ("🔥 Databricks", "Reads via UniForm", PURPLE), ("🔺 Trino/Presto", "Full connector", LIGHT_BLUE),
           ("❄️ Snowflake", "External table support", LIGHT_BLUE), ("🌊 Flink", "Streaming read/write", LIGHT_BLUE)]
for i, (name, desc, color) in enumerate(engines):
    x = 0.8 + (i % 4) * 2.1; y = 3.15 + (i // 4) * 0.75
    cd(slide, x, y, 1.9, 0.65, color, ICE_BG)
    t(slide, x + 0.05, y + 0.05, 1.8, 0.25, name, 10, color, True, PP_ALIGN.CENTER)
    t(slide, x + 0.05, y + 0.3, 1.8, 0.25, desc, 8, MUTED, False, PP_ALIGN.CENTER)

cd(slide, 0.8, 4.7, 8.0, 1.8, SOFT_RED, ICE_BG)
t(slide, 1.0, 4.75, 7.5, 0.25, "Iceberg vs. Delta Lake — Format Comparison", 12, SOFT_RED, True)
comps = [("Governance", "Apache Foundation (community)", "Databricks (vendor-managed)"),
         ("Engine support", "Spark, Trino, Flink, Redshift, Athena, Snowflake", "Optimized for Databricks runtime"),
         ("AWS integration", "First-class in Glue, Athena, Redshift, EMR", "Requires Databricks or UniForm"),
         ("Portability", "Move data between engines freely", "Best experience within Databricks")]
for i, (label, ice, delta) in enumerate(comps):
    y = 5.1 + i * 0.3
    t(slide, 1.0, y, 1.5, 0.25, label, 9, MUTED)
    t(slide, 2.6, y, 3.0, 0.25, ice, 9, SOFT_GREEN)
    t(slide, 5.8, y, 3.0, 0.25, delta, 9, PURPLE)

side_b = [("🔓 True Data Portability", "Any engine today and tomorrow. No conversion needed."),
          ("⏪ Time Travel Queries", "Query data at any point in time. Roll back bad writes."),
          ("📐 Schema Evolution", "Add/rename/drop columns without rewriting data."),
          ("📊 Partition Evolution", "Change partitioning without rewriting existing data."),
          ("💰 S3 Economics", "Iceberg on S3 at $0.023/GB. No proprietary storage layer.")]
for i, (title, body) in enumerate(side_b):
    y = 1.3 + i * 1.1
    cd(slide, 9.2, y, 3.8, 0.95, LIGHT_BLUE, ICE_BG)
    t(slide, 9.4, y + 0.05, 3.4, 0.25, title, 11, LIGHT_BLUE, True)
    t(slide, 9.4, y + 0.35, 3.4, 0.5, body, 9, WHITE)

# ==================== SLIDE 6: REDSHIFT + OPENSEARCH ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide, RGBColor(0xFF, 0xFF, 0xFF))
ln(slide, 0, 0, 13.333, 0.65, DARK)
t(slide, 0.8, 0.1, 9, 0.45, "AWS Analytics: Redshift + OpenSearch", 22, WHITE, True)

cd(slide, 0.5, 0.9, 6.0, 5.2, RGBColor(0x93, 0xC5, 0xFD), RGBColor(0xEF, 0xF6, 0xFF))
t(slide, 0.7, 1.0, 5.5, 0.4, "📊 Amazon Redshift", 20, RGBColor(0x1E, 0x40, 0xAF), True)
t(slide, 0.7, 1.4, 5.5, 0.3, "Serverless SQL Analytics & AI Query Layer", 11, RGBColor(0x64, 0x74, 0x8B))
t(slide, 0.7, 1.8, 5.5, 0.6, "Cloud data warehouse with native Bedrock integration. Users ask questions in plain English — Bedrock generates SQL and queries live data.", 11, RGBColor(0x33, 0x41, 0x55))
rf = ["Native Bedrock Knowledge Base data source", "Serverless — pay per query", "Zero-ETL from Aurora, DynamoDB, S3",
      "Queries S3 via Redshift Spectrum", "Federated queries across RDS, Aurora, S3", "ML inference in SQL"]
txBox = t(slide, 0.7, 2.5, 5.5, 0.3, "✓ " + rf[0], 10, RGBColor(0x47, 0x55, 0x69))
for f in rf[1:]: ap(txBox.text_frame, "✓ " + f, 10, RGBColor(0x47, 0x55, 0x69), sb=3)
t(slide, 0.7, 4.8, 5.5, 0.6, "Extends Databricks: Native Bedrock connection for AI queries on live structured data — no ETL needed.", 10, RGBColor(0x06, 0x5F, 0x46), True)

cd(slide, 6.8, 0.9, 6.0, 5.2, RGBColor(0xD8, 0xB4, 0xFE), RGBColor(0xFD, 0xF4, 0xFF))
t(slide, 7.0, 1.0, 5.5, 0.4, "🔍 Amazon OpenSearch", 20, RGBColor(0x7C, 0x3A, 0xED), True)
t(slide, 7.0, 1.4, 5.5, 0.3, "Full-Text Search, Vector Search & RAG", 11, RGBColor(0x64, 0x74, 0x8B))
t(slide, 7.0, 1.8, 5.5, 0.6, "Enterprise search and vector database for RAG. Powers Bedrock Knowledge Bases for semantic search across millions of documents.", 11, RGBColor(0x33, 0x41, 0x55))
of = ["Native Bedrock Knowledge Base vector store", "Full-text + vector hybrid search", "Serverless — scales to zero",
      "k-NN vector search for semantic similarity", "Ingests from S3, Kinesis, DynamoDB", "Dashboards built-in"]
txBox = t(slide, 7.0, 2.5, 5.5, 0.3, "✓ " + of[0], 10, RGBColor(0x47, 0x55, 0x69))
for f in of[1:]: ap(txBox.text_frame, "✓ " + f, 10, RGBColor(0x47, 0x55, 0x69), sb=3)
t(slide, 7.0, 4.8, 5.5, 0.6, "Extends Databricks: Full-text + vector search engine for RAG — a capability Databricks doesn't have natively.", 10, RGBColor(0x06, 0x5F, 0x46), True)

ln(slide, 0, 6.3, 13.333, 0.7, DARK)
t(slide, 0.8, 6.4, 12, 0.4, "🔗 Together: Redshift handles structured queries. OpenSearch handles unstructured search. Both feed into Bedrock natively.", 12, WHITE)

# ==================== SLIDE 7: QUICKSIGHT + SAGEMAKER ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide, RGBColor(0xFF, 0xFF, 0xFF))
ln(slide, 0, 0, 13.333, 0.65, DARK)
t(slide, 0.8, 0.1, 10, 0.45, "AWS Analytics: QuickSight + SageMaker Unified Studio", 22, WHITE, True)

cd(slide, 0.5, 0.9, 6.0, 5.2, RGBColor(0xFD, 0xE0, 0x47), RGBColor(0xFE, 0xFC, 0xE8))
t(slide, 0.7, 1.0, 5.5, 0.4, "📈 Amazon QuickSight", 20, RGBColor(0xA1, 0x62, 0x07), True)
t(slide, 0.7, 1.4, 5.5, 0.3, "Serverless BI with Generative AI (Q)", 11, RGBColor(0x64, 0x74, 0x8B))
t(slide, 0.7, 1.8, 5.5, 0.6, "Cloud-native BI with natural language querying. Embed dashboards in any app. Pay per session — not per named user.", 11, RGBColor(0x33, 0x41, 0x55))
qf = ["QuickSight Q: plain English → charts", "Connects to Redshift, S3, Athena, RDS", "Embeddable dashboards in web apps",
      "Per-session: $0.30 vs $70/user/mo", "SPICE in-memory engine", "Paginated reports for compliance"]
txBox = t(slide, 0.7, 2.5, 5.5, 0.3, "✓ " + qf[0], 10, RGBColor(0x47, 0x55, 0x69))
for f in qf[1:]: ap(txBox.text_frame, "✓ " + f, 10, RGBColor(0x47, 0x55, 0x69), sb=3)
t(slide, 0.7, 4.8, 5.5, 0.6, "Extends Databricks: Embeddable BI with per-session pricing and natural language queries.", 10, RGBColor(0x06, 0x5F, 0x46), True)

cd(slide, 6.8, 0.9, 6.0, 5.2, RGBColor(0x86, 0xEF, 0xAC), RGBColor(0xF0, 0xFD, 0xF4))
t(slide, 7.0, 1.0, 5.5, 0.4, "🧪 SageMaker Unified Studio", 20, RGBColor(0x15, 0x80, 0x3D), True)
t(slide, 7.0, 1.4, 5.5, 0.3, "End-to-End ML & AI Development Platform", 11, RGBColor(0x64, 0x74, 0x8B))
t(slide, 7.0, 1.8, 5.5, 0.6, "One IDE for data engineering, analytics, ML, and GenAI — all connected to your S3 data lake and Bedrock.", 11, RGBColor(0x33, 0x41, 0x55))
sf = ["Unified workspace: SQL, Python, Spark, ML", "Native Bedrock integration for GenAI", "Built-in MLOps: train, tune, deploy, monitor",
      "Access S3, Redshift, Glue from one IDE", "Fine-tune foundation models on your data", "Governed collaboration with Lake Formation"]
txBox = t(slide, 7.0, 2.5, 5.5, 0.3, "✓ " + sf[0], 10, RGBColor(0x47, 0x55, 0x69))
for f in sf[1:]: ap(txBox.text_frame, "✓ " + f, 10, RGBColor(0x47, 0x55, 0x69), sb=3)
t(slide, 7.0, 4.8, 5.5, 0.6, "Extends Databricks: Full MLOps with native Bedrock fine-tuning and deployment.", 10, RGBColor(0x06, 0x5F, 0x46), True)

ln(slide, 0, 6.3, 13.333, 0.7, DARK)
t(slide, 0.8, 6.4, 12, 0.4, "🔗 Together: QuickSight delivers BI at per-session cost. SageMaker gives one IDE for analytics, ML, and GenAI.", 12, WHITE)
