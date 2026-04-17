"""Convert AWS-Analytics-vs-Databricks.html to PowerPoint."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Colors
DARK_BG = RGBColor(0x23, 0x2F, 0x3E)
ORANGE = RGBColor(0xFF, 0x99, 0x00)
WHITE = RGBColor(0xE2, 0xE8, 0xF0)
GRAY = RGBColor(0xA0, 0xAE, 0xC0)
DARK_GRAY = RGBColor(0x71, 0x80, 0x96)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
RED = RGBColor(0xDC, 0x26, 0x26)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
LIGHT_BLUE = RGBColor(0x60, 0xA5, 0xFA)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
LIGHT_PURPLE = RGBColor(0xA7, 0x8B, 0xFA)
PALANTIR_GREEN = RGBColor(0x22, 0xC5, 0x5E)
LIGHT_GREEN = RGBColor(0x86, 0xEF, 0xAC)
MUTED = RGBColor(0x94, 0xA3, 0xB8)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color=DARK_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def tb(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox

def ap(tf, text, size=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT, sb=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    if sb: p.space_before = Pt(sb)
    return p

def card(slide, left, top, width, height, border_color=ORANGE, fill_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or RGBColor(0x1A, 0x23, 0x32)
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    return shape

def bar(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

# ==================== SLIDE 1: TITLE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
bar(slide, 0, 0, 13.333, 0.06, ORANGE)

tb(slide, 1.0, 1.0, 7, 1.2, "Break Free from Vendor Lock-In:\nThe AWS Analytics Advantage", 32, ORANGE, True)
tb(slide, 1.0, 2.5, 6.5, 1.2, "Build on Amazon S3 as your open data lake. Access best-of-breed analytics, search, BI, ML, and generative AI services — all natively integrated, all pay-per-use.", 15, GRAY)
tb(slide, 1.0, 4.0, 6, 1.2, "Keep Databricks for what it does well.\nStop forcing everything through one vendor.\nLet AWS services work together natively.", 14, ORANGE, True)

# Problem box
card(slide, 7.5, 1.0, 5.3, 5.5, ORANGE)
tb(slide, 7.7, 1.1, 5.0, 0.3, "THE DATABRICKS-ONLY PROBLEM", 11, ORANGE, True)
problems = [
    "✗ All data imported into Databricks = expensive storage at premium rates",
    "✗ AI/ML models purchased through Databricks marketplace",
    "✗ No native connection to Bedrock, OpenSearch, QuickSight, or SageMaker",
    "✗ Scaling data = scaling Databricks spend linearly",
    "✗ Separate identity, networking, billing from your AWS environment",
    "✗ Vendor lock-in: migrating away means rewriting everything",
]
txBox = tb(slide, 7.7, 1.5, 5.0, 0.3, problems[0], 12, RGBColor(0xFC, 0x81, 0x81))
for p in problems[1:]:
    ap(txBox.text_frame, p, 12, RGBColor(0xFC, 0x81, 0x81), sb=6)

tb(slide, 0.5, 7.0, 3, 0.3, "Amazon Web Services", 11, ORANGE, True)
tb(slide, 7, 7.0, 6, 0.3, "AWS Analytics Suite | Open Data Lake Strategy", 11, DARK_GRAY, False, PP_ALIGN.RIGHT)

# ==================== SLIDE 2: LOCK-IN vs OPEN ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, RGBColor(0xFF, 0xFF, 0xFF))
bar(slide, 0, 0, 13.333, 0.7, DARK_BG)
tb(slide, 0.8, 0.1, 9, 0.5, "Vendor Lock-In vs. Open Data Lake Architecture", 24, WHITE, True)
tb(slide, 10.5, 0.15, 2.5, 0.4, "SIDE BY SIDE", 11, DARK_BG, True, PP_ALIGN.CENTER)

# Locked column
card(slide, 0.5, 1.0, 5.8, 5.5, RGBColor(0xFC, 0xA5, 0xA5), RGBColor(0xFE, 0xF2, 0xF2))
tb(slide, 0.7, 1.1, 4.5, 0.4, "Databricks-Only Approach", 16, RGBColor(0x1E, 0x29, 0x3B), True)
tb(slide, 4.5, 1.15, 1.5, 0.3, "Locked In", 10, RGBColor(0x7F, 0x1D, 0x1D), True, PP_ALIGN.CENTER)

locked_flow = [
    ("📥 Import ALL Data into Databricks", "Premium storage rates, data duplication"),
    ("⚙️ Databricks Compute for Everything", "ETL, analytics, ML — all one vendor"),
    ("🤖 Databricks AI/ML Marketplace", "Models purchased through Databricks"),
    ("📊 Databricks Dashboards", "Limited BI, no embedded analytics"),
    ("👤 End Users", "Single vendor dependency for everything"),
]
for i, (box, note) in enumerate(locked_flow):
    y = 1.6 + i * 0.85
    tb(slide, 0.8, y, 5.2, 0.3, box, 12, RGBColor(0x99, 0x1B, 0x1B), True, PP_ALIGN.CENTER)
    tb(slide, 0.8, y + 0.3, 5.2, 0.25, note, 9, RED, False, PP_ALIGN.CENTER)
    if i < 4:
        tb(slide, 3.2, y + 0.55, 0.5, 0.2, "↓", 14, RED, True, PP_ALIGN.CENTER)

# Open column
card(slide, 7.0, 1.0, 5.8, 5.5, RGBColor(0x86, 0xEF, 0xAC), RGBColor(0xF0, 0xFD, 0xF4))
tb(slide, 7.2, 1.1, 4.5, 0.4, "AWS Open Data Lake", 16, RGBColor(0x1E, 0x29, 0x3B), True)
tb(slide, 11.0, 1.15, 1.5, 0.3, "Open & Flexible", 10, RGBColor(0x14, 0x53, 0x2D), True, PP_ALIGN.CENTER)

open_flow = [
    ("🪣 Amazon S3 Data Lake (Open Formats)", "Lowest cost, open formats, infinite scale"),
    ("📊 Redshift | 🔍 OpenSearch | 📈 QuickSight | 🧪 SageMaker", "Best tool for each job"),
    ("🤖 Amazon Bedrock — Unified AI Layer", "All services feed into Bedrock natively"),
    ("👤 End Users", "Best tool for each job, one AWS bill"),
]
for i, (box, note) in enumerate(open_flow):
    y = 1.6 + i * 1.05
    tb(slide, 7.3, y, 5.4, 0.35, box, 12, RGBColor(0x06, 0x5F, 0x46), True, PP_ALIGN.CENTER)
    tb(slide, 7.3, y + 0.35, 5.4, 0.25, note, 9, GREEN, False, PP_ALIGN.CENTER)
    if i < 3:
        tb(slide, 9.7, y + 0.6, 0.5, 0.2, "↓", 14, GREEN, True, PP_ALIGN.CENTER)

# Bottom callout
bar(slide, 0, 6.7, 13.333, 0.8, DARK_BG)
tb(slide, 0.8, 6.8, 12, 0.5, "💡 Co-existence strategy: Keep Databricks for Spark workloads. Add AWS-native services for analytics, search, BI, and AI — all reading from the same S3 data lake.", 13, WHITE)

# ==================== SLIDE 3: S3 DATA LAKE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
bar(slide, 0, 0, 13.333, 0.06, ORANGE)
tb(slide, 0.8, 0.3, 11, 0.5, "Amazon S3: The Foundation That Changes Everything", 28, ORANGE, True)
tb(slide, 0.8, 0.8, 11, 0.4, "Store once in open formats. Query from any service. Scale without limits. Pay only for what you store.", 13, GRAY)

# Center lake box
card(slide, 1.5, 1.4, 7.5, 2.0, LIGHT_BLUE, RGBColor(0x1E, 0x40, 0xAF))
tb(slide, 1.7, 1.5, 7.0, 0.4, "🪣 Amazon S3 Data Lake + Apache Iceberg", 20, WHITE, True, PP_ALIGN.CENTER)
tb(slide, 1.7, 2.0, 7.0, 0.6, "Open table format on S3. ACID transactions, time travel, schema evolution.\nYour data is never locked in. Any engine reads it. $0.023/GB/month.", 12, RGBColor(0xBF, 0xDB, 0xFE), False, PP_ALIGN.CENTER)

# Stats
stats = [("90%", "Cost Reduction"), ("∞", "Scale"), ("0", "Vendor Lock-In"), ("5+", "AWS Services")]
for i, (num, lbl) in enumerate(stats):
    x = 2.0 + i * 1.8
    tb(slide, x, 2.7, 1.5, 0.3, num, 18, ORANGE, True, PP_ALIGN.CENTER)
    tb(slide, x, 3.0, 1.5, 0.2, lbl, 9, LIGHT_BLUE, False, PP_ALIGN.CENTER)

# Service chips
services = [
    ("📊 Redshift", "SQL Analytics"), ("🔍 OpenSearch", "Search & RAG"), ("📈 QuickSight", "BI & Dashboards"),
    ("🧪 SageMaker", "ML & AI Studio"), ("🤖 Bedrock", "Generative AI"), ("🧊 Iceberg", "Open Table Format"),
    ("📋 Glue", "Data Catalog"), ("🔒 Lake Formation", "Governance"), ("🔄 Athena", "Ad-hoc Query"),
]
for i, (name, role) in enumerate(services):
    x = 0.8 + (i % 5) * 2.4
    y = 3.5 + (i // 5) * 0.8
    card(slide, x, y, 2.2, 0.7, ORANGE)
    tb(slide, x + 0.1, y + 0.05, 2.0, 0.25, name, 11, ORANGE, True, PP_ALIGN.CENTER)
    tb(slide, x + 0.1, y + 0.35, 2.0, 0.25, role, 9, GRAY, False, PP_ALIGN.CENTER)

# Benefits
benefits = [
    ("Cost at Scale", "S3: $0.023/GB/month vs Databricks Delta: $0.06+/GB. At petabyte scale, the gap is millions/yr."),
    ("Apache Iceberg: True Openness", "Readable by Spark, Trino, Redshift, Athena, Flink, and Databricks itself."),
    ("Decouple Storage from Compute", "Scale independently. Don't pay for idle compute to keep data accessible."),
    ("Native AWS Integration", "Every AWS analytics service reads S3 natively. No connectors, no ETL, no data movement."),
]
for i, (title, body) in enumerate(benefits):
    y = 5.2 + i * 0.55
    tb(slide, 0.8, y, 2.5, 0.25, title, 11, ORANGE, True)
    tb(slide, 3.5, y, 9.0, 0.25, body, 10, WHITE)

# ==================== SLIDE 4: APACHE ICEBERG ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, RGBColor(0x0C, 0x19, 0x29))
bar(slide, 0, 0, 13.333, 0.06, ORANGE)
tb(slide, 0.8, 0.3, 11, 0.5, "Apache Iceberg on S3: Your Data Goes Everywhere", 26, LIGHT_BLUE, True)
tb(slide, 0.8, 0.8, 11, 0.4, "Open table format means zero lock-in. Every major engine reads Iceberg natively — including Databricks.", 13, RGBColor(0x94, 0xA3, 0xB8))

# Center box
card(slide, 0.8, 1.3, 8.0, 1.3, LIGHT_BLUE, RGBColor(0x1E, 0x3A, 0x5F))
tb(slide, 1.0, 1.35, 7.5, 0.35, "🧊 Apache Iceberg Tables on Amazon S3", 18, WHITE, True, PP_ALIGN.CENTER)
tb(slide, 1.0, 1.75, 7.5, 0.6, "ACID transactions · Time travel · Schema evolution · Partition evolution\nOne copy of data in S3. Every engine reads it. No proprietary format.", 11, RGBColor(0xBF, 0xDB, 0xFE), False, PP_ALIGN.CENTER)

# Engines grid
tb(slide, 0.8, 2.8, 8, 0.3, "ENGINES THAT READ ICEBERG NATIVELY", 10, LIGHT_BLUE, True, PP_ALIGN.CENTER)
engines = [
    ("📊 Redshift", "Zero-copy via Spectrum", ORANGE), ("🔄 Athena", "Serverless SQL", ORANGE),
    ("⚡ EMR/Spark", "Native read/write", ORANGE), ("📋 Glue", "Catalog + ETL", ORANGE),
    ("🔥 Databricks", "Reads via UniForm", LIGHT_BLUE), ("🔺 Trino/Presto", "Full connector", LIGHT_BLUE),
    ("❄️ Snowflake", "External table support", LIGHT_BLUE), ("🌊 Flink", "Streaming read/write", LIGHT_BLUE),
]
for i, (name, desc, color) in enumerate(engines):
    x = 0.8 + (i % 4) * 2.1
    y = 3.15 + (i // 4) * 0.75
    card(slide, x, y, 1.9, 0.65, color, RGBColor(0x0C, 0x19, 0x29))
    tb(slide, x + 0.05, y + 0.05, 1.8, 0.25, name, 10, color, True, PP_ALIGN.CENTER)
    tb(slide, x + 0.05, y + 0.3, 1.8, 0.25, desc, 8, RGBColor(0x94, 0xA3, 0xB8), False, PP_ALIGN.CENTER)

# Iceberg vs Delta comparison
card(slide, 0.8, 4.7, 8.0, 1.8, RGBColor(0xFC, 0xA5, 0xA5), RGBColor(0x0C, 0x19, 0x29))
tb(slide, 1.0, 4.75, 7.5, 0.25, "Iceberg vs. Databricks Delta Lake — Format Comparison", 12, RGBColor(0xFC, 0xA5, 0xA5), True)
comparisons = [
    ("Governance", "Apache Foundation (open)", "Databricks (vendor-controlled)"),
    ("Engine support", "Spark, Trino, Flink, Redshift, Athena, Snowflake", "Best in Databricks; limited elsewhere"),
    ("AWS integration", "First-class in Glue, Athena, Redshift, EMR", "Requires Databricks runtime"),
    ("Portability", "Move data between engines freely", "Tied to Databricks ecosystem"),
]
for i, (label, iceberg, delta) in enumerate(comparisons):
    y = 5.1 + i * 0.3
    tb(slide, 1.0, y, 1.5, 0.25, label, 9, RGBColor(0x94, 0xA3, 0xB8))
    tb(slide, 2.6, y, 3.0, 0.25, iceberg, 9, LIGHT_GREEN)
    tb(slide, 5.8, y, 3.0, 0.25, delta, 9, RGBColor(0xFC, 0xA5, 0xA5))

# Side benefits
side_benefits = [
    ("🔓 True Data Portability", "Any engine today and tomorrow. No conversion needed."),
    ("⏪ Time Travel Queries", "Query data at any point in time. Roll back bad writes."),
    ("📐 Schema Evolution", "Add/rename/drop columns without rewriting data."),
    ("📊 Partition Evolution", "Change partitioning without rewriting existing data."),
    ("💰 S3 Economics", "Iceberg on S3 at $0.023/GB. No proprietary storage layer."),
]
for i, (title, body) in enumerate(side_benefits):
    y = 1.3 + i * 1.1
    card(slide, 9.2, y, 3.8, 0.95, LIGHT_BLUE, RGBColor(0x0C, 0x19, 0x29))
    tb(slide, 9.4, y + 0.05, 3.4, 0.25, title, 11, LIGHT_BLUE, True)
    tb(slide, 9.4, y + 0.35, 3.4, 0.5, body, 9, WHITE)

# ==================== SLIDE 5: REDSHIFT + OPENSEARCH ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, RGBColor(0xFF, 0xFF, 0xFF))
bar(slide, 0, 0, 13.333, 0.65, DARK_BG)
tb(slide, 0.8, 0.1, 9, 0.45, "AWS Analytics Services: Redshift + OpenSearch", 22, WHITE, True)

# Redshift card
card(slide, 0.5, 0.9, 6.0, 5.2, RGBColor(0x93, 0xC5, 0xFD), RGBColor(0xEF, 0xF6, 0xFF))
tb(slide, 0.7, 1.0, 5.5, 0.4, "📊 Amazon Redshift", 20, RGBColor(0x1E, 0x40, 0xAF), True)
tb(slide, 0.7, 1.4, 5.5, 0.3, "Serverless SQL Analytics & AI Query Layer", 11, RGBColor(0x64, 0x74, 0x8B))
tb(slide, 0.7, 1.8, 5.5, 0.6, "The fastest cloud data warehouse with native Bedrock integration. Users ask questions in plain English — Bedrock generates SQL and queries live data.", 11, RGBColor(0x33, 0x41, 0x55))
redshift_features = ["Native Bedrock Knowledge Base data source", "Serverless — pay per query", "Zero-ETL from Aurora, DynamoDB, S3", "Queries S3 via Redshift Spectrum", "Federated queries across RDS, Aurora, S3", "ML inference in SQL with Redshift ML"]
txBox = tb(slide, 0.7, 2.5, 5.5, 0.3, "✓ " + redshift_features[0], 10, RGBColor(0x47, 0x55, 0x69))
for f in redshift_features[1:]:
    ap(txBox.text_frame, "✓ " + f, 10, RGBColor(0x47, 0x55, 0x69), sb=3)
tb(slide, 0.7, 4.8, 5.5, 0.6, "Databricks gap: No native Bedrock connection. Must ETL to S3, then Bedrock reads stale copies.", 10, RED, True)

# OpenSearch card
card(slide, 6.8, 0.9, 6.0, 5.2, RGBColor(0xD8, 0xB4, 0xFE), RGBColor(0xFD, 0xF4, 0xFF))
tb(slide, 7.0, 1.0, 5.5, 0.4, "🔍 Amazon OpenSearch", 20, RGBColor(0x7C, 0x3A, 0xED), True)
tb(slide, 7.0, 1.4, 5.5, 0.3, "Full-Text Search, Vector Search & RAG", 11, RGBColor(0x64, 0x74, 0x8B))
tb(slide, 7.0, 1.8, 5.5, 0.6, "Enterprise search and vector database for RAG applications. Powers Bedrock Knowledge Bases for semantic search across millions of documents.", 11, RGBColor(0x33, 0x41, 0x55))
os_features = ["Native Bedrock Knowledge Base vector store", "Full-text + vector hybrid search", "Serverless — scales to zero", "k-NN vector search for semantic similarity", "Ingests from S3, Kinesis, DynamoDB", "Dashboards built-in"]
txBox = tb(slide, 7.0, 2.5, 5.5, 0.3, "✓ " + os_features[0], 10, RGBColor(0x47, 0x55, 0x69))
for f in os_features[1:]:
    ap(txBox.text_frame, "✓ " + f, 10, RGBColor(0x47, 0x55, 0x69), sb=3)
tb(slide, 7.0, 4.8, 5.5, 0.6, "Databricks gap: No equivalent search engine. Must use third-party for RAG, adding cost and complexity.", 10, RED, True)

bar(slide, 0, 6.3, 13.333, 0.7, DARK_BG)
tb(slide, 0.8, 6.4, 12, 0.4, "🔗 Together: Redshift handles structured queries. OpenSearch handles unstructured search. Both feed into Bedrock as native data sources.", 12, WHITE)

# ==================== SLIDE 6: QUICKSIGHT + SAGEMAKER ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, RGBColor(0xFF, 0xFF, 0xFF))
bar(slide, 0, 0, 13.333, 0.65, DARK_BG)
tb(slide, 0.8, 0.1, 10, 0.45, "AWS Analytics Services: QuickSight + SageMaker", 22, WHITE, True)

# QuickSight card
card(slide, 0.5, 0.9, 6.0, 5.2, RGBColor(0xFD, 0xE0, 0x47), RGBColor(0xFE, 0xFC, 0xE8))
tb(slide, 0.7, 1.0, 5.5, 0.4, "📈 Amazon QuickSight", 20, RGBColor(0xA1, 0x62, 0x07), True)
tb(slide, 0.7, 1.4, 5.5, 0.3, "Serverless BI with Generative AI (Q)", 11, RGBColor(0x64, 0x74, 0x8B))
tb(slide, 0.7, 1.8, 5.5, 0.6, "Cloud-native BI with natural language querying. Embed interactive dashboards in any application. Pay per session — not per named user.", 11, RGBColor(0x33, 0x41, 0x55))
qs_features = ["QuickSight Q: plain English → charts", "Connects to Redshift, S3, Athena, RDS", "Embeddable dashboards in web apps", "Per-session: $0.30/session vs $70/user/mo", "SPICE in-memory engine", "Paginated reports for compliance"]
txBox = tb(slide, 0.7, 2.5, 5.5, 0.3, "✓ " + qs_features[0], 10, RGBColor(0x47, 0x55, 0x69))
for f in qs_features[1:]:
    ap(txBox.text_frame, "✓ " + f, 10, RGBColor(0x47, 0x55, 0x69), sb=3)
tb(slide, 0.7, 4.8, 5.5, 0.6, "Databricks gap: Basic dashboards, not embeddable, require Databricks compute. Every viewer needs a license.", 10, RED, True)

# SageMaker card
card(slide, 6.8, 0.9, 6.0, 5.2, RGBColor(0x86, 0xEF, 0xAC), RGBColor(0xF0, 0xFD, 0xF4))
tb(slide, 7.0, 1.0, 5.5, 0.4, "🧪 SageMaker Unified Studio", 20, RGBColor(0x15, 0x80, 0x3D), True)
tb(slide, 7.0, 1.4, 5.5, 0.3, "End-to-End ML & AI Development Platform", 11, RGBColor(0x64, 0x74, 0x8B))
tb(slide, 7.0, 1.8, 5.5, 0.6, "One IDE for data engineering, analytics, ML model development, and generative AI — all connected to your S3 data lake and Bedrock.", 11, RGBColor(0x33, 0x41, 0x55))
sm_features = ["Unified workspace: SQL, Python, Spark, ML", "Native Bedrock integration for GenAI", "Built-in MLOps: train, tune, deploy, monitor", "Access S3, Redshift, Glue from one IDE", "Fine-tune foundation models on your data", "Governed collaboration with Lake Formation"]
txBox = tb(slide, 7.0, 2.5, 5.5, 0.3, "✓ " + sm_features[0], 10, RGBColor(0x47, 0x55, 0x69))
for f in sm_features[1:]:
    ap(txBox.text_frame, "✓ " + f, 10, RGBColor(0x47, 0x55, 0x69), sb=3)
tb(slide, 7.0, 4.8, 5.5, 0.6, "Databricks gap: Notebooks siloed from AWS AI services. Models can't natively deploy to Bedrock or integrate with QuickSight Q.", 10, RED, True)

bar(slide, 0, 6.3, 13.333, 0.7, DARK_BG)
tb(slide, 0.8, 6.4, 12, 0.4, "🔗 Together: QuickSight delivers BI at per-session cost. SageMaker gives data teams one IDE for analytics, ML, and GenAI — all on the same data lake.", 12, WHITE)

# ==================== SLIDE 7: MULTI-VENDOR ARCHITECTURE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, RGBColor(0x0F, 0x17, 0x2A))
bar(slide, 0, 0, 13.333, 0.06, ORANGE)
tb(slide, 0.8, 0.3, 11, 0.5, "Open Data, Open Choice: Best Service for Every Use Case", 26, ORANGE, True)
tb(slide, 0.8, 0.8, 11, 0.3, "With Iceberg tables on S3, your data is accessible to every platform. Pick the best tool for each job.", 13, RGBColor(0x94, 0xA3, 0xB8))

# Foundation
card(slide, 1.5, 1.3, 10.3, 0.8, LIGHT_BLUE, RGBColor(0x1E, 0x3A, 0x5F))
tb(slide, 1.7, 1.35, 10.0, 0.3, "🧊 Apache Iceberg Tables on Amazon S3", 16, WHITE, True, PP_ALIGN.CENTER)
tb(slide, 1.7, 1.65, 10.0, 0.3, "One copy of data. Open format. Readable by every platform below — simultaneously.", 11, RGBColor(0xBF, 0xDB, 0xFE), False, PP_ALIGN.CENTER)

tb(slide, 6.2, 2.2, 1.0, 0.3, "↓  Same data, read by all  ↓", 10, LIGHT_BLUE, False, PP_ALIGN.CENTER)

# Three vendor cards
vendors = [
    ("🔥 Databricks", "Data Engineering & Spark", PURPLE, LIGHT_PURPLE,
     ["Large-scale Spark ETL", "Complex data transformations", "Data science notebooks", "Delta-to-Iceberg via UniForm", "Existing Spark workloads"]),
    ("☁️ AWS Analytics Suite", "AI, Search, BI & Real-Time", ORANGE, ORANGE,
     ["Redshift: SQL + Bedrock AI queries", "OpenSearch: Full-text + RAG", "QuickSight: Embedded BI", "SageMaker: ML training + deploy", "Bedrock: GenAI across all data", "Athena: Serverless SQL on Iceberg"]),
    ("🛡️ Palantir Foundry", "Operational Intelligence", PALANTIR_GREEN, RGBColor(0x4A, 0xDE, 0x80),
     ["Ontology-based operational apps", "Cross-domain data fusion", "Mission-critical workflows", "Reads Iceberg from S3 directly", "GovCloud compatible"]),
]

for i, (name, role, border, text_color, items) in enumerate(vendors):
    x = 0.8 + i * 4.2
    card(slide, x, 2.6, 3.8, 3.5, border, RGBColor(0x0F, 0x17, 0x2A))
    tb(slide, x + 0.2, 2.7, 3.4, 0.35, name, 16, text_color, True, PP_ALIGN.CENTER)
    tb(slide, x + 0.2, 3.05, 3.4, 0.25, role, 10, RGBColor(0x94, 0xA3, 0xB8), False, PP_ALIGN.CENTER)
    txBox = tb(slide, x + 0.2, 3.4, 3.4, 0.3, "→ " + items[0], 10, WHITE)
    for item in items[1:]:
        ap(txBox.text_frame, "→ " + item, 10, WHITE, sb=3)

# Key message
card(slide, 0.8, 6.3, 11.7, 0.7, ORANGE, RGBColor(0x0F, 0x17, 0x2A))
tb(slide, 1.0, 6.35, 11.3, 0.3, "The data stays in S3. The format stays open. The choice stays yours.", 14, ORANGE, True, PP_ALIGN.CENTER)
tb(slide, 1.0, 6.65, 11.3, 0.25, "No vendor can hold your data hostage. Add or remove platforms without data migration.", 11, WHITE, False, PP_ALIGN.CENTER)

# ==================== SLIDE 8: USE CASE MATRIX ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, RGBColor(0xFF, 0xFF, 0xFF))
bar(slide, 0, 0, 13.333, 0.65, DARK_BG)
tb(slide, 0.8, 0.1, 9, 0.45, "Use Case Alignment: Best Platform for Each Workload", 22, WHITE, True)

# Table headers
header_y = 0.85
tb(slide, 0.5, header_y, 2.5, 0.3, "USE CASE", 9, RGBColor(0x47, 0x55, 0x69), True)
tb(slide, 3.2, header_y, 2.3, 0.3, "DATABRICKS", 9, RGBColor(0x7C, 0x3A, 0xED), True)
tb(slide, 5.7, header_y, 2.5, 0.3, "AWS ANALYTICS", 9, RGBColor(0xB4, 0x53, 0x09), True)
tb(slide, 8.4, header_y, 2.3, 0.3, "PALANTIR", 9, RGBColor(0x15, 0x80, 0x3D), True)
tb(slide, 10.9, header_y, 2.0, 0.3, "BEST FIT", 9, RGBColor(0x1E, 0x40, 0xAF), True)

matrix = [
    ("Large-Scale ETL / Spark", "✓ Native Spark, DLT", "EMR/Glue Spark", "Not primary", "Databricks", PURPLE),
    ("Generative AI / RAG", "No native Bedrock", "✓ Bedrock + OS + Redshift", "AIP for ops AI", "AWS (Bedrock)", ORANGE),
    ("SQL Analytics / DW", "DB SQL (costly)", "✓ Redshift Serverless", "Not primary", "AWS (Redshift)", ORANGE),
    ("Business Intelligence", "Basic, not embeddable", "✓ QuickSight embedded", "Foundry dashboards", "AWS (QuickSight)", ORANGE),
    ("Full-Text & Semantic Search", "No search engine", "✓ OpenSearch hybrid", "Object search only", "AWS (OpenSearch)", ORANGE),
    ("ML Model Training", "✓ MLflow, Feature Store", "✓ SageMaker + Bedrock", "Palantir AIP", "Both", ORANGE),
    ("Operational Decision Support", "Not designed for this", "Step Functions + Bedrock", "✓ Foundry / AIP", "Palantir", PALANTIR_GREEN),
    ("Data Science Notebooks", "✓ Great notebook UX", "✓ SageMaker Studio", "Code Workbook", "Both", ORANGE),
]

for i, (uc, db, aws, pal, best, best_color) in enumerate(matrix):
    y = 1.2 + i * 0.6
    bg_color = RGBColor(0xFA, 0xFA, 0xFA) if i % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
    tb(slide, 0.5, y, 2.5, 0.3, uc, 10, RGBColor(0x1E, 0x29, 0x3B), True)
    tb(slide, 3.2, y, 2.3, 0.3, db, 9, RGBColor(0x47, 0x55, 0x69))
    tb(slide, 5.7, y, 2.5, 0.3, aws, 9, RGBColor(0x47, 0x55, 0x69))
    tb(slide, 8.4, y, 2.3, 0.3, pal, 9, RGBColor(0x47, 0x55, 0x69))
    tb(slide, 10.9, y, 2.0, 0.3, best, 10, best_color, True)

bar(slide, 0, 6.3, 13.333, 0.7, DARK_BG)
tb(slide, 0.8, 6.4, 12, 0.4, "🧊 Iceberg advantage: All three platforms read the same tables on S3. No data duplication. Align the best service to each use case.", 12, WHITE)

# ==================== SLIDE 9: BEDROCK AI LAYER ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
bar(slide, 0, 0, 13.333, 0.06, ORANGE)
tb(slide, 0.8, 0.3, 11, 0.5, "Amazon Bedrock: The AI Layer That Ties It All Together", 26, ORANGE, True)
tb(slide, 0.8, 0.8, 11, 0.3, "Every AWS analytics service feeds into Bedrock natively — one AI layer for all your data", 13, GRAY)

# Bedrock center
card(slide, 1.5, 1.3, 10.3, 0.8, ORANGE, RGBColor(0xFF, 0x99, 0x00))
tb(slide, 1.7, 1.35, 10.0, 0.3, "Amazon Bedrock Knowledge Bases", 18, DARK_BG, True, PP_ALIGN.CENTER)
tb(slide, 1.7, 1.65, 10.0, 0.3, "Natural language interface to ALL your data — structured and unstructured — through one unified API", 11, RGBColor(0x45, 0x1A, 0x03), False, PP_ALIGN.CENTER)

# Data sources
sources = [
    ("📊 Redshift", "Structured Data", "Bedrock generates SQL, queries live tables. Real-time answers. Zero ETL.", "NATIVE SOURCE"),
    ("🔍 OpenSearch", "Vector Store / RAG", "Semantic search across millions of documents. Vector embeddings for RAG.", "NATIVE SOURCE"),
    ("🪣 S3 Data Lake", "Documents & Files", "PDFs, images, CSVs ingested directly. Auto-chunked and indexed.", "NATIVE SOURCE"),
    ("🧪 SageMaker", "Custom Models", "Fine-tuned models deployed as Bedrock endpoints. Custom ML inference.", "DIRECT INTEGRATION"),
]
for i, (name, stype, desc, tag) in enumerate(sources):
    x = 0.8 + i * 3.1
    card(slide, x, 2.4, 2.9, 2.5, ORANGE)
    tb(slide, x + 0.1, 2.5, 2.7, 0.3, name, 13, ORANGE, True, PP_ALIGN.CENTER)
    tb(slide, x + 0.1, 2.8, 2.7, 0.25, stype, 10, GRAY, False, PP_ALIGN.CENTER)
    tb(slide, x + 0.1, 3.15, 2.7, 0.8, desc, 9, WHITE)
    tag_color = GREEN if "NATIVE" in tag else BLUE
    tb(slide, x + 0.1, 4.3, 2.7, 0.25, tag, 8, tag_color, True, PP_ALIGN.CENTER)

# What Databricks Can't Do vs What AWS Delivers
card(slide, 0.8, 5.2, 5.8, 1.5, RGBColor(0xFC, 0x81, 0x81))
tb(slide, 1.0, 5.25, 5.4, 0.3, "🚫 What Databricks Can't Do", 12, ORANGE, True)
tb(slide, 1.0, 5.6, 5.4, 0.9, "No native Bedrock connection. Must export data to S3, then Bedrock reads stale copies. AI models purchased through Databricks marketplace — not Bedrock foundation models.", 10, WHITE)

card(slide, 7.0, 5.2, 5.8, 1.5, GREEN)
tb(slide, 7.2, 5.25, 5.4, 0.3, "✅ What AWS Delivers", 12, ORANGE, True)
tb(slide, 7.2, 5.6, 5.4, 0.9, "One Bedrock Knowledge Base queries Redshift (structured), OpenSearch (vectors), and S3 (documents) simultaneously. One question, all data sources. One API, one bill.", 10, WHITE)

# ==================== SLIDE 10: CURRENT STATE (DATABRICKS-ONLY) ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, RGBColor(0x1A, 0x16, 0x25))
bar(slide, 0, 0, 13.333, 0.06, PURPLE)
tb(slide, 0.8, 0.3, 11, 0.5, "Current State: Databricks-Only Architecture", 24, LIGHT_PURPLE, True)
tb(slide, 0.8, 0.7, 11, 0.3, "All data flows through one vendor — creating cost, scale, governance, and AI limitations", 12, RGBColor(0x94, 0xA3, 0xB8))

# Flow: Sources → Glue → Data Lake → Databricks → AI → BI
flow_items = [
    ("Postgres / Oracle / Other", "Data Sources", RGBColor(0x64, 0x74, 0x8B)),
    ("Glue / Informatica", "ETL to S3", PURPLE),
    ("Data Lake S3", "Databricks Catalog", LIGHT_BLUE),
    ("Databricks", "Delta / Proton Query", PURPLE),
    ("Databricks AI?", "Limited models", PURPLE),
    ("BI Layer", "QuickSight? Tableau?", PURPLE),
]
for i, (name, desc, color) in enumerate(flow_items):
    x = 0.5 + i * 2.1
    card(slide, x, 1.2, 1.9, 0.9, color, RGBColor(0x1A, 0x16, 0x25))
    tb(slide, x + 0.05, 1.25, 1.8, 0.3, name, 10, LIGHT_PURPLE, True, PP_ALIGN.CENTER)
    tb(slide, x + 0.05, 1.6, 1.8, 0.25, desc, 8, RGBColor(0xC4, 0xB5, 0xFD), False, PP_ALIGN.CENTER)
    if i < 5:
        tb(slide, x + 1.8, 1.4, 0.3, 0.3, "→", 16, PURPLE, True)

# Problem callouts
problems = [
    ("⚠ Near Real-Time Gap", "Dumping to S3 is not real time. ETL latency = stale data."),
    ("⚠ Governance Lock-In", "Governance through Databricks catalog only. Vendor lock-in for access control."),
    ("⚠ AI Model Limitations", "Not all models available. Bedrock has more models at better prices."),
    ("⚠ Cost & Scale Challenges", "Loading all data into Databricks is costly. Scaling challenges. BI connectors uncertain."),
]
for i, (title, body) in enumerate(problems):
    x = 0.5 + i * 3.15
    card(slide, x, 2.5, 2.9, 1.5, RGBColor(0xFC, 0xA5, 0xA5), RGBColor(0x1A, 0x16, 0x25))
    tb(slide, x + 0.1, 2.55, 2.7, 0.25, title, 10, RGBColor(0xFC, 0xA5, 0xA5), True)
    tb(slide, x + 0.1, 2.85, 2.7, 0.9, body, 9, RGBColor(0xFE, 0xCA, 0xCA))

tb(slide, 0.8, 4.3, 11.5, 0.5, "Core Question: Should one vendor own your data movement, governance, warehousing, AI, and BI?", 13, RGBColor(0xFC, 0xA5, 0xA5), True, PP_ALIGN.CENTER)

# ==================== SLIDE 11: OPEN ARCHITECTURE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, RGBColor(0x0F, 0x17, 0x2A))
bar(slide, 0, 0, 13.333, 0.06, ORANGE)
tb(slide, 0.8, 0.2, 11, 0.4, "Open Architecture: Iceberg + Best Service for Each Use Case", 22, ORANGE, True)
tb(slide, 0.8, 0.6, 11, 0.3, "zETL from Postgres/Oracle → Redshift → Iceberg on S3. Three swim lanes. Each vendor stays in its lane.", 11, RGBColor(0x94, 0xA3, 0xB8))

# Left: Sources → zETL → Iceberg
tb(slide, 0.5, 1.1, 1.0, 0.3, "Sources", 9, DARK_GRAY, True, PP_ALIGN.CENTER)
sources_list = ["Postgres", "Oracle", "Other"]
for i, s in enumerate(sources_list):
    card(slide, 0.3, 1.5 + i * 0.55, 1.2, 0.45, RGBColor(0x64, 0x74, 0x8B), RGBColor(0x0F, 0x17, 0x2A))
    tb(slide, 0.35, 1.55 + i * 0.55, 1.1, 0.25, s, 10, WHITE, True, PP_ALIGN.CENTER)

tb(slide, 1.5, 2.0, 0.3, 0.3, "→", 16, ORANGE, True)

card(slide, 1.8, 1.7, 1.0, 0.8, ORANGE, RGBColor(0x0F, 0x17, 0x2A))
tb(slide, 1.85, 1.8, 0.9, 0.25, "zETL", 12, ORANGE, True, PP_ALIGN.CENTER)
tb(slide, 1.85, 2.1, 0.9, 0.2, "Real-time", 8, RGBColor(0xFC, 0xD3, 0x4D), False, PP_ALIGN.CENTER)

tb(slide, 2.8, 2.0, 0.3, 0.3, "→", 16, LIGHT_BLUE, True)

card(slide, 3.1, 1.5, 1.5, 0.5, BLUE, RGBColor(0x0F, 0x17, 0x2A))
tb(slide, 3.15, 1.55, 1.4, 0.3, "Redshift Tables", 10, LIGHT_BLUE, True, PP_ALIGN.CENTER)
tb(slide, 3.8, 2.05, 0.3, 0.2, "↓", 12, LIGHT_BLUE, True, PP_ALIGN.CENTER)
card(slide, 3.1, 2.3, 1.5, 0.6, LIGHT_BLUE, RGBColor(0x0F, 0x17, 0x2A))
tb(slide, 3.15, 2.35, 1.4, 0.25, "🧊 Iceberg on S3", 10, LIGHT_BLUE, True, PP_ALIGN.CENTER)
tb(slide, 3.15, 2.6, 1.4, 0.2, "Open · Any engine", 8, RGBColor(0x93, 0xC5, 0xFD), False, PP_ALIGN.CENTER)

# Three swim lanes
lanes = [
    ("Databricks", PURPLE, LIGHT_PURPLE, [("Spark ETL · Delta", ""), ("Databricks AI", "MLflow"), ("DB SQL Dash", "")]),
    ("Palantir", PALANTIR_GREEN, RGBColor(0x4A, 0xDE, 0x80), [("Foundry", "Ontology"), ("Palantir AIP", "Ops AI"), ("Foundry Apps", "")]),
    ("AWS Analytics", ORANGE, ORANGE, [("Redshift + OpenSearch", "SQL · Search"), ("Bedrock + SageMaker", "GenAI · ML"), ("QuickSight + BI", "Tableau/PBI")]),
]

for i, (name, border, text_color, boxes) in enumerate(lanes):
    y = 3.3 + i * 1.2
    # Lane background
    card(slide, 4.8, y, 8.0, 1.0, border, RGBColor(0x0F, 0x17, 0x2A))
    tb(slide, 4.9, y + 0.05, 0.5, 0.3, "→", 14, border, True)
    # Three boxes in lane
    for j, (box_name, box_desc) in enumerate(boxes):
        x = 5.4 + j * 2.6
        tb(slide, x, y + 0.1, 2.3, 0.25, box_name, 10, text_color, True, PP_ALIGN.CENTER)
        if box_desc:
            tb(slide, x, y + 0.4, 2.3, 0.2, box_desc, 8, RGBColor(0x94, 0xA3, 0xB8), False, PP_ALIGN.CENTER)
        if j < 2:
            tb(slide, x + 2.2, y + 0.15, 0.3, 0.2, "→", 12, border, True)

# Bottom
tb(slide, 0.8, 6.8, 11.5, 0.4, "🧊 Iceberg is the key: One copy of data on S3. Three vendors read it natively. Change your mind later — the data stays open.", 12, RGBColor(0x93, 0xC5, 0xFD), False, PP_ALIGN.CENTER)

# ==================== SLIDE 12: OPEN ARCHITECTURE WITH ICEBERG (detailed flow) ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, RGBColor(0x0F, 0x17, 0x2A))
bar(slide, 0, 0, 13.333, 0.06, ORANGE)
tb(slide, 0.8, 0.15, 11, 0.4, "Open Architecture: Iceberg Tables + Best Service for Each Use Case", 22, ORANGE, True)
tb(slide, 0.8, 0.5, 11, 0.3, "zETL from Postgres/Oracle → Redshift Tables → Iceberg on S3. Every engine reads the same data.", 11, RGBColor(0x94, 0xA3, 0xB8))

# Column headers
hdrs = [("DATA SOURCES", 0.3, 1.1), ("ZETL", 1.6, 0.8), ("REDSHIFT + ICEBERG ON S3", 2.6, 2.2),
        ("DATA WAREHOUSE / ANALYTICS", 5.0, 3.0), ("AI LAYER", 8.3, 2.0), ("BI", 10.5, 1.5)]
for label, x, w in hdrs:
    tb(slide, x, 0.85, w, 0.2, label, 7, RGBColor(0x94, 0xA3, 0xB8), True, PP_ALIGN.CENTER)

# Sources
for i, src in enumerate(["Postgres\nAWS RDS", "Oracle\nOn-Prem", "Other\nSources"]):
    card(slide, 0.3, 1.2 + i * 0.65, 1.1, 0.55, RGBColor(0x64, 0x74, 0x8B), RGBColor(0x0F, 0x17, 0x2A))
    tb(slide, 0.35, 1.22 + i * 0.65, 1.0, 0.45, src, 9, WHITE, True, PP_ALIGN.CENTER)

tb(slide, 1.4, 1.7, 0.3, 0.3, "→", 16, ORANGE, True)

# zETL
card(slide, 1.6, 1.4, 0.8, 0.7, ORANGE, RGBColor(0x0F, 0x17, 0x2A))
tb(slide, 1.62, 1.45, 0.76, 0.25, "zETL", 11, ORANGE, True, PP_ALIGN.CENTER)
tb(slide, 1.62, 1.75, 0.76, 0.2, "Near real-time", 7, RGBColor(0xFC, 0xD3, 0x4D), False, PP_ALIGN.CENTER)

tb(slide, 2.4, 1.7, 0.3, 0.3, "→", 14, ORANGE, True)

# Redshift Tables + Iceberg
card(slide, 2.7, 1.15, 1.8, 0.55, BLUE, RGBColor(0x0F, 0x17, 0x2A))
tb(slide, 2.75, 1.18, 1.7, 0.25, "Redshift Tables", 10, LIGHT_BLUE, True, PP_ALIGN.CENTER)
tb(slide, 2.75, 1.43, 1.7, 0.2, "Live query + auto-publish", 7, RGBColor(0x93, 0xC5, 0xFD), False, PP_ALIGN.CENTER)
tb(slide, 3.5, 1.72, 0.3, 0.2, "↓", 12, LIGHT_BLUE, True, PP_ALIGN.CENTER)
card(slide, 2.7, 1.95, 1.8, 0.6, LIGHT_BLUE, RGBColor(0x0F, 0x17, 0x2A))
tb(slide, 2.75, 1.98, 1.7, 0.25, "🧊 Iceberg Tables on S3", 10, LIGHT_BLUE, True, PP_ALIGN.CENTER)
tb(slide, 2.75, 2.25, 1.7, 0.2, "Open format · Any engine reads", 7, RGBColor(0x93, 0xC5, 0xFD), False, PP_ALIGN.CENTER)

tb(slide, 4.55, 1.7, 0.3, 0.3, "→", 14, RGBColor(0x16, 0xA3, 0x4A), True)

# Data Warehouse / Analytics column
dw_items = [("Databricks", "Spark / ML", PURPLE, LIGHT_PURPLE),
            ("Palantir", "Foundry / AIP", PALANTIR_GREEN, RGBColor(0x4A, 0xDE, 0x80)),
            ("Redshift", "SQL Analytics", ORANGE, ORANGE),
            ("Snowflake", "If needed", RGBColor(0x64, 0x74, 0x8B), WHITE)]
for i, (name, role, border, tc) in enumerate(dw_items[:2]):
    x = 5.0 + i * 1.5
    card(slide, x, 1.15, 1.3, 0.55, border, RGBColor(0x0F, 0x17, 0x2A))
    tb(slide, x + 0.05, 1.18, 1.2, 0.25, name, 9, tc, True, PP_ALIGN.CENTER)
    tb(slide, x + 0.05, 1.43, 1.2, 0.2, role, 7, RGBColor(0x94, 0xA3, 0xB8), False, PP_ALIGN.CENTER)
for i, (name, role, border, tc) in enumerate(dw_items[2:]):
    x = 5.0 + i * 1.5
    card(slide, x, 1.8, 1.3, 0.55, border, RGBColor(0x0F, 0x17, 0x2A))
    tb(slide, x + 0.05, 1.83, 1.2, 0.25, name, 9, tc, True, PP_ALIGN.CENTER)
    tb(slide, x + 0.05, 2.08, 1.2, 0.2, role, 7, RGBColor(0x94, 0xA3, 0xB8), False, PP_ALIGN.CENTER)

card(slide, 5.0, 2.45, 2.8, 0.45, ORANGE, RGBColor(0x0F, 0x17, 0x2A))
tb(slide, 5.05, 2.48, 2.7, 0.3, "SageMaker Unified Studio", 9, ORANGE, True, PP_ALIGN.CENTER)

tb(slide, 7.9, 1.7, 0.3, 0.3, "→", 14, RGBColor(0x16, 0xA3, 0x4A), True)

# AI Layer
ai_items = [("Bedrock", "GenAI / RAG", ORANGE), ("SageMaker", "Custom ML", ORANGE),
            ("Palantir AIP", "Ontology AI", PALANTIR_GREEN), ("Databricks AI", "If best fit", PURPLE)]
for i, (name, role, color) in enumerate(ai_items):
    card(slide, 8.3, 1.1 + i * 0.6, 1.8, 0.5, color, RGBColor(0x0F, 0x17, 0x2A))
    tb(slide, 8.35, 1.13 + i * 0.6, 1.7, 0.2, name, 9, color, True, PP_ALIGN.CENTER)
    tb(slide, 8.35, 1.33 + i * 0.6, 1.7, 0.2, role, 7, RGBColor(0x94, 0xA3, 0xB8), False, PP_ALIGN.CENTER)

tb(slide, 10.2, 1.7, 0.3, 0.3, "→", 14, RGBColor(0x16, 0xA3, 0x4A), True)

# BI column
bi_items = ["QuickSight", "Tableau", "Power BI", "MicroStrategy"]
for i, name in enumerate(bi_items):
    color = ORANGE if i == 0 else RGBColor(0x64, 0x74, 0x8B)
    card(slide, 10.5, 1.1 + i * 0.6, 1.5, 0.5, color, RGBColor(0x0F, 0x17, 0x2A))
    tb(slide, 10.55, 1.2 + i * 0.6, 1.4, 0.25, name, 10, WHITE if i > 0 else ORANGE, True, PP_ALIGN.CENTER)

# Bottom benefit cards
flow_bens = [
    ("✓ Near Real-Time with zETL", "Postgres/Oracle connectors create Iceberg tables via zero-ETL through Redshift. Near real-time data availability.", 3.5),
    ("✓ Best AI Tool for Each Use Case", "Databricks ML for some. Bedrock for GenAI. Palantir AIP for ontology AI. SageMaker for custom models.", 3.8),
    ("✓ Open Governance + BI Freedom", "Lake Formation for governance. Any BI tool connects to Iceberg — QuickSight, Tableau, Power BI, MicroStrategy.", 3.8),
]
for i, (title, body, w) in enumerate(flow_bens):
    x = 0.3 + i * 4.1
    card(slide, x, 3.3, w, 1.3, RGBColor(0x86, 0xEF, 0xAC), RGBColor(0x0F, 0x17, 0x2A))
    tb(slide, x + 0.1, 3.35, w - 0.2, 0.25, title, 9, RGBColor(0x86, 0xEF, 0xAC), True)
    tb(slide, x + 0.1, 3.65, w - 0.2, 0.8, body, 8, RGBColor(0xD1, 0xFA, 0xE5))

# Key principle
card(slide, 0.3, 4.85, 12.7, 0.6, ORANGE, RGBColor(0x0F, 0x17, 0x2A))
tb(slide, 0.5, 4.9, 12.3, 0.4, "Key Principle: Iceberg tables on S3 make your data open to every engine. Databricks, Palantir, Redshift, Snowflake, SageMaker — they all read the same tables. May the best service be aligned to the appropriate use case.", 10, RGBColor(0xFC, 0xD3, 0x4D), False, PP_ALIGN.CENTER)

# ==================== SLIDE 13: AWS SOLVES EVERY PROBLEM ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, RGBColor(0x0A, 0x16, 0x28))
bar(slide, 0, 0, 13.333, 0.06, ORANGE)
tb(slide, 0.8, 0.2, 10, 0.5, "AWS Solves Every Problem in the Pipeline", 24, ORANGE, True)
tb(slide, 0.8, 0.6, 10, 0.3, "Same architecture flow — with an AWS benefit replacing each Databricks-only problem", 12, RGBColor(0x94, 0xA3, 0xB8))

# Column headers
headers = ["DATA SOURCES", "DATA MOVEMENT", "DATA LAKE / GOVERNANCE", "DATA WAREHOUSE", "AI", "BI"]
widths = [1.4, 1.8, 2.2, 2.2, 2.2, 1.5]
x_pos = 0.3
for hdr, w in zip(headers, widths):
    tb(slide, x_pos, 0.95, w, 0.2, hdr, 7, MUTED, True, PP_ALIGN.CENTER)
    x_pos += w + 0.1

# Flow boxes row
flow = [
    ("Postgres\nOracle\nOther", 1.4, RGBColor(0x64, 0x74, 0x8B)),
    ("Redshift zETL\nZero-ETL connectors", 1.8, ORANGE),
    ("🧊 Iceberg on S3\nGlue Data Catalog", 2.2, LIGHT_BLUE),
    ("Redshift Serverless\n+ OpenSearch\nAthena · EMR", 2.2, ORANGE),
    ("Amazon Bedrock\n+ SageMaker\nNova · Titan · 100+ models", 2.2, ORANGE),
    ("QuickSight\n+ Tableau / PBI", 1.5, ORANGE),
]
x_pos = 0.3
for text, w, color in flow:
    card(slide, x_pos, 1.2, w, 1.1, color, RGBColor(0x0A, 0x16, 0x28))
    tb(slide, x_pos + 0.05, 1.25, w - 0.1, 0.9, text, 9, color, True, PP_ALIGN.CENTER)
    x_pos += w + 0.1

# Green benefit boxes
benefits_row = [
    ("✓ Eliminate Glue/Informatica", "No ETL tool to buy. Redshift zETL auto-replicates in near real-time.", 1.8),
    ("✓ Open Governance", "Lake Formation + Glue Catalog. IAM-integrated. No vendor catalog lock-in.", 2.2),
    ("✓ Serverless, Scales, No Load Cost", "Redshift Serverless pay-per-query. Queries Iceberg on S3 directly.", 2.2),
    ("✓ 100+ Models, Best Price", "Bedrock: 100+ foundation models. Native KB with Redshift + OpenSearch.", 2.2),
    ("✓ Native + Any BI Tool", "QuickSight embedded, per-session. Plus Tableau, PBI connect natively.", 1.5),
]
x_pos = 1.7  # skip data sources column
for title, body, w in benefits_row:
    card(slide, x_pos, 2.6, w, 1.4, GREEN, RGBColor(0x0A, 0x16, 0x28))
    tb(slide, x_pos + 0.05, 2.65, w - 0.1, 0.25, title, 9, GREEN, True)
    tb(slide, x_pos + 0.05, 2.95, w - 0.1, 0.8, body, 8, RGBColor(0xD1, 0xFA, 0xE5))
    x_pos += w + 0.1

# "Was" problem boxes
was_items = [
    ("Was: Glue/Informatica", "Expensive licenses. Batch only.", 1.8),
    ("Was: Databricks Catalog", "Vendor lock-in for governance.", 2.2),
    ("Was: Load into Databricks", "Costly. Scaling challenges.", 2.2),
    ("Was: Databricks AI Only", "Limited models. Higher prices.", 2.2),
    ("Was: DB Dashboards", "Basic. Not embeddable.", 1.5),
]
x_pos = 1.7
for title, body, w in was_items:
    card(slide, x_pos, 4.3, w, 0.9, RGBColor(0xFC, 0xA5, 0xA5), RGBColor(0x0A, 0x16, 0x28))
    tb(slide, x_pos + 0.05, 4.35, w - 0.1, 0.2, title, 8, RGBColor(0xFC, 0xA5, 0xA5), True)
    tb(slide, x_pos + 0.05, 4.6, w - 0.1, 0.5, body, 7, RGBColor(0xFE, 0xCA, 0xCA))
    x_pos += w + 0.1

# Bottom summary
bar(slide, 0.3, 5.5, 12.7, 0.6, RGBColor(0x15, 0x25, 0x38))
tb(slide, 0.5, 5.55, 12.3, 0.4, "Every problem has an AWS-native answer. Eliminate ETL tools, remove vendor lock-in, get serverless scale, access 100+ AI models, and connect any BI tool — all on one bill, one IAM, one VPC.", 11, RGBColor(0xFC, 0xD3, 0x4D), False, PP_ALIGN.CENTER)

# ==================== SLIDE 14: SWIM LANE FLOW ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, RGBColor(0x0F, 0x17, 0x2A))
bar(slide, 0, 0, 13.333, 0.06, ORANGE)
tb(slide, 0.8, 0.15, 11, 0.4, "Open Architecture: Iceberg on S3 → Best Service for Each Use Case", 22, ORANGE, True)
tb(slide, 0.8, 0.5, 11, 0.3, "zETL from Postgres/Oracle → Redshift → Iceberg Tables on S3. Three paths fan out. Each vendor stays in its lane.", 11, RGBColor(0x94, 0xA3, 0xB8))

# Left: Sources → zETL → Iceberg
for i, src in enumerate(["Postgres", "Oracle", "Other"]):
    card(slide, 0.3, 1.3 + i * 0.55, 1.0, 0.45, RGBColor(0x64, 0x74, 0x8B), RGBColor(0x0F, 0x17, 0x2A))
    tb(slide, 0.35, 1.35 + i * 0.55, 0.9, 0.25, src, 9, WHITE, True, PP_ALIGN.CENTER)

tb(slide, 1.3, 1.7, 0.3, 0.3, "→", 14, ORANGE, True)

card(slide, 1.6, 1.4, 0.8, 0.7, ORANGE, RGBColor(0x0F, 0x17, 0x2A))
tb(slide, 1.62, 1.45, 0.76, 0.25, "zETL", 11, ORANGE, True, PP_ALIGN.CENTER)
tb(slide, 1.62, 1.75, 0.76, 0.2, "Real-time", 7, RGBColor(0xFC, 0xD3, 0x4D), False, PP_ALIGN.CENTER)

tb(slide, 2.4, 1.7, 0.3, 0.3, "→", 14, LIGHT_BLUE, True)

card(slide, 2.7, 1.2, 1.4, 0.45, BLUE, RGBColor(0x0F, 0x17, 0x2A))
tb(slide, 2.75, 1.23, 1.3, 0.25, "Redshift Tables", 9, LIGHT_BLUE, True, PP_ALIGN.CENTER)
tb(slide, 3.3, 1.68, 0.3, 0.2, "↓", 12, LIGHT_BLUE, True, PP_ALIGN.CENTER)
card(slide, 2.7, 1.9, 1.4, 0.55, LIGHT_BLUE, RGBColor(0x0F, 0x17, 0x2A))
tb(slide, 2.75, 1.93, 1.3, 0.25, "🧊 Iceberg on S3", 9, LIGHT_BLUE, True, PP_ALIGN.CENTER)
tb(slide, 2.75, 2.18, 1.3, 0.2, "Open · Any engine", 7, RGBColor(0x93, 0xC5, 0xFD), False, PP_ALIGN.CENTER)

# Lane labels
tb(slide, 4.9, 2.7, 2.0, 0.2, "ANALYTICS / WAREHOUSE", 7, RGBColor(0x94, 0xA3, 0xB8), True, PP_ALIGN.CENTER)
tb(slide, 7.2, 2.7, 2.0, 0.2, "AI / ML", 7, RGBColor(0x94, 0xA3, 0xB8), True, PP_ALIGN.CENTER)
tb(slide, 9.5, 2.7, 1.5, 0.2, "BI", 7, RGBColor(0x94, 0xA3, 0xB8), True, PP_ALIGN.CENTER)

# Databricks lane
card(slide, 4.5, 3.0, 7.0, 0.85, PURPLE, RGBColor(0x0F, 0x17, 0x2A))
tb(slide, 4.6, 3.05, 0.3, 0.3, "→", 14, PURPLE, True)
card(slide, 5.0, 3.1, 2.0, 0.6, PURPLE, RGBColor(0x1E, 0x1B, 0x2E))
tb(slide, 5.05, 3.13, 1.9, 0.2, "Databricks", 10, LIGHT_PURPLE, True, PP_ALIGN.CENTER)
tb(slide, 5.05, 3.35, 1.9, 0.2, "Spark ETL · Delta", 8, RGBColor(0xC4, 0xB5, 0xFD), False, PP_ALIGN.CENTER)
tb(slide, 7.1, 3.2, 0.3, 0.2, "→", 12, PURPLE, True)
card(slide, 7.4, 3.1, 1.8, 0.6, PURPLE, RGBColor(0x1E, 0x1B, 0x2E))
tb(slide, 7.45, 3.13, 1.7, 0.2, "Databricks AI", 10, LIGHT_PURPLE, True, PP_ALIGN.CENTER)
tb(slide, 7.45, 3.35, 1.7, 0.2, "MLflow · Mosaic", 8, RGBColor(0xC4, 0xB5, 0xFD), False, PP_ALIGN.CENTER)
tb(slide, 9.3, 3.2, 0.3, 0.2, "→", 12, PURPLE, True)
card(slide, 9.6, 3.1, 1.3, 0.6, PURPLE, RGBColor(0x1E, 0x1B, 0x2E))
tb(slide, 9.65, 3.2, 1.2, 0.25, "DB SQL Dash", 9, LIGHT_PURPLE, True, PP_ALIGN.CENTER)

# Palantir lane
card(slide, 4.5, 4.0, 7.0, 0.85, PALANTIR_GREEN, RGBColor(0x0F, 0x17, 0x2A))
tb(slide, 4.6, 4.05, 0.3, 0.3, "→", 14, PALANTIR_GREEN, True)
card(slide, 5.0, 4.1, 2.0, 0.6, PALANTIR_GREEN, RGBColor(0x0F, 0x23, 0x18))
tb(slide, 5.05, 4.13, 1.9, 0.2, "Palantir Foundry", 10, RGBColor(0x4A, 0xDE, 0x80), True, PP_ALIGN.CENTER)
tb(slide, 5.05, 4.35, 1.9, 0.2, "Ontology · Pipelines", 8, LIGHT_GREEN, False, PP_ALIGN.CENTER)
tb(slide, 7.1, 4.2, 0.3, 0.2, "→", 12, PALANTIR_GREEN, True)
card(slide, 7.4, 4.1, 1.8, 0.6, PALANTIR_GREEN, RGBColor(0x0F, 0x23, 0x18))
tb(slide, 7.45, 4.13, 1.7, 0.2, "Palantir AIP", 10, RGBColor(0x4A, 0xDE, 0x80), True, PP_ALIGN.CENTER)
tb(slide, 7.45, 4.35, 1.7, 0.2, "Ontology-driven AI", 8, LIGHT_GREEN, False, PP_ALIGN.CENTER)
tb(slide, 9.3, 4.2, 0.3, 0.2, "→", 12, PALANTIR_GREEN, True)
card(slide, 9.6, 4.1, 1.3, 0.6, PALANTIR_GREEN, RGBColor(0x0F, 0x23, 0x18))
tb(slide, 9.65, 4.2, 1.2, 0.25, "Foundry Apps", 9, RGBColor(0x4A, 0xDE, 0x80), True, PP_ALIGN.CENTER)

# AWS lane
card(slide, 4.5, 5.0, 7.0, 0.85, ORANGE, RGBColor(0x0F, 0x17, 0x2A))
tb(slide, 4.6, 5.05, 0.3, 0.3, "→", 14, ORANGE, True)
card(slide, 5.0, 5.1, 1.0, 0.6, ORANGE, RGBColor(0x0F, 0x17, 0x2A))
tb(slide, 5.02, 5.13, 0.96, 0.2, "Redshift", 9, ORANGE, True, PP_ALIGN.CENTER)
tb(slide, 5.02, 5.35, 0.96, 0.2, "SQL", 7, RGBColor(0xFC, 0xD3, 0x4D), False, PP_ALIGN.CENTER)
card(slide, 6.1, 5.1, 1.0, 0.6, ORANGE, RGBColor(0x0F, 0x17, 0x2A))
tb(slide, 6.12, 5.13, 0.96, 0.2, "OpenSearch", 9, ORANGE, True, PP_ALIGN.CENTER)
tb(slide, 6.12, 5.35, 0.96, 0.2, "Search · RAG", 7, RGBColor(0xFC, 0xD3, 0x4D), False, PP_ALIGN.CENTER)
tb(slide, 7.2, 5.2, 0.3, 0.2, "→", 12, ORANGE, True)
card(slide, 7.5, 5.1, 0.9, 0.6, ORANGE, RGBColor(0x0F, 0x17, 0x2A))
tb(slide, 7.52, 5.13, 0.86, 0.2, "Bedrock", 9, ORANGE, True, PP_ALIGN.CENTER)
tb(slide, 7.52, 5.35, 0.86, 0.2, "GenAI", 7, RGBColor(0xFC, 0xD3, 0x4D), False, PP_ALIGN.CENTER)
card(slide, 8.5, 5.1, 0.9, 0.6, ORANGE, RGBColor(0x0F, 0x17, 0x2A))
tb(slide, 8.52, 5.13, 0.86, 0.2, "SageMaker", 9, ORANGE, True, PP_ALIGN.CENTER)
tb(slide, 8.52, 5.35, 0.86, 0.2, "Custom ML", 7, RGBColor(0xFC, 0xD3, 0x4D), False, PP_ALIGN.CENTER)
tb(slide, 9.5, 5.2, 0.3, 0.2, "→", 12, ORANGE, True)
card(slide, 9.8, 5.1, 0.7, 0.28, ORANGE, RGBColor(0x0F, 0x17, 0x2A))
tb(slide, 9.82, 5.12, 0.66, 0.2, "QuickSight", 8, ORANGE, True, PP_ALIGN.CENTER)
card(slide, 9.8, 5.42, 0.7, 0.28, RGBColor(0x64, 0x74, 0x8B), RGBColor(0x0F, 0x17, 0x2A))
tb(slide, 9.82, 5.44, 0.66, 0.2, "Tableau/PBI", 8, WHITE, False, PP_ALIGN.CENTER)

# Bottom callouts
lane_bens = [
    ("Databricks Path", "Best for large-scale Spark ETL, data engineering, and ML workloads.", LIGHT_PURPLE),
    ("Palantir Path", "Best for ontology-driven operations, mission-critical decision-making.", RGBColor(0x4A, 0xDE, 0x80)),
    ("AWS Analytics Path", "Best for SQL analytics, GenAI/RAG, semantic search, BI, custom ML. Serverless.", ORANGE),
]
for i, (title, body, color) in enumerate(lane_bens):
    x = 0.3 + i * 4.1
    card(slide, x, 6.0, 3.8, 0.7, color, RGBColor(0x0F, 0x17, 0x2A))
    tb(slide, x + 0.1, 6.05, 3.6, 0.2, title, 10, color, True)
    tb(slide, x + 0.1, 6.3, 3.6, 0.3, body, 8, LIGHT_GREEN if i == 1 else (RGBColor(0xFC, 0xD3, 0x4D) if i == 2 else RGBColor(0xC4, 0xB5, 0xFD)))

card(slide, 0.3, 6.85, 12.7, 0.5, LIGHT_BLUE, RGBColor(0x0F, 0x17, 0x2A))
tb(slide, 0.5, 6.9, 12.3, 0.35, "🧊 Iceberg is the key: One copy of data on S3. Three vendors read it natively. May the best service be aligned to the appropriate use case.", 10, RGBColor(0x93, 0xC5, 0xFD), False, PP_ALIGN.CENTER)

# ==================== SLIDE 15: SECOND USE CASE MATRIX ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, RGBColor(0xFF, 0xFF, 0xFF))
bar(slide, 0, 0, 13.333, 0.65, DARK_BG)
tb(slide, 0.8, 0.1, 9, 0.45, "Best Service for Each Use Case — Let the Data Decide", 22, WHITE, True)
tb(slide, 10.5, 0.15, 2.5, 0.4, "USE CASE ALIGNMENT", 10, DARK_BG, True, PP_ALIGN.CENTER)

hy = 0.85
tb(slide, 0.5, hy, 2.5, 0.3, "USE CASE", 9, RGBColor(0x47, 0x55, 0x69), True)
tb(slide, 3.2, hy, 2.3, 0.3, "DATABRICKS", 9, PURPLE, True)
tb(slide, 5.7, hy, 2.5, 0.3, "AWS ANALYTICS", 9, ORANGE, True)
tb(slide, 8.4, hy, 2.3, 0.3, "PALANTIR", 9, PALANTIR_GREEN, True)
tb(slide, 10.9, hy, 2.0, 0.3, "BEST FIT", 9, RGBColor(0x1E, 0x40, 0xAF), True)

matrix2 = [
    ("Large-Scale Spark ETL", "✓ Strong — native Spark", "EMR / Glue Spark jobs", "Not primary", "Databricks", PURPLE),
    ("SQL Analytics & Warehousing", "Databricks SQL (costly)", "✓ Redshift Serverless — native Bedrock, zETL", "Not primary", "Redshift", ORANGE),
    ("Generative AI / RAG", "Limited models, no native RAG", "✓ Bedrock + OpenSearch + Redshift — native KB", "AIP (ontology-based)", "Bedrock", ORANGE),
    ("Full-Text & Semantic Search", "No native search engine", "✓ OpenSearch — full-text + vector + serverless", "Not primary", "OpenSearch", ORANGE),
    ("Business Intelligence", "Basic SQL dashboards only", "✓ QuickSight — embedded, per-session, Q", "Foundry apps", "QuickSight", ORANGE),
    ("Custom ML Model Training", "✓ MLflow notebooks", "✓ SageMaker — full MLOps, Bedrock fine-tuning", "Foundry ML", "Both", ORANGE),
    ("Operational Decision-Making", "Not designed for this", "Step Functions + Bedrock Agents", "✓ Foundry / AIP — purpose-built", "Palantir", PALANTIR_GREEN),
    ("Near Real-Time Ingestion", "Requires custom streaming", "✓ Redshift zETL — auto-replication", "Requires connectors", "Redshift zETL", ORANGE),
]

for i, (uc, db, aws, pal, best, best_color) in enumerate(matrix2):
    y = 1.2 + i * 0.6
    tb(slide, 0.5, y, 2.5, 0.3, uc, 10, RGBColor(0x1E, 0x29, 0x3B), True)
    tb(slide, 3.2, y, 2.3, 0.3, db, 9, RGBColor(0x47, 0x55, 0x69))
    tb(slide, 5.7, y, 2.5, 0.3, aws, 9, RGBColor(0x47, 0x55, 0x69))
    tb(slide, 8.4, y, 2.3, 0.3, pal, 9, RGBColor(0x47, 0x55, 0x69))
    tb(slide, 10.9, y, 2.0, 0.3, best, 10, best_color, True)

bar(slide, 0, 6.3, 13.333, 0.7, DARK_BG)
tb(slide, 0.8, 6.4, 12, 0.4, "🧊 Iceberg makes this possible: All services read the same tables on S3. Align the best service to each use case — change your mind later without migrating data.", 12, WHITE)

# ==================== SLIDE 16: SUMMARY / CTA ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
bar(slide, 0, 0, 13.333, 0.06, ORANGE)
tb(slide, 0.8, 0.3, 11, 0.5, "The AWS Analytics Advantage: Why Add Native Services", 26, ORANGE, True)
tb(slide, 0.8, 0.8, 11, 0.3, "Six reasons to complement Databricks with AWS-native analytics on your S3 data lake", 13, GRAY)

summary_cards = [
    ("💰", "Cost at Scale", "S3 at $0.023/GB vs importing into Databricks. Serverless everything — pay only when you query.", "Up to 90% storage cost reduction"),
    ("🤖", "Native AI Integration", "Redshift, OpenSearch, S3 are all native Bedrock data sources. One AI layer, zero ETL.", "Zero custom integration code"),
    ("🔓", "No Vendor Lock-In", "Data stays in S3 in open formats. Any engine reads it. Swap services without migration.", "Open formats, open architecture"),
    ("📈", "Best Tool for Each Job", "Redshift for SQL. OpenSearch for search. QuickSight for BI. SageMaker for ML. Bedrock for GenAI.", "5 services, one data lake"),
    ("🔒", "Unified Security", "One IAM, one VPC, one Lake Formation governance. Data never leaves your AWS account.", "Single identity, single governance"),
    ("🤝", "Co-Existence, Not Replace", "Keep Databricks for Spark. Add AWS-native for analytics, search, BI, AI. Better together.", "Databricks + AWS = better together"),
]

for i, (icon, title, body, stat) in enumerate(summary_cards):
    col = i % 3
    row = i // 3
    x = 0.8 + col * 4.0
    y = 1.3 + row * 2.6
    card(slide, x, y, 3.7, 2.3, ORANGE)
    tb(slide, x + 0.15, y + 0.1, 3.4, 0.3, icon, 22, WHITE, False, PP_ALIGN.LEFT)
    tb(slide, x + 0.15, y + 0.45, 3.4, 0.3, title, 14, ORANGE, True)
    tb(slide, x + 0.15, y + 0.8, 3.4, 0.8, body, 10, WHITE)
    tb(slide, x + 0.15, y + 1.7, 3.4, 0.3, stat, 10, ORANGE, True)

# CTA bar
bar(slide, 0, 6.6, 13.333, 0.9, ORANGE)
tb(slide, 0.8, 6.7, 8, 0.35, "Ready to Unlock the Full AWS Analytics Suite?", 18, DARK_BG, True)
tb(slide, 0.8, 7.05, 8, 0.3, "Start with a POC: S3 data lake → Redshift + OpenSearch + Bedrock. See results in days.", 12, DARK_BG)
tb(slide, 10.0, 6.8, 3.0, 0.5, "Let's Build a POC →", 14, ORANGE, True, PP_ALIGN.CENTER)

# ==================== SAVE ====================
output_path = "AWS-Analytics-vs-Databricks-all-slides.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
