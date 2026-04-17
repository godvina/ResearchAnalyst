"""Convert investigative-intelligence-gtm.html to PowerPoint."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ORANGE = RGBColor(0xFF, 0x99, 0x00)
LIGHT_ORANGE = RGBColor(0xFF, 0xB8, 0x4D)
WHITE = RGBColor(0xE2, 0xE8, 0xF0)
GRAY = RGBColor(0xA0, 0xAE, 0xC0)
DARK_GRAY = RGBColor(0x71, 0x80, 0x96)
GREEN = RGBColor(0x48, 0xBB, 0x78)
RED = RGBColor(0xFC, 0x81, 0x81)
YELLOW = RGBColor(0xF6, 0xE0, 0x5E)
BLUE = RGBColor(0x63, 0xB3, 0xED)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox

def add_paragraph(text_frame, text, font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, space_before=0):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    if space_before:
        p.space_before = Pt(space_before)
    return p

def add_card(slide, left, top, width, height, border_color=ORANGE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x3A)
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    return shape

# ==================== SLIDE 1: TITLE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

add_text_box(slide, 1.2, 1.0, 10, 1.0, "Investigative Intelligence:", 40, WHITE, False)
add_text_box(slide, 1.2, 1.8, 10, 1.0, "AI for Federal Law Enforcement", 40, ORANGE, True)
add_text_box(slide, 1.2, 3.0, 8, 0.5, "AWS Investigative Intelligence Platform 2026", 16, GRAY)
add_text_box(slide, 1.2, 3.8, 10, 0.5, "Department of Justice — Antitrust Division", 18, WHITE, True)
add_text_box(slide, 1.2, 4.3, 10, 0.8, "How generative AI, knowledge graphs, and intelligent search transform\n500TB of case evidence into actionable intelligence", 14, DARK_GRAY)

# Divider line
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(5.3), Inches(4), Pt(1))
shape.fill.solid()
shape.fill.fore_color.rgb = ORANGE
shape.line.fill.background()

add_text_box(slide, 1.2, 5.5, 4, 0.4, "David Eyre", 14, ORANGE, True)
add_text_box(slide, 1.2, 5.9, 4, 0.4, "Emerging Tech Solutions — Amazon Web Services", 12, GRAY)

# ==================== SLIDE 2: INDUSTRY CHALLENGE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.3, 11, 0.7, "The Two Challenges Defining Federal Investigations", 28, ORANGE, True)

# Challenge 1
card1 = add_card(slide, 0.8, 1.2, 5.6, 4.8, RED)
txBox = add_text_box(slide, 1.1, 1.4, 5.0, 0.5, "🏛️ Challenge 1: Do More With Less", 18, RED, True)
tf = txBox.text_frame
add_paragraph(tf, "", 6)
add_paragraph(tf, "Federal workforce reductions mean fewer investigators handling the same — or growing — caseload. Institutional knowledge walks out the door.", 12, WHITE, space_before=4)
add_paragraph(tf, "", 6)
add_paragraph(tf, "THE COMPOUNDING EFFECT", 10, RED, True, space_before=8)
add_paragraph(tf, "📉 Workforce reductions with no new backfills", 11, WHITE, space_before=4)
add_paragraph(tf, "📚 Thousands of cases queued", 11, WHITE, space_before=2)
add_paragraph(tf, "⏱️ Investigators spend 60–80% of time on manual review", 11, WHITE, space_before=2)
add_paragraph(tf, "📅 Cases take 18–24 months average", 11, WHITE, space_before=2)
add_paragraph(tf, "📝 Brief preparation is manual and error-prone", 11, WHITE, space_before=2)

# 500TB stat
add_text_box(slide, 1.5, 6.2, 4.0, 0.8, "500TB+", 36, ORANGE, True, PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 6.7, 4.0, 0.4, "of case evidence per major federal investigation", 11, GRAY, False, PP_ALIGN.CENTER)

# Challenge 2
card2 = add_card(slide, 6.8, 1.2, 5.6, 5.5, YELLOW)
txBox2 = add_text_box(slide, 7.1, 1.4, 5.0, 0.5, "🔍 Challenge 2: Find the Needle You Don't Know Exists", 18, YELLOW, True)
tf2 = txBox2.text_frame
add_paragraph(tf2, "", 6)
add_paragraph(tf2, "The hardest cases aren't about finding a known document — they're about discovering connections no one knew to look for.", 12, WHITE, space_before=4)
add_paragraph(tf2, "", 6)
add_paragraph(tf2, "WHAT'S MISSING TODAY", 10, YELLOW, True, space_before=8)
add_paragraph(tf2, "🚫 No AI-assisted pattern discovery across evidence", 11, WHITE, space_before=4)
add_paragraph(tf2, "🔒 Evidence silos across agencies", 11, WHITE, space_before=2)
add_paragraph(tf2, "🕸️ No knowledge graph linking entities across cases", 11, WHITE, space_before=2)
add_paragraph(tf2, "📦 Legacy tools find documents, not intelligence", 11, WHITE, space_before=2)
add_paragraph(tf2, "🎯 No prosecution readiness assessment until months in", 11, WHITE, space_before=2)

add_text_box(slide, 7.1, 6.2, 5.0, 0.5, "AI doesn't replace investigators — it gives each one the throughput of ten.", 13, GREEN, True, PP_ALIGN.CENTER)

# ==================== SLIDE 3: SOLUTION OVERVIEW ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.3, 11, 0.7, "Solution Overview — Two Personas", 28, ORANGE, True)

# Investigator card
add_card(slide, 0.8, 1.3, 4.5, 4.5, GREEN)
txBox = add_text_box(slide, 1.1, 1.5, 4.0, 0.5, "🔍 Investigator Intelligence", 18, GREEN, True)
tf = txBox.text_frame
for item in ["🤖 AI Briefing & Case Summary", "🕸️ Entity Network Discovery", "📋 Evidence Triage", "📖 Investigative Playbooks", "🗺️ Geospatial Analysis", "⏳ Timeline Intelligence"]:
    add_paragraph(tf, item, 13, WHITE, space_before=6)

# Shared center
add_text_box(slide, 5.6, 1.5, 2.0, 0.4, "SHARED", 10, DARK_GRAY, True, PP_ALIGN.CENTER)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(2.0), Pt(2), Inches(2.5))
shape.fill.solid()
shape.fill.fore_color.rgb = ORANGE
shape.line.fill.background()
add_text_box(slide, 5.6, 4.7, 2.0, 0.3, "Knowledge Graph", 11, GREEN, False, PP_ALIGN.CENTER)
add_text_box(slide, 5.6, 5.0, 2.0, 0.3, "Semantic Search", 11, GREEN, False, PP_ALIGN.CENTER)
add_text_box(slide, 5.6, 5.3, 2.0, 0.3, "Cross-Case Analysis", 11, GREEN, False, PP_ALIGN.CENTER)

# Prosecutor card
add_card(slide, 7.8, 1.3, 4.5, 4.5, YELLOW)
txBox = add_text_box(slide, 8.1, 1.5, 4.0, 0.5, "⚖️ Prosecutor Intelligence", 18, YELLOW, True)
tf = txBox.text_frame
for item in ["📊 Case Strength Scorecard", "✅ Element Assessment", "📚 Precedent Matching", "📄 Court Document Assembly", "🎯 Prosecution Funnel"]:
    add_paragraph(tf, item, 13, WHITE, space_before=6)

# Bottom banner
add_text_box(slide, 0.8, 6.3, 11.5, 0.5, "One platform, two workflows — from first lead to court-ready case", 15, GREEN, False, PP_ALIGN.CENTER)

# ==================== SLIDE 4: PROCESS FLOW ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.3, 11, 0.7, "Functional Process Flow", 28, ORANGE, True)
add_text_box(slide, 0.8, 0.9, 11, 0.4, "The investigation lifecycle — from raw evidence to prosecution readiness", 14, GRAY)

steps = [
    ("1", "INGEST", "Documents → S3\nTextract / Rekognition\nParse & Extract text\nOCR scanned docs"),
    ("2", "ENRICH", "Entity Extraction\nNeptune Graph Load\nOpenSearch Index\nBedrock Embeddings"),
    ("3", "ANALYZE", "AI Briefing\nPattern Discovery\nCross-Case Analysis\nLead Generation"),
    ("4", "INVESTIGATE", "Playbook Workflow\nEntity Tracking\nGraph Exploration\nTimeline Analysis"),
    ("5", "PROSECUTE", "Case Strength Score\nSubject Evaluation\nCourt Doc Assembly\nProsecution Readiness"),
]

for i, (num, title, items) in enumerate(steps):
    x = 0.8 + i * 2.4
    # Number circle
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.7), Inches(1.5), Inches(0.5), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ORANGE
    shape.line.fill.background()
    shape.text_frame.paragraphs[0].text = num
    shape.text_frame.paragraphs[0].font.size = Pt(16)
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].font.color.rgb = DARK_BG
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text_box(slide, x, 2.2, 2.2, 0.4, title, 14, LIGHT_ORANGE, True, PP_ALIGN.CENTER)
    add_text_box(slide, x, 2.7, 2.2, 1.5, items, 11, GRAY, False, PP_ALIGN.CENTER)

    # Arrow between steps
    if i < 4:
        add_text_box(slide, x + 2.1, 1.55, 0.4, 0.5, "→", 20, ORANGE, True, PP_ALIGN.CENTER)

# Bottom layer cards
layers = [
    ("📥", "S3 + Textract + Rekognition", "Data ingestion layer", GREEN),
    ("🧠", "Bedrock + Neptune + OpenSearch", "Intelligence layer", YELLOW),
    ("⚖️", "AI Scoring + Document Assembly", "Prosecution layer", BLUE),
]
for i, (icon, label, sub, color) in enumerate(layers):
    x = 0.8 + i * 4.0
    add_card(slide, x, 4.8, 3.6, 1.2, color)
    add_text_box(slide, x, 4.9, 3.6, 0.4, icon, 20, color, False, PP_ALIGN.CENTER)
    add_text_box(slide, x, 5.3, 3.6, 0.3, label, 11, color, True, PP_ALIGN.CENTER)
    add_text_box(slide, x, 5.6, 3.6, 0.3, sub, 10, DARK_GRAY, False, PP_ALIGN.CENTER)

# ==================== SLIDE 5: ARCHITECTURE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.3, 11, 0.7, "Architecture — AWS Lego Blocks", 28, ORANGE, True)

arch_layers = [
    ("Data Layer", [("Amazon S3", "Data Lake — $11,500/mo"), ("Amazon S3", "Frontend Hosting")]),
    ("Ingestion", [("AWS Step Functions", "Orchestration — $500/mo"), ("AWS Lambda", "Compute — $2,400/mo"), ("Amazon Textract", "Doc Parsing — $4,500/mo"), ("Amazon Rekognition", "Image Analysis — $2,800/mo")]),
    ("Intelligence", [("OpenSearch Serverless", "Vector + Full-text — $8,200/mo"), ("Amazon Neptune", "Knowledge Graph — $3,800/mo"), ("Amazon Bedrock", "Claude/Titan AI — $12,000/mo")]),
    ("API Layer", [("API Gateway", "REST API — $350/mo"), ("AWS Lambda", "API Compute")]),
    ("Security", [("AWS IAM", "Access Control"), ("Amazon Cognito", "Auth (future)"), ("AWS KMS", "Encryption — $100/mo")]),
    ("Infrastructure", [("AWS CDK", "IaC — $0"), ("CloudWatch", "Monitoring — $200/mo")]),
]

for row_i, (layer_name, services) in enumerate(arch_layers):
    y = 1.2 + row_i * 0.9
    label_color = YELLOW if layer_name == "Intelligence" else DARK_GRAY
    add_text_box(slide, 0.3, y, 1.8, 0.5, layer_name.upper(), 10, label_color, True, PP_ALIGN.RIGHT)

    for col_i, (svc, desc) in enumerate(services):
        x = 2.5 + col_i * 2.7
        box_color = LIGHT_ORANGE if layer_name == "Intelligence" else (BLUE if layer_name == "Security" else ORANGE)
        card = add_card(slide, x, y, 2.5, 0.7, box_color)
        add_text_box(slide, x + 0.1, y + 0.05, 2.3, 0.3, svc, 11, WHITE, True)
        add_text_box(slide, x + 0.1, y + 0.35, 2.3, 0.3, desc, 9, DARK_GRAY)

add_text_box(slide, 0.8, 6.8, 11.5, 0.4, "Estimated total: $46,350/mo for a 500TB deployment", 13, DARK_GRAY, False, PP_ALIGN.CENTER)

# ==================== SLIDE 6: INGESTION PIPELINE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.3, 11, 0.7, "Data Ingestion Pipeline — Today & Tomorrow", 28, ORANGE, True)

# Phase 1 label
add_text_box(slide, 0.8, 1.1, 6, 0.3, "PROVEN  Custom Pipeline — 30,000 Documents Loaded", 13, GREEN, True)
add_text_box(slide, 0.8, 1.4, 10, 0.3, "Phase 1 — Document Processing Pipeline (Step Functions)", 11, GREEN, True)

phase1_steps = [
    ("📁 S3 Upload", "Raw docs into data lake", GREEN),
    ("📄 Textract Parse", "OCR scanned docs", ORANGE),
    ("🧠 Entity Extraction", "People, orgs, locations", LIGHT_ORANGE),
    ("🔢 Embeddings", "Titan vectors for search", BLUE),
    ("🔍 OpenSearch Index", "Full-text + kNN", RGBColor(0x9F, 0x7A, 0xEA)),
    ("📸 Rekognition", "Face & label detection", RED),
]

for i, (title, desc, color) in enumerate(phase1_steps):
    x = 0.8 + i * 2.0
    add_card(slide, x, 1.8, 1.8, 0.9, color)
    add_text_box(slide, x + 0.05, 1.85, 1.7, 0.3, title, 10, color, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.05, 2.2, 1.7, 0.4, desc, 8, GRAY, False, PP_ALIGN.CENTER)
    if i < 5:
        add_text_box(slide, x + 1.7, 2.0, 0.3, 0.3, "→", 14, ORANGE, True)

# Phase 2
add_text_box(slide, 0.8, 3.0, 10, 0.3, "Phase 2 — Knowledge Graph Build (CSV Bulk Load → Neptune)", 11, RGBColor(0xF6, 0xAD, 0x55), True)

phase2_steps = [
    ("📊 Entity Index", "Canonical entity dedup"),
    ("📝 CSV Generation", "Nodes + edges files"),
    ("🕸️ Neptune Bulk Load", "CSV → graph"),
    ("✅ Intelligence Ready", "Full platform operational"),
]

for i, (title, desc) in enumerate(phase2_steps):
    x = 0.8 + i * 3.0
    color = GREEN if i == 3 else RGBColor(0xF6, 0xAD, 0x55)
    add_card(slide, x, 3.4, 2.8, 0.9, color)
    add_text_box(slide, x + 0.05, 3.45, 2.7, 0.3, title, 10, color, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.05, 3.8, 2.7, 0.4, desc, 8, GRAY, False, PP_ALIGN.CENTER)
    if i < 3:
        add_text_box(slide, x + 2.7, 3.6, 0.3, 0.3, "→", 14, RGBColor(0xF6, 0xAD, 0x55), True)

# Future managed service
add_text_box(slide, 0.8, 4.7, 6, 0.3, "FUTURE  Managed Batch Loader — AWS Service", 13, BLUE, True)

future_steps = ["⚙️ Config", "📥 Ingest", "🔄 Transform", "🔍 Index"]
for i, step in enumerate(future_steps):
    x = 0.8 + i * 3.0
    add_card(slide, x, 5.1, 2.8, 0.7, BLUE)
    add_text_box(slide, x + 0.1, 5.2, 2.6, 0.4, step, 12, BLUE, True, PP_ALIGN.CENTER)
    if i < 3:
        add_text_box(slide, x + 2.7, 5.2, 0.3, 0.3, "→", 14, BLUE, True)

# Alignment path
add_text_box(slide, 0.8, 6.2, 11.5, 0.8, "🔄 ALIGNMENT PATH: Our custom pipeline code aligns with the future managed Batch Loader. ATR requirements from this pilot feed directly into the service team roadmap.", 11, GRAY, False, PP_ALIGN.CENTER)

# ==================== SLIDE 7: DOJ PILOT ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.3, 11, 0.7, "DOJ Antitrust Pilot — 500TB", 28, ORANGE, True)

# Pilot Scope card
add_card(slide, 0.8, 1.2, 5.6, 2.8, GREEN)
add_text_box(slide, 1.1, 1.4, 5.0, 0.4, "Pilot Scope", 16, LIGHT_ORANGE, True)
scope_items = [
    ("Data Volume", "500 TB of case evidence"),
    ("Division", "DOJ Antitrust Division"),
    ("Timeline", "12-week pilot"),
    ("Budget", "~$200K (ProServe + AWS services)"),
]
for i, (label, value) in enumerate(scope_items):
    y = 1.9 + i * 0.45
    add_text_box(slide, 1.1, y, 1.8, 0.4, label, 12, DARK_GRAY)
    add_text_box(slide, 3.0, y, 3.0, 0.4, value, 12, WHITE, True)

# Beta customer card
add_card(slide, 0.8, 4.3, 5.6, 1.2, BLUE)
add_text_box(slide, 1.1, 4.5, 5.0, 0.3, "Beta Customer Benefit", 14, BLUE, True)
add_text_box(slide, 1.1, 4.9, 5.0, 0.5, "Requirements feed directly to AWS service team for OpenSearch Batch Loader. DOJ shapes the managed service roadmap.", 12, WHITE)

# Expected Outcomes
add_text_box(slide, 6.8, 1.2, 5.5, 0.4, "Expected Outcomes", 16, LIGHT_ORANGE, True)
add_text_box(slide, 7.0, 1.8, 2.0, 0.8, "70%", 40, GREEN, True)
add_text_box(slide, 9.0, 2.0, 3.5, 0.5, "reduction in manual\ndocument review time", 14, WHITE)

outcomes = [
    "• AI-discovered patterns humans miss",
    "• Cross-case intelligence across historical matters",
    "• Prosecution readiness scoring",
    "• Entity network discovery across 500TB",
]
txBox = add_text_box(slide, 7.0, 2.8, 5.0, 0.3, outcomes[0], 12, WHITE)
for item in outcomes[1:]:
    add_paragraph(txBox.text_frame, item, 12, WHITE, space_before=4)

# Pilot to production
add_card(slide, 6.8, 4.3, 5.6, 1.2, YELLOW)
add_text_box(slide, 7.1, 4.5, 5.0, 0.3, "PILOT → PRODUCTION PATH", 12, YELLOW, True, PP_ALIGN.CENTER)
add_text_box(slide, 7.1, 4.9, 5.0, 0.4, "12 weeks to prove value → expand to full division", 13, GRAY, False, PP_ALIGN.CENTER)

# ==================== SLIDE 8: ISV COEXISTENCE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.3, 11, 0.7, "ISV Coexistence — Intelligence Layer, Not Replacement", 28, ORANGE, True)

# Existing tools layer
add_card(slide, 1.5, 1.2, 10.3, 1.5, BLUE)
add_text_box(slide, 1.8, 1.3, 9.5, 0.3, "EXISTING eDISCOVERY & CASE MANAGEMENT TOOLS", 11, BLUE, True, PP_ALIGN.CENTER)
tools = "Relativity    •    Nuix    •    CasePoint    •    Concordance    •    DISCO"
add_text_box(slide, 1.8, 1.7, 9.5, 0.4, tools, 14, WHITE, False, PP_ALIGN.CENTER)
add_text_box(slide, 1.8, 2.1, 9.5, 0.3, "These tools find documents — they don't discover intelligence, patterns, or hidden connections", 10, DARK_GRAY, False, PP_ALIGN.CENTER)

add_text_box(slide, 6.2, 2.8, 1.0, 0.4, "↕️", 20, ORANGE, False, PP_ALIGN.CENTER)

# Our intelligence layer
add_card(slide, 1.5, 3.3, 10.3, 1.2, GREEN)
add_text_box(slide, 1.8, 3.4, 9.5, 0.3, "⚡ OUR INTELLIGENCE LAYER", 12, GREEN, True, PP_ALIGN.CENTER)
intel = "AI Briefing  •  Entity Graph  •  Cross-Case Analysis  •  Pattern Discovery  •  Prosecution Scoring"
add_text_box(slide, 1.8, 3.8, 9.5, 0.3, intel, 13, WHITE, False, PP_ALIGN.CENTER)
add_text_box(slide, 1.8, 4.1, 9.5, 0.3, "Enriches, not replaces. Data flows both ways.", 11, GRAY, False, PP_ALIGN.CENTER)

add_text_box(slide, 6.2, 4.6, 1.0, 0.4, "↕️", 20, ORANGE, False, PP_ALIGN.CENTER)

# AWS Foundation
add_card(slide, 1.5, 5.1, 10.3, 1.0, YELLOW)
add_text_box(slide, 1.8, 5.2, 9.5, 0.3, "AWS FOUNDATION SERVICES", 11, YELLOW, True, PP_ALIGN.CENTER)
aws_svc = "Bedrock  •  Neptune  •  OpenSearch  •  S3  •  Lambda"
add_text_box(slide, 1.8, 5.5, 9.5, 0.3, aws_svc, 13, WHITE, False, PP_ALIGN.CENTER)

# Bottom cards
add_card(slide, 0.8, 6.3, 5.5, 0.8, GREEN)
add_text_box(slide, 1.0, 6.35, 5.0, 0.3, "Key Message", 12, GREEN, True)
add_text_box(slide, 1.0, 6.6, 5.0, 0.3, "\"We add AI intelligence ON TOP of existing case management tools\"", 11, WHITE)

add_card(slide, 6.8, 6.3, 5.5, 0.8, YELLOW)
add_text_box(slide, 7.0, 6.35, 5.0, 0.3, "ISV Benefit", 12, YELLOW, True)
add_text_box(slide, 7.0, 6.6, 5.0, 0.3, "Their customers consume MORE AWS services, not fewer.", 11, WHITE)

# ==================== SLIDE 9: RISKS ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.3, 11, 0.7, "Risks and Mitigation", 28, ORANGE, True)

risks = [
    ("RISK 1", "Service teams know 1 of 3 services. Neptune SA doesn't know OpenSearch. Bedrock SA doesn't know Neptune.",
     "Cross-train SAs across Neptune / OpenSearch / Bedrock. Sell as a solution, not individual services."),
    ("RISK 2", "Partners / ProServe lack solution expertise across all three services.",
     "Solution blueprints + assisted deployment. Reduce dependency on deep multi-service expertise."),
    ("RISK 3", "OpenSearch Batch Loader not in GovCloud / FedRAMP yet.",
     "Service team builds interim pipeline code. DOJ becomes beta customer with requirements flowing to service team."),
]

for i, (risk_label, risk_text, mitigation_text) in enumerate(risks):
    x = 0.8 + i * 4.0
    add_card(slide, x, 1.2, 3.7, 5.0, RGBColor(0x33, 0x33, 0x50))

    add_text_box(slide, x + 0.2, 1.4, 3.3, 0.3, risk_label, 10, RED, True)
    add_text_box(slide, x + 0.2, 1.8, 3.3, 1.5, risk_text, 12, WHITE)

    add_text_box(slide, x + 0.2, 3.5, 3.3, 0.3, "MITIGATION", 10, ORANGE, True)
    add_text_box(slide, x + 0.2, 3.9, 3.3, 1.5, mitigation_text, 12, WHITE)

add_text_box(slide, 0.8, 6.5, 11.5, 0.5, "Every risk has a clear mitigation path. The biggest risk is doing nothing — competitors are building this.", 14, GREEN, False, PP_ALIGN.CENTER)

# ==================== SLIDE 10: GTM TRANSFORMATION ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.3, 11, 0.7, "GTM Transformation — Service Focus → Solution Focus", 28, ORANGE, True)

# Today: Siloed
add_card(slide, 0.8, 1.3, 5.0, 4.0, RED)
add_text_box(slide, 1.1, 1.5, 4.5, 0.4, "Today: Siloed", 18, RED, True)

today_items = [
    "Seller + SA → Neptune",
    "Seller + SA → OpenSearch",
    "Seller + SA → Bedrock",
]
for i, item in enumerate(today_items):
    add_text_box(slide, 1.3, 2.2 + i * 0.6, 4.0, 0.4, item, 13, WHITE)

add_text_box(slide, 1.3, 4.2, 4.0, 0.5, "6 people. Customer sees 3 separate conversations.", 12, RED)

# Arrow
add_text_box(slide, 5.9, 2.8, 1.5, 0.8, "→", 48, ORANGE, True, PP_ALIGN.CENTER)

# Tomorrow: Unified
add_card(slide, 7.5, 1.3, 5.0, 4.0, GREEN)
add_text_box(slide, 7.8, 1.5, 4.5, 0.4, "Tomorrow: Unified", 18, GREEN, True)

tomorrow_items = [
    "Cross-trained SA → Full Solution",
    "Cross-trained SA → Full Solution",
    "Cross-trained SA → Full Solution",
]
for i, item in enumerate(tomorrow_items):
    add_text_box(slide, 8.0, 2.2 + i * 0.6, 4.0, 0.4, item, 13, WHITE)

add_text_box(slide, 8.0, 4.2, 4.0, 0.5, "3 SAs. Customer sees ONE conversation.", 12, GREEN)

# Bottom stat
add_text_box(slide, 0.8, 5.8, 11.5, 0.8, "Same headcount → 3× coverage → Solution-led selling", 24, ORANGE, True, PP_ALIGN.CENTER)

# ==================== SLIDE 11: THE ASK — SOLUTION-LED GTM ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.3, 11, 0.7, "The Ask — From Service Selling to Solution Selling", 28, ORANGE, True)

# TOP: Before / After
# Before card
add_card(slide, 0.8, 1.1, 5.3, 3.0, RED)
add_text_box(slide, 1.0, 1.15, 5.0, 0.25, "❌ TODAY — SERVICE-LED", 10, RED, True)
before_items = [
    "SAs know one service, not the customer's problem",
    "Every deal starts from scratch — no reusable solutions",
    "6 people per engagement, customer sees 3 conversations",
    "No cross-service expertise — silos everywhere",
    "Competitors ship solutions while we sell components",
]
txBox = add_text_box(slide, 1.0, 1.5, 5.0, 0.3, before_items[0], 11, WHITE)
for item in before_items[1:]:
    add_paragraph(txBox.text_frame, item, 11, WHITE, space_before=5)

# Arrow
add_text_box(slide, 6.2, 2.0, 0.8, 0.6, "→", 32, ORANGE, True, PP_ALIGN.CENTER)

# After card
add_card(slide, 7.0, 1.1, 5.5, 3.0, GREEN)
add_text_box(slide, 7.2, 1.15, 5.0, 0.25, "✅ TOMORROW — SOLUTION-LED", 10, GREEN, True)
after_items = [
    "Start with the customer problem, build a repeatable solution",
    "Solutions built once, deployed many times across accounts",
    "One cross-trained SA owns the full conversation",
    "Code is free — customers pay only for AWS services",
    "AWS leads with outcomes, not SKUs",
]
txBox = add_text_box(slide, 7.2, 1.5, 5.0, 0.3, after_items[0], 11, WHITE)
for item in after_items[1:]:
    add_paragraph(txBox.text_frame, item, 11, WHITE, space_before=5)

# BOTTOM: Three-tier escalation
tiers = [
    ("Phase 1 — Prove It", "Investigative Intelligence",
     ["→ DOJ Antitrust pilot — 500TB, 12 weeks", "→ 61 deployments across FedCiv + DoD", "→ Validate solution-led model works"],
     "US Federal only", "$101M ARR", ORANGE),
    ("Phase 2 — Scale It", "Federal Solution Factory",
     ["→ 3–4 solutions/yr from real needs", "→ Fraud Detection — HHS Medicare, DOJ, Treasury", "→ Each solution: $25–50M Federal ARR"],
     "US Federal — multiple solutions", "$200–300M ARR", LIGHT_ORANGE),
    ("Phase 3 — Transform It", "Global Solution-Led GTM",
     ["→ Investigative Intel goes global", "→ Fraud Detection → FinServ, Insurance, Healthcare", "→ New solutions per sector from regional teams"],
     "Global cross-sector", "$500M+ ARR", GREEN),
]

for i, (phase, title, items, scope, arr, color) in enumerate(tiers):
    x = 0.8 + i * 4.2
    add_card(slide, x, 4.4, 3.8, 2.8, color)
    add_text_box(slide, x + 0.15, 4.45, 3.5, 0.2, phase, 9, color, True)
    add_text_box(slide, x + 0.15, 4.7, 3.5, 0.35, title, 14, WHITE, True)
    txBox = add_text_box(slide, x + 0.15, 5.15, 3.5, 0.3, items[0], 10, GRAY)
    for item in items[1:]:
        add_paragraph(txBox.text_frame, item, 10, GRAY, space_before=3)
    # Scope + ARR at bottom
    add_text_box(slide, x + 0.15, 6.5, 3.5, 0.2, scope, 8, DARK_GRAY)
    add_text_box(slide, x + 0.15, 6.7, 3.5, 0.3, arr, 14, color, True)

    # Arrows between tiers
    if i < 2:
        add_text_box(slide, x + 3.7, 5.2, 0.5, 0.5, "→", 20, ORANGE, True, PP_ALIGN.CENTER)

add_text_box(slide, 0.8, 7.1, 11.5, 0.3, "One pilot proves the model. The model changes how AWS sells.", 13, DARK_GRAY, False, PP_ALIGN.CENTER)

# ==================== SLIDE 12: BILL OF MATERIALS ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.3, 11, 0.7, "AWS Services Bill of Materials", 28, ORANGE, True)
add_text_box(slide, 0.8, 0.9, 11, 0.4, "Per deployment — 500TB case evidence workload", 13, GRAY)

# Table header
header_y = 1.4
add_text_box(slide, 0.8, header_y, 4.0, 0.4, "Service", 12, ORANGE, True)
add_text_box(slide, 5.0, header_y, 4.5, 0.4, "Role", 12, ORANGE, True)
add_text_box(slide, 10.0, header_y, 2.5, 0.4, "Est. Monthly Cost", 12, ORANGE, True, PP_ALIGN.RIGHT)

# Divider
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.75), Inches(11.5), Pt(1.5))
shape.fill.solid()
shape.fill.fore_color.rgb = ORANGE
shape.line.fill.background()

bom_items = [
    ("Amazon S3", "Data Lake + Frontend Hosting", "$11,500"),
    ("Amazon OpenSearch Serverless", "Vector search + full-text", "$8,200"),
    ("Amazon Bedrock (Claude 3)", "AI analysis, embeddings, RAG", "$12,000"),
    ("Amazon Textract", "Document parsing / OCR", "$4,500"),
    ("Amazon Neptune", "Knowledge graph", "$3,800"),
    ("Amazon Rekognition", "Image / face analysis", "$2,800"),
    ("AWS Lambda", "API + Ingestion compute", "$2,400"),
    ("AWS Step Functions", "Pipeline orchestration", "$500"),
    ("Amazon API Gateway", "REST API", "$350"),
    ("Amazon CloudWatch", "Monitoring & logging", "$200"),
    ("AWS KMS", "Encryption at rest", "$100"),
    ("AWS CDK", "Infrastructure as code", "$0"),
]

for i, (svc, role, cost) in enumerate(bom_items):
    y = 1.9 + i * 0.35
    add_text_box(slide, 0.8, y, 4.0, 0.3, svc, 11, WHITE)
    add_text_box(slide, 5.0, y, 4.5, 0.3, role, 11, GRAY)
    add_text_box(slide, 10.0, y, 2.5, 0.3, cost, 11, WHITE, False, PP_ALIGN.RIGHT)

# Total row
total_y = 1.9 + len(bom_items) * 0.35 + 0.1
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(total_y), Inches(11.5), Pt(1.5))
shape.fill.solid()
shape.fill.fore_color.rgb = ORANGE
shape.line.fill.background()

add_text_box(slide, 0.8, total_y + 0.05, 4.0, 0.3, "Total per deployment", 12, ORANGE, True)
add_text_box(slide, 10.0, total_y + 0.05, 2.5, 0.3, "$46,350/mo", 12, ORANGE, True, PP_ALIGN.RIGHT)

add_text_box(slide, 0.8, total_y + 0.4, 4.0, 0.3, "Annual per deployment", 12, ORANGE, True)
add_text_box(slide, 10.0, total_y + 0.4, 2.5, 0.3, "$556,200/yr", 12, ORANGE, True, PP_ALIGN.RIGHT)

# ==================== SLIDE 13: ARR OPPORTUNITY ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.3, 11, 0.7, "Federal ARR Opportunity — Path to $100M", 28, ORANGE, True)

# Deal Tier Model
add_text_box(slide, 0.8, 1.0, 6, 0.3, "DEAL TIER MODEL — BASED ON ACTUAL CBP ANCHOR", 10, ORANGE, True)

tier_header_y = 1.35
add_text_box(slide, 0.8, tier_header_y, 2.5, 0.3, "Deal Tier", 10, ORANGE, True)
add_text_box(slide, 3.3, tier_header_y, 1.0, 0.3, "Avg $/mo", 10, ORANGE, True, PP_ALIGN.RIGHT)
add_text_box(slide, 4.5, tier_header_y, 1.0, 0.3, "Annual", 10, ORANGE, True, PP_ALIGN.RIGHT)
add_text_box(slide, 5.6, tier_header_y, 1.5, 0.3, "Example", 10, ORANGE, True)

tiers = [
    ("🔴 Tier 1: High-Volume", "$600K", "$7.2M", "CBP Biometrics", RED),
    ("🟠 Tier 2: Large Investigation", "$150K", "$1.8M", "CBP TASPD/PSPD", LIGHT_ORANGE),
    ("🟡 Tier 3: Standard", "$46K", "$556K", "DOJ Antitrust, SEC", ORANGE),
]

for i, (tier, avg, annual, example, color) in enumerate(tiers):
    y = 1.65 + i * 0.3
    add_text_box(slide, 0.8, y, 2.5, 0.3, tier, 9, color, True)
    add_text_box(slide, 3.3, y, 1.0, 0.3, avg, 9, WHITE, False, PP_ALIGN.RIGHT)
    add_text_box(slide, 4.5, y, 1.0, 0.3, annual, 9, WHITE, False, PP_ALIGN.RIGHT)
    add_text_box(slide, 5.6, y, 1.5, 0.3, example, 8, DARK_GRAY)

# Department breakdown
add_text_box(slide, 0.8, 2.8, 6, 0.3, "DEPARTMENT BREAKDOWN", 10, ORANGE, True)

dept_header_y = 3.1
for col, (label, x) in enumerate([("Department", 0.8), ("T1", 3.5), ("T2", 4.0), ("T3", 4.5), ("Annual", 5.2)]):
    align = PP_ALIGN.CENTER if col in [1,2,3] else (PP_ALIGN.RIGHT if col == 4 else PP_ALIGN.LEFT)
    add_text_box(slide, x, dept_header_y, 1.0 if col < 4 else 1.5, 0.25, label, 9, ORANGE, True, align)

depts = [
    ("DHS/CBP", "2", "5", "4", "$25.6M", LIGHT_ORANGE),
    ("DHS (ICE, SS, CISA, TSA)", "1", "3", "4", "$14.8M", WHITE),
    ("DOJ (Antitrust, FBI, DEA)", "1", "4", "8", "$18.9M", WHITE),
    ("Treasury (IRS-CI, FinCEN)", "1", "2", "3", "$12.5M", WHITE),
    ("HHS (OIG, FDA, CMS)", "0", "2", "3", "$5.3M", WHITE),
    ("Other FedCiv (SEC, FTC...)", "0", "2", "5", "$6.4M", WHITE),
    ("FedCiv Total", "5", "18", "27", "$83.5M", ORANGE),
    ("DoD (JAG, CID, NCIS...)", "1", "4", "6", "$17.7M", BLUE),
    ("Combined Federal", "6", "22", "33", "$101.2M", ORANGE),
]

for i, (dept, t1, t2, t3, annual, color) in enumerate(depts):
    y = 3.35 + i * 0.28
    bold = dept in ["FedCiv Total", "Combined Federal"]
    add_text_box(slide, 0.8, y, 2.5, 0.25, dept, 8, color, bold)
    add_text_box(slide, 3.5, y, 0.5, 0.25, t1, 8, color, bold, PP_ALIGN.CENTER)
    add_text_box(slide, 4.0, y, 0.5, 0.25, t2, 8, color, bold, PP_ALIGN.CENTER)
    add_text_box(slide, 4.5, y, 0.5, 0.25, t3, 8, color, bold, PP_ALIGN.CENTER)
    add_text_box(slide, 5.2, y, 1.5, 0.25, annual, 8, color, bold, PP_ALIGN.RIGHT)

# Right side - anchor, use case, bar chart, total
add_card(slide, 7.2, 1.0, 5.3, 1.0, RED)
add_text_box(slide, 7.4, 1.05, 5.0, 0.25, "🔴 ANCHOR: CBP BIOMETRICS (PROVEN)", 9, RED, True)
add_text_box(slide, 7.4, 1.35, 5.0, 0.25, "Current: $300K/mo → doubling to $600K/mo", 10, WHITE)
add_text_box(slide, 7.4, 1.6, 5.0, 0.25, "+ 5 additional deals at TASPD, PSPD, BEMS, Cargo", 9, GRAY)

add_card(slide, 7.2, 2.2, 5.3, 1.2, LIGHT_ORANGE)
add_text_box(slide, 7.4, 2.25, 5.0, 0.25, "🟠 USE CASE: FENTANYL DETECTION PIPELINE", 9, LIGHT_ORANGE, True)
add_text_box(slide, 7.4, 2.55, 5.0, 0.7, "Millions of daily transactions → OpenSearch pattern filtering → Bedrock AI risk scoring → Neptune graph identifies hidden networks → Border officers act in real-time", 9, WHITE)

# Bar chart
add_text_box(slide, 7.2, 3.6, 5.0, 0.25, "ANNUAL AWS SPEND BY DEPARTMENT", 9, ORANGE, True)

bars = [
    ("DHS/CBP", 100, "$25.6M", ORANGE),
    ("DOJ", 73.8, "$18.9M", ORANGE),
    ("DHS Other", 57.8, "$14.8M", ORANGE),
    ("Treasury", 48.8, "$12.5M", ORANGE),
    ("HHS", 20.7, "$5.3M", ORANGE),
    ("Other", 25, "$6.4M", ORANGE),
    ("DoD", 69.1, "$17.7M", BLUE),
]

for i, (label, pct, value, color) in enumerate(bars):
    y = 3.9 + i * 0.32
    add_text_box(slide, 7.2, y, 1.0, 0.25, label, 8, GRAY, False, PP_ALIGN.RIGHT)
    bar_width = 3.5 * (pct / 100)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.3), Inches(y + 0.02), Inches(bar_width), Inches(0.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    add_text_box(slide, 8.3, y, bar_width, 0.25, value, 8, WHITE, True)

# Total card
add_card(slide, 7.2, 6.2, 5.3, 1.0, ORANGE)
add_text_box(slide, 7.4, 6.25, 5.0, 0.5, "$101.2M", 32, ORANGE, True, PP_ALIGN.CENTER)
add_text_box(slide, 7.4, 6.7, 5.0, 0.3, "Combined Federal ARR Opportunity", 11, GRAY, False, PP_ALIGN.CENTER)

# ==================== SLIDE 14: CALL TO ACTION ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.3, 11, 0.7, "Call to Action — Next Steps", 28, ORANGE, True)

# Timeline
timeline_items = [
    ("⚡", "Immediate", "Complete DOJ Antitrust questionnaire. Finalize 500TB pilot scope."),
    ("30", "30 Days", "Deploy pilot environment. Begin data ingestion with service team support."),
    ("90", "90 Days", "Pilot complete. Measure outcomes. Build case study."),
    ("6m", "6 Months", "Expand to 3 additional FedCiv departments."),
    ("1y", "12 Months", "Target 10+ deployments. $5.5M+ ARR."),
]

for i, (icon, time_label, desc) in enumerate(timeline_items):
    y = 1.2 + i * 1.1
    # Dot
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), Inches(y), Inches(0.5), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x40)
    shape.line.color.rgb = ORANGE
    shape.line.width = Pt(2)
    shape.text_frame.paragraphs[0].text = icon
    shape.text_frame.paragraphs[0].font.size = Pt(10)
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].font.color.rgb = ORANGE
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Connecting line
    if i < len(timeline_items) - 1:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.22), Inches(y + 0.5), Pt(2), Inches(0.6))
        line.fill.solid()
        line.fill.fore_color.rgb = ORANGE
        line.line.fill.background()

    add_text_box(slide, 1.7, y - 0.05, 2.0, 0.3, time_label, 13, LIGHT_ORANGE, True)
    add_text_box(slide, 1.7, y + 0.25, 4.5, 0.4, desc, 12, WHITE)

# Right side cards
# The Ask
add_card(slide, 7.5, 1.2, 5.0, 1.6, GREEN)
add_text_box(slide, 7.7, 1.3, 4.6, 0.4, "🎯", 24, GREEN, False, PP_ALIGN.CENTER)
add_text_box(slide, 7.7, 1.7, 4.6, 0.3, "The Ask", 16, GREEN, True, PP_ALIGN.CENTER)
add_text_box(slide, 7.7, 2.1, 4.6, 0.5, "Approve the DOJ Antitrust pilot.\n12 weeks. ~$200K. Prove the platform.", 13, WHITE, False, PP_ALIGN.CENTER)

# The Prize
add_card(slide, 7.5, 3.1, 5.0, 1.6, YELLOW)
add_text_box(slide, 7.7, 3.2, 4.6, 0.4, "💰", 24, YELLOW, False, PP_ALIGN.CENTER)
add_text_box(slide, 7.7, 3.6, 4.6, 0.3, "The Prize", 16, LIGHT_ORANGE, True, PP_ALIGN.CENTER)
add_text_box(slide, 7.7, 4.0, 4.6, 0.5, "$101M federal ARR opportunity.\n61 deployments across FedCiv + DoD.", 13, WHITE, False, PP_ALIGN.CENTER)

# The Advantage
add_card(slide, 7.5, 5.0, 5.0, 1.6, BLUE)
add_text_box(slide, 7.7, 5.1, 4.6, 0.4, "🚀", 24, BLUE, False, PP_ALIGN.CENTER)
add_text_box(slide, 7.7, 5.5, 4.6, 0.3, "The Advantage", 16, BLUE, True, PP_ALIGN.CENTER)
add_text_box(slide, 7.7, 5.9, 4.6, 0.5, "Solution is built. Code is free.\nCustomers pay only for AWS services.", 13, WHITE, False, PP_ALIGN.CENTER)

# ==================== SAVE ====================
output_path = "investigative-intelligence-gtm-v3.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
