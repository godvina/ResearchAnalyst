"""Convert deployment challenge/solution slides to PowerPoint."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_textbox(slide, left, top, width, height, text, size=11, color=(0xE2,0xE8,0xF0), bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = RGBColor(*color)
    p.font.bold = bold
    p.alignment = align
    return tf

def add_box(slide, left, top, width, height, fill_rgb, border_rgb=None):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill_rgb)
    if border_rgb:
        shape.line.color.rgb = RGBColor(*border_rgb)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

# ============ SLIDE 1: THE CHALLENGE ============
s1 = prs.slides.add_slide(prs.slide_layouts[6])
s1.background.fill.solid()
s1.background.fill.fore_color.rgb = RGBColor(0x0a, 0x16, 0x28)

add_textbox(s1, 0.5, 0.3, 3, 0.3, "⚠ ESCALATION REQUIRED", 10, (0xFC,0x81,0x81), True)
add_textbox(s1, 0.5, 0.6, 12, 0.5, "Investigative Intelligence Platform — Deployment Challenge", 26, (0xFF,0xFF,0xFF), True)
add_textbox(s1, 0.5, 1.1, 12, 0.4, "A cross-service solution is ready for customer deployment, but our organizational structure creates gaps no single team can fill.", 12, (0xA0,0xAE,0xC0))

# Team cards
teams = [
    ("1", "OpenSearch Service Team", "Building configurable Data Loader\nas separate AWS service"),
    ("2", "Solution Acceleration Team", "GitHub repo, testing, approval\nfor internal → customer (3-4 months)"),
    ("3", "Account SA", "Rebuilding app using Kiro prompts\nfor GovCloud deployment"),
    ("4", "WWSO Service Team", "Neptune SA + OpenSearch SA\narchitecture guidance"),
]
for i, (num, name, role) in enumerate(teams):
    x = 0.5 + (i % 2) * 3.2
    y = 1.7 + (i // 2) * 1.3
    add_box(s1, x, y, 3.0, 1.1, (0x0D,0x1F,0x3C), (0x2D,0x55,0x8A))
    tf = add_textbox(s1, x+0.1, y+0.08, 2.8, 1.0, f"{num}  {name}", 12, (0x90,0xCD,0xF4), True)
    p2 = tf.add_paragraph()
    p2.text = role
    p2.font.size = Pt(10)
    p2.font.color.rgb = RGBColor(0xA0, 0xAE, 0xC0)
    p2.space_before = Pt(4)

# Challenge boxes (right side)
challenges = [
    ("🔴 The Deployment Gap", "ProServe lacks specialists in OpenSearch, Neptune, and cross-service architecture. Partners don't have this expertise either. No single team owns end-to-end deployment of a multi-service solution.", (0xFC,0x81,0x81)),
    ("🔴 IP Fragmentation Risk", "Account SA is rebuilding from scratch using Kiro prompts instead of deploying the existing CDK/CloudFormation template. Risks losing months of battle-tested code, 30+ resolved deployment issues, and config-driven GovCloud support.", (0xFC,0x81,0x81)),
    ("🔴 Organizational Misalignment", "AWS is organized by service teams. This solution requires 6+ services (Bedrock, Aurora, Neptune, OpenSearch, Lambda, Step Functions). No team owns 'cross-service investigative solutions.'", (0xFC,0x81,0x81)),
]
for i, (title, desc, color) in enumerate(challenges):
    y = 1.7 + i * 1.55
    add_box(s1, 6.8, y, 6.0, 1.4, (0x1A,0x0D,0x0D), (0x5C,0x2D,0x2D))
    tf = add_textbox(s1, 7.0, y+0.08, 5.6, 1.3, title, 12, color, True)
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(10)
    p2.font.color.rgb = RGBColor(0xCB, 0xD5, 0xE0)
    p2.space_before = Pt(4)

# Timeline risk
add_box(s1, 0.5, 4.5, 6.0, 2.2, (0x1A,0x1A,0x0D), (0x5C,0x4D,0x2D))
tf = add_textbox(s1, 0.7, 4.6, 5.6, 2.0, "⚠ Timeline Risk", 12, (0xF6,0xAD,0x55), True)
risks = [
    "⏱ Customer expects pilot in 8-12 weeks",
    "⏱ Solution Acceleration approval: 3-4 months",
    "⏱ Kiro rebuild: unknown timeline, untested at scale",
    "⏱ ProServe staffing gap: no cross-service architects",
]
for r in risks:
    p = tf.add_paragraph()
    p.text = r
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0xCB, 0xD5, 0xE0)
    p.space_before = Pt(3)

add_textbox(s1, 0.5, 7.1, 6, 0.3, "AWS Emerging Tech Solutions", 9, (0x4A,0x55,0x68))
add_textbox(s1, 7, 7.1, 6, 0.3, "Confidential — Internal Leadership Review", 9, (0x4A,0x55,0x68), align=PP_ALIGN.RIGHT)

# ============ SLIDE 2: THE SOLUTION ============
s2 = prs.slides.add_slide(prs.slide_layouts[6])
s2.background.fill.solid()
s2.background.fill.fore_color.rgb = RGBColor(0x0a, 0x16, 0x28)

add_textbox(s2, 0.5, 0.3, 3, 0.3, "✓ PROPOSED PATH FORWARD", 10, (0x68,0xD3,0x91), True)
add_textbox(s2, 0.5, 0.6, 12, 0.5, "Solution-in-a-Box: Deploy What's Built, Not Rebuild", 26, (0xFF,0xFF,0xFF), True)
add_textbox(s2, 0.5, 1.1, 12, 0.4, "The application is production-tested with 82K documents, config-driven for GovCloud, and deployable via CloudFormation in 20 minutes.", 12, (0xA0,0xAE,0xC0))

# What exists
add_box(s2, 0.5, 1.6, 6.0, 2.4, (0x0D,0x1F,0x1A), (0x2D,0x5C,0x3D))
tf = add_textbox(s2, 0.7, 1.68, 5.6, 2.2, "✅ What Already Exists", 13, (0x48,0xBB,0x78), True)
items = [
    "Config-driven CDK: Demo → GovCloud Test → Production",
    "CloudFormation template — deployable via console, no CDK needed",
    "82,529 docs, 115K entities, 175K relationships processed",
    "Well-Architected review, FedRAMP model registry built",
    "30+ deployment issues documented and resolved",
    "Label-based access control, audit trail, data governance",
]
for item in items:
    p = tf.add_paragraph()
    p.text = "• " + item
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0xCB, 0xD5, 0xE0)
    p.space_before = Pt(2)

# Why Kiro rebuild is risky
add_box(s2, 0.5, 4.2, 6.0, 1.3, (0x1A,0x0D,0x1A), (0x5C,0x2D,0x5C))
tf = add_textbox(s2, 0.7, 4.28, 5.6, 1.2, "⚠ Why Kiro Rebuild Is Risky", 12, (0xB7,0x94,0xF4), True)
p = tf.add_paragraph()
p.text = "Kiro-generated code won't include 30+ battle-tested fixes (VPC endpoint SG rules, Neptune HTTP API, chunked extraction, embed truncation). These took weeks to discover. A rebuild starts from zero on all of them."
p.font.size = Pt(10)
p.font.color.rgb = RGBColor(0xCB, 0xD5, 0xE0)
p.space_before = Pt(4)

# Timeline
add_textbox(s2, 6.8, 1.55, 6, 0.3, "PROPOSED TIMELINE", 10, (0x90,0xCD,0xF4), True)
weeks = [
    ("W1", "Now", "Deploy template to\nGovCloud Isengard"),
    ("W2-3", "Pilot", "Customer GovCloud\nwith synthetic data"),
    ("W4-8", "Scale", "150GB customer data\nAuth + RBAC"),
    ("W8-12", "Production", "ATO process\nFull deployment"),
]
for i, (w, label, desc) in enumerate(weeks):
    x = 6.8 + i * 1.55
    add_box(s2, x, 1.9, 1.4, 1.6, (0x0D,0x1F,0x3C), (0x2D,0x55,0x8A))
    tf = add_textbox(s2, x+0.1, 1.95, 1.2, 1.5, w, 20, (0x63,0xB3,0xED), True, PP_ALIGN.CENTER)
    p = tf.add_paragraph()
    p.text = label
    p.font.size = Pt(8)
    p.font.color.rgb = RGBColor(0x71, 0x80, 0x96)
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = desc
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(0xCB, 0xD5, 0xE0)
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(6)

# The Ask
add_box(s2, 6.8, 3.7, 6.0, 3.0, (0x1A,0x1A,0x0D), (0x8C,0x7A,0x3D))
tf = add_textbox(s2, 7.0, 3.78, 5.6, 2.8, "📋 The Ask", 14, (0xF6,0xAD,0x55), True)
asks = [
    "1️⃣  Align Account SA to deploy existing template vs. Kiro rebuild — preserve months of IP",
    "2️⃣  Fast-track Solution Acceleration review — application is ready now, not in 3-4 months",
    "3️⃣  Assign one cross-service architect (ProServe or WWSO) for GovCloud deployment",
    "4️⃣  Sponsor 'Solution-in-a-Box' model — deploy the integrated asset, not individual service components",
]
for ask in asks:
    p = tf.add_paragraph()
    p.text = ask
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    p.space_before = Pt(8)

# Solution-in-a-box concept
add_box(s2, 0.5, 5.7, 6.0, 1.0, (0x0D,0x1F,0x3C), (0x48,0xBB,0x78))
tf = add_textbox(s2, 0.7, 5.78, 5.6, 0.9, "🎯 Solution-in-a-Box Model", 12, (0x48,0xBB,0x78), True)
p = tf.add_paragraph()
p.text = "One CloudFormation template + one Lambda zip + one frontend file. A single SA deploys in a day. No Neptune specialist, no OpenSearch specialist needed. This is how we scale cross-service solutions."
p.font.size = Pt(10)
p.font.color.rgb = RGBColor(0xCB, 0xD5, 0xE0)
p.space_before = Pt(4)

add_textbox(s2, 0.5, 7.1, 6, 0.3, "AWS Emerging Tech Solutions", 9, (0x4A,0x55,0x68))
add_textbox(s2, 7, 7.1, 6, 0.3, "Confidential — Internal Leadership Review", 9, (0x4A,0x55,0x68), align=PP_ALIGN.RIGHT)

out = os.path.join("docs", "presentation", "Deployment-Challenge-Solution.pptx")
prs.save(out)
print(f"Saved: {out}")
