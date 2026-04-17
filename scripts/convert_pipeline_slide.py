"""Convert the AI pipeline slide to PowerPoint."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = RGBColor(0x0a, 0x16, 0x28)

# Title
txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "How AI Transforms Raw Documents into Investigative Intelligence"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# Subtitle
txBox2 = slide.shapes.add_textbox(Inches(0.6), Inches(0.8), Inches(12), Inches(0.3))
tf2 = txBox2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "FROM 82,000 UNSTRUCTURED DOCUMENTS TO SEARCHABLE, CONNECTED INTELLIGENCE"
p2.font.size = Pt(12)
p2.font.bold = True
p2.font.color.rgb = RGBColor(0x63, 0xB3, 0xED)

# Pipeline stages
stages = [
    ("1", "📄", "INGEST", "Raw documents (PDFs, scans, emails)\nuploaded to S3. OCR extracts\ntext from scanned images.", "S3 · Textract · Step Functions"),
    ("2", "🧠", "EMBED", "AI converts each document's meaning\ninto a 1,536-number coordinate —\nits position in 'meaning space.'", "Bedrock · Titan Embed · pgvector"),
    ("3", "🔍", "EXTRACT", "AI identifies every person, org,\nlocation, date, and financial\namount in each document.", "Bedrock · Claude Haiku · NER"),
    ("4", "🕸️", "CONNECT", "Entities become nodes in a\nknowledge graph. Relationships\nbetween people and places mapped.", "Neptune · Gremlin · Graph DB"),
]

x_start = 0.5
for i, (num, icon, title, desc, tech) in enumerate(stages):
    x = x_start + i * 3.2
    # Stage box
    shape = slide.shapes.add_shape(1, Inches(x), Inches(1.4), Inches(2.8), Inches(2.2))  # rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1A, 0x36, 0x5D)
    shape.line.color.rgb = RGBColor(0x2D, 0x55, 0x8A)
    shape.line.width = Pt(1)

    # Stage content
    txBox = slide.shapes.add_textbox(Inches(x + 0.15), Inches(1.5), Inches(2.5), Inches(2.0))
    tf = txBox.text_frame
    tf.word_wrap = True

    # Icon + Title
    p = tf.paragraphs[0]
    p.text = f"{icon}  {title}"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x63, 0xB3, 0xED)

    # Description
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(10)
    p2.font.color.rgb = RGBColor(0xA0, 0xAE, 0xC0)
    p2.space_before = Pt(8)

    # Tech badges
    p3 = tf.add_paragraph()
    p3.text = tech
    p3.font.size = Pt(8)
    p3.font.color.rgb = RGBColor(0x90, 0xCD, 0xF4)
    p3.font.bold = True
    p3.space_before = Pt(10)

    # Arrow between stages
    if i < 3:
        arrow_x = x + 2.85
        txA = slide.shapes.add_textbox(Inches(arrow_x), Inches(2.2), Inches(0.3), Inches(0.4))
        tfA = txA.text_frame
        pA = tfA.paragraphs[0]
        pA.text = "→"
        pA.font.size = Pt(24)
        pA.font.color.rgb = RGBColor(0x4A, 0x9E, 0xFF)
        pA.font.bold = True

# Bottom section - Embeddings box
emb_shape = slide.shapes.add_shape(1, Inches(0.5), Inches(3.9), Inches(6.0), Inches(3.0))
emb_shape.fill.solid()
emb_shape.fill.fore_color.rgb = RGBColor(0x0D, 0x1F, 0x2D)
emb_shape.line.color.rgb = RGBColor(0x48, 0xBB, 0x78)
emb_shape.line.width = Pt(2)

txEmb = slide.shapes.add_textbox(Inches(0.7), Inches(4.0), Inches(5.6), Inches(2.8))
tf = txEmb.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "📐  Vector Embeddings — Semantic Search"
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = RGBColor(0x48, 0xBB, 0x78)

p2 = tf.add_paragraph()
p2.text = "Every document gets a mathematical fingerprint that captures what it's about, not just what words it contains. Documents about similar topics cluster together — enabling search by meaning."
p2.font.size = Pt(11)
p2.font.color.rgb = RGBColor(0xCB, 0xD5, 0xE0)
p2.space_before = Pt(10)

p3 = tf.add_paragraph()
p3.text = 'Search: "financial fraud involving shell companies"'
p3.font.size = Pt(10)
p3.font.italic = True
p3.font.color.rgb = RGBColor(0x63, 0xB3, 0xED)
p3.space_before = Pt(14)

p4 = tf.add_paragraph()
p4.text = '✓ Finds documents about "suspicious wire transfers to offshore accounts" — even though no words match'
p4.font.size = Pt(10)
p4.font.color.rgb = RGBColor(0x68, 0xD3, 0x91)
p4.space_before = Pt(4)

# Bottom section - Entities box
ent_shape = slide.shapes.add_shape(1, Inches(6.8), Inches(3.9), Inches(6.0), Inches(3.0))
ent_shape.fill.solid()
ent_shape.fill.fore_color.rgb = RGBColor(0x0D, 0x1F, 0x2D)
ent_shape.line.color.rgb = RGBColor(0xF6, 0xAD, 0x55)
ent_shape.line.width = Pt(2)

txEnt = slide.shapes.add_textbox(Inches(7.0), Inches(4.0), Inches(5.6), Inches(2.8))
tf = txEnt.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "🏷️  Entity Extraction — Knowledge Graph"
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = RGBColor(0xF6, 0xAD, 0x55)

p2 = tf.add_paragraph()
p2.text = "AI reads every document and highlights the important stuff — who, what, where, when. These become connected nodes in a graph that reveals hidden patterns across thousands of documents."
p2.font.size = Pt(11)
p2.font.color.rgb = RGBColor(0xCB, 0xD5, 0xE0)
p2.space_before = Pt(10)

p3 = tf.add_paragraph()
p3.text = "Extracted from one document:"
p3.font.size = Pt(9)
p3.font.color.rgb = RGBColor(0x71, 0x80, 0x96)
p3.space_before = Pt(14)

entities = [
    ("👤 Jeffrey Epstein", RGBColor(0x90, 0xCD, 0xF4)),
    ("🏢 JP Morgan Chase", RGBColor(0xFB, 0xD3, 0x8D)),
    ("🏢 LSJE LLC", RGBColor(0xFB, 0xD3, 0x8D)),
    ("📍 Palm Beach, FL", RGBColor(0x9A, 0xE6, 0xB4)),
    ("📅 March 2005", RGBColor(0xD6, 0xBC, 0xFA)),
]

p4 = tf.add_paragraph()
p4.text = "  ".join(e[0] for e in entities)
p4.font.size = Pt(10)
p4.font.bold = True
p4.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
p4.space_before = Pt(6)

# Footer
txFoot = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(6), Inches(0.3))
tf = txFoot.text_frame
p = tf.paragraphs[0]
p.text = "AWS Investigative Intelligence Platform — Bedrock · Aurora pgvector · Neptune · S3"
p.font.size = Pt(9)
p.font.color.rgb = RGBColor(0x4A, 0x55, 0x68)

txFoot2 = slide.shapes.add_textbox(Inches(7), Inches(7.0), Inches(6), Inches(0.3))
tf2 = txFoot2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "82,529 documents · 115,422 entities · 175,690 relationships"
p2.font.size = Pt(9)
p2.font.color.rgb = RGBColor(0x4A, 0x55, 0x68)
p2.alignment = PP_ALIGN.RIGHT

out = os.path.join("docs", "presentation", "AI-Intelligence-Pipeline.pptx")
prs.save(out)
print(f"Saved: {out}")
